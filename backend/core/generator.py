import json
import os
import re
import asyncio
import copy
import math
import time
import secrets
import random
import zipfile
from datetime import datetime, timezone
from difflib import SequenceMatcher
from pathlib import Path
from docx import Document
from docx.shared import Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from google import genai
from google.genai import types as genai_types
import logging
from pypdf import PdfReader
from core.security import sanitize_filename

from core.base_docs_loader import BaseDocsLoader
from core.docx_parser import DocxParser
from core.presentation_analyzer import PresentationAnalyzer


logger = logging.getLogger(__name__)
os.makedirs("storage", exist_ok=True)



# Стильные палитры для презентаций
PALETTES = [
    {"bg": (248, 250, 252), "primary": (37, 99, 235), "text": (15, 23, 42), "accent": (219, 234, 254)},
    {"bg": (255, 251, 235), "primary": (217, 119, 6), "text": (69, 26, 3), "accent": (254, 243, 199)},
    {"bg": (240, 253, 244), "primary": (16, 185, 129), "text": (6, 78, 59), "accent": (209, 250, 229)},
    {"bg": (250, 245, 255), "primary": (147, 51, 234), "text": (49, 46, 129), "accent": (243, 232, 255)},
    {"bg": (30, 41, 59), "primary": (56, 189, 248), "text": (248, 250, 252), "accent": (15, 23, 42)}
]

BULLET_PREFIX_RE = re.compile(
    r"^\s*(?:"
    r"[\u2022\u00B7\u25CF\u25AA\u25A0\u25E6\u2023\u2013\u2014\-]+|"
    r"\*+|"
    r"\d+[\.\)]|"
    r"[A-Za-z\u0400-\u04FF]\)|"
    r"\([0-9A-Za-z\u0400-\u04FF]+\)"
    r")\s*"
)
LEADING_ORDERED_TITLE_RE = re.compile(r"^\d+[\.\)]\s*(.+)$")
PLACEHOLDER_TEXTS = {
    "немає",
    "нема",
    "відсутнє",
    "відсутні",
    "не вказано",
    "не зазначено",
    "n/a",
    "-",
    "—",
}

def clean_text(value):
    text = str(value or "").strip()
    if not text:
        return ""
    text = text.replace("\u200b", "")
    text = text.replace("**", "").replace("__", "")
    text = re.sub(r'^[*_`]+|[*_`]+$', '', text).strip()
    for _ in range(2):
        new_text = BULLET_PREFIX_RE.sub("", text).strip()
        if new_text == text:
            break
        text = new_text
    text = re.sub(r"\s+", " ", text).strip()
    return text

def clean_text_preserve_prefix(value):
    text = str(value or "").strip()
    if not text:
        return ""
    text = text.replace("\u200b", "")
    text = text.replace("**", "").replace("__", "")
    text = re.sub(r'^[`]+|[`]+$', '', text).strip()
    text = re.sub(r"^[•*]+", "", text).strip()
    text = re.sub(r"\s+", " ", text).strip()
    return text

def normalize_sentence_punctuation(value):
    text = clean_text_preserve_prefix(value)
    if not text:
        return ""
    text = re.sub(r"\s+([,;:\.\!\?])", r"\1", text)
    text = re.sub(r"([,;:\.\!\?]){2,}", r"\1", text)
    text = re.sub(r"\.{2,}", ".", text)
    text = re.sub(r";{2,}", ";", text)
    text = re.sub(r":{2,}", ":", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip()

def is_placeholder_text(value):
    text = clean_text_preserve_prefix(value).lower().strip()
    text = text.strip(" .,:;!?")
    return text in PLACEHOLDER_TEXTS

def normalize_list(val):
    if not val:
        return []
    items = []
    if isinstance(val, list):
        items = val
    elif isinstance(val, str):
        raw = val.strip()
        if "\n" in raw:
            items = [line for line in raw.splitlines() if line.strip()]
        else:
            items = [raw]
    else:
        items = [val]

    cleaned = []
    for item in items:
        if item is None:
            continue
        if isinstance(item, (list, tuple)):
            for sub in item:
                cleaned_text = clean_text(sub)
                if cleaned_text and not is_placeholder_text(cleaned_text):
                    cleaned.append(cleaned_text)
        else:
            cleaned_text = clean_text(item)
            if cleaned_text and not is_placeholder_text(cleaned_text):
                cleaned.append(cleaned_text)

    if len(cleaned) == 1 and ";" in cleaned[0]:
        cleaned = [clean_text(part) for part in cleaned[0].split(";")]

    return [c for c in cleaned if c]

def normalize_block_list(val):
    if not val:
        return []
    if isinstance(val, list):
        items = val
    elif isinstance(val, str):
        raw = val.strip()
        items = [line for line in raw.splitlines() if line.strip()] if "\n" in raw else [raw]
    else:
        items = [val]

    cleaned = []
    for item in items:
        if item is None:
            continue
        if isinstance(item, (list, tuple)):
            for sub in item:
                cleaned_text = clean_text_preserve_prefix(sub)
                if cleaned_text and not is_placeholder_text(cleaned_text):
                    cleaned.append(cleaned_text)
        else:
            cleaned_text = clean_text_preserve_prefix(item)
            if cleaned_text and not is_placeholder_text(cleaned_text):
                cleaned.append(cleaned_text)
    return [c for c in cleaned if c]

def normalize_expected_results(val):
    if isinstance(val, dict):
        knowledge = normalize_list(val.get("knowledge") or val.get("знання") or val.get("знаннєві"))
        skills = normalize_list(val.get("skills") or val.get("уміння") or val.get("діяльнісні") or val.get("навички"))
        values = normalize_list(val.get("values") or val.get("цінності") or val.get("ціннісні"))
        if not (knowledge or skills or values):
            flat = []
            for v in val.values():
                flat.extend(normalize_list(v))
            skills = flat
        return {"knowledge": knowledge, "skills": skills, "values": values}

    items = normalize_list(val)
    knowledge, skills, values = [], [], []
    for item in items:
        low = item.lower()
        if low.startswith(("знання", "знанн", "знаннєві")):
            knowledge.append(clean_text(item.split(":", 1)[-1] if ":" in item else item))
        elif low.startswith(("уміння", "діяльн", "навички")):
            skills.append(clean_text(item.split(":", 1)[-1] if ":" in item else item))
        elif low.startswith(("цінніс", "цінності")):
            values.append(clean_text(item.split(":", 1)[-1] if ":" in item else item))
        else:
            skills.append(item)
    return {"knowledge": knowledge, "skills": skills, "values": values}

def format_time(value):
    text = clean_text(value)
    if not text:
        return ""
    if re.search(r"\d", text) and "хв" not in text:
        return f"{text} хв"
    return text

class SlideRenderer:
    def __init__(self, prs, palette):
        self.prs = prs
        self.bg = PptxRGB(*palette["bg"])
        self.primary = PptxRGB(*palette["primary"])
        self.text_col = PptxRGB(*palette["text"])
        self.accent = PptxRGB(*palette["accent"])
        self.MARGIN = Inches(0.8)
        self.W = Inches(13.333)
        self.H = Inches(7.5)
        self.FONT_TITLE = "Calibri"
        self.FONT_BODY = "Calibri"
        self.FONT_SMALL = "Calibri"

    def _make_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Повністю пустий слайд
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.bg
        self._add_accents(slide)
        return slide

    def _add_accents(self, slide):
        # Ліва акцентна смуга
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), self.H)
        band.fill.solid()
        band.fill.fore_color.rgb = self.primary
        band.line.fill.background()

        # Верхня тонка лінія
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), 0, self.W - Inches(0.28), Inches(0.08))
        top.fill.solid()
        top.fill.fore_color.rgb = self.accent
        top.line.fill.background()
        top.fill.transparency = 0.2

        # М'яке коло-акцент
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(0.2), Inches(2.8), Inches(2.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.accent
        circle.fill.transparency = 0.78
        circle.line.fill.background()

    def render_title(self, data):
        slide = self._make_slide()
        # Малий бейдж
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.4), Inches(3.2), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.primary
        badge.fill.transparency = 0.1
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.text = "КОНСПЕКТ УРОКУ"
        badge_tf.paragraphs[0].font.size = PptxPt(12)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.name = self.FONT_SMALL
        badge_tf.paragraphs[0].font.color.rgb = PptxRGB(255, 255, 255)

        tx = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.333), Inches(2.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "Без назви")
        p.font.size = PptxPt(54)
        p.font.bold = True
        p.font.name = self.FONT_TITLE
        p.font.color.rgb = self.primary

        subtitle = data.get("subtitle")
        if subtitle:
            st = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(10.5), Inches(0.8)).text_frame
            st.word_wrap = True
            sp = st.paragraphs[0]
            sp.text = subtitle
            sp.font.size = PptxPt(22)
            sp.font.name = self.FONT_BODY
            sp.font.color.rgb = self.text_col

        # Бренд у футері
        footer = slide.shapes.add_textbox(Inches(10.2), Inches(6.9), Inches(3.0), Inches(0.4)).text_frame
        fp = footer.paragraphs[0]
        fp.text = "METODIST AI"
        fp.font.size = PptxPt(10)
        fp.font.name = self.FONT_SMALL
        fp.font.color.rgb = self.text_col
        fp.font.bold = True

    def render_bullets(self, data, index=None, total=None):
        slide = self._make_slide()
        
        # Заголовок слайду
        tx_title = slide.shapes.add_textbox(self.MARGIN, Inches(0.6), Inches(11.0), Inches(1.2))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        emoji = str(data.get("emoji", "")).strip()
        title_text = clean_text(data.get("title", ""))
        p_title.text = f"{emoji} {title_text}".strip()
        p_title.font.size = PptxPt(38)
        p_title.font.bold = True
        p_title.font.name = self.FONT_TITLE
        p_title.font.color.rgb = self.primary

        # Акцентна лінія під заголовком
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.MARGIN, Inches(1.55), Inches(5.5), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = self.primary
        line.line.fill.background()

        # Легка панель під контент
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.85), Inches(12.0), Inches(5.2))
        panel.fill.solid()
        panel.fill.fore_color.rgb = PptxRGB(255, 255, 255)
        panel.fill.transparency = 0.15
        panel.line.fill.background()

        # Блок для тексту (список)
        body_shape = slide.shapes.add_textbox(self.MARGIN, Inches(2.0), Inches(11.5), Inches(4.8))
        tf_body = body_shape.text_frame
        tf_body.word_wrap = True # Текст не буде вилазити за екран!
        tf_body.text = ""

        bullets = data.get("bullets", [])
        for bullet in bullets:
            clean_bullet = clean_text(bullet)
            if not clean_bullet:
                continue
            p = tf_body.add_paragraph()
            p.text = f"• {clean_bullet}" # Додаємо красивий маркер
            p.font.size = PptxPt(26)
            p.font.name = self.FONT_BODY
            p.font.color.rgb = self.text_col
            p.space_after = PptxPt(14) # Відступ між пунктами

        if index is not None and total is not None:
            brand = slide.shapes.add_textbox(Inches(0.9), Inches(6.9), Inches(3.0), Inches(0.35)).text_frame
            bp = brand.paragraphs[0]
            bp.text = "METODIST AI"
            bp.font.size = PptxPt(10)
            bp.font.name = self.FONT_SMALL
            bp.font.color.rgb = self.text_col
            bp.font.bold = True

            footer = slide.shapes.add_textbox(Inches(11.3), Inches(6.9), Inches(1.8), Inches(0.35)).text_frame
            fp = footer.paragraphs[0]
            fp.text = f"{index}/{total}"
            fp.font.size = PptxPt(10)
            fp.font.name = self.FONT_SMALL
            fp.font.color.rgb = self.text_col

class LessonGenerator:
    OPTIONAL_SECTION_FIELD_MAP = {
        "Завдання": "tasks",
        "Очікувані результати": "expected_results",
        "Ключові компетентності": "key_competencies",
        "Наскрізні вміння": "cross_cutting_skills",
        "Ціннісні орієнтири": "values",
        "Інтеграція": "integration",
        "Міжпредметні зв'язки": "integration",
        "Методи": "methods",
        "Форми роботи": "forms",
        "Формувальне оцінювання": "assessment",
        "Диференціація": "differentiation",
        "Ресурси": "resources",
        "Домашнє завдання": "homework",
    }
    DEFAULT_GENERIC_HEADER_ORDER = ("Тема", "Клас", "Предмет", "Мета", "Обладнання")
    DEFAULT_GENERIC_STAGE_SUBSTEPS = ("Вчитель", "Учні", "Діяльність")
    TRIVIAL_EQUIPMENT_NAMES = {
        "крейда",
        "дошка",
        "класна дошка",
        "ручка",
        "ручки",
        "олівець",
        "олівці",
        "простий олівець",
        "кольорові олівці",
        "зошит",
        "зошити",
        "лінійка",
        "лінійки",
    }
    EXPLICIT_STAGE_LABEL_MAP = {
        "мета": "goal",
        "вчитель": "teacher_actions",
        "учні": "student_actions",
        "діяльність": "activities",
        "оцінювання": "assessment",
        "самооцінювання": "assessment",
        "взаємооцінювання": "assessment",
        "рефлексія": "reflection",
        "диференціація": "differentiation",
        "матеріали": "materials",
    }
    GOAL_LINE_PREFIXES = (
        "ознайомити",
        "формувати",
        "закріпити",
        "розвивати",
        "виховувати",
        "перевірити",
        "узагальнити",
        "повторити",
        "створити",
        "навчити",
        "вчити",
        "відпрацювати",
        "налаштувати",
        "пояснити",
    )
    STUDENT_ACTION_PREFIXES = (
        "виконують",
        "відповідають",
        "слухають",
        "розглядають",
        "читають",
        "записують",
        "працюють",
        "називають",
        "будують",
        "висловлюють",
        "малюють",
        "обговорюють",
        "рухаються",
        "спостерігають",
        "складають",
        "повторюють",
        "розв'язують",
        "запам'ятовують",
        "викладають",
        "порівнюють",
        "характеризують",
        "слухають",
        "вивчають",
        "визначають",
        "пишуть",
    )
    TEACHER_ACTION_PREFIXES = (
        "подивіться",
        "послухайте",
        "пригадайте",
        "прочитайте",
        "відкрийте",
        "запишіть",
        "скажіть",
        "назвіть",
        "поміркуйте",
        "розгляньте",
        "погляньте",
        "спробуйте",
        "давайте",
        "зараз",
        "поясніть",
        "порівняйте",
        "знайдіть",
        "виконайте",
        "послухаймо",
        "пригадаймо",
        "пояснює",
        "ставить",
        "пропонує",
        "організовує",
        "демонструє",
        "оголошує",
        "підводить",
        "просить",
        "зачитує",
        "показує",
        "вмикає",
        "роздає",
        "звертає увагу",
    )
    MATERIAL_HINT_TOKENS = (
        "картк",
        "підруч",
        "зошит",
        "презентац",
        "буквар",
        "аудіозапис",
        "відео",
        "схем",
        "фішк",
        "малюн",
        "ілюстрац",
        "карта",
        "герб",
        "прапор",
        "олів",
        "папір",
        "роздат",
        "таблиц",
        "наочн",
        "плакат",
    )
    STAGE_RENDER_SUBSTEP_ORDER = (
        "Вчитель",
        "Учні",
        "Діяльність",
    )

    WEAK_GENERIC_PATTERNS = (
        "провести обговорення",
        "ознайомити учнів",
        "виконати завдання",
        "провести рефлексію",
        "обговорення теми",
        "відбувається",
        "проводиться",
        "здійснюється",
        "опрацювати матеріал",
        "попрацювати над темою",
        "виконують вправи",
        "учні працюють",
    )
    TASK_SIGNAL_TOKENS = (
        "вправа",
        "завдання",
        "питання",
        "запитання",
        "приклад",
        "кейс",
        "картка",
        "гра",
        "дослід",
        "проблемне",
        "перевір",
        "обговоріть",
        "поясніть",
        "порівняйте",
        "знайдіть",
        "складіть",
        "побудуйте",
    )
    REPORT_LABEL_PREFIXES = (
        "вчитель:",
        "учні:",
        "діяльність:",
        "оцінювання:",
        "рефлексія:",
        "мета:",
    )
    META_REPORT_MARKERS = (
        "є конкретні дії",
        "є математичний сигнал",
        "є елемент швидкої перевірки",
        "є короткий діалоговий формат",
        "є запитання на розуміння",
        "учитель коротко озвучує мету етапу",
        "учитель ставить уточнювальне запитання",
        "на етапі «",
        "у цьому етапі",
        "цей етап містить",
        "тут відбувається",
    )
    MATH_OFFTOPIC_TOKENS = (
        "каліграф",
        "буква",
        "літера",
        "іменник",
        "прикметник",
        "правопис",
        "читання",
        "вірш",
        "тексту",
        "речення",
        "звук",
    )
    MATH3_GENERIC_PHRASES = (
        "подаруйте гарний настрій",
        "подорож світом математики",
        "подорож світом чисел",
        "пам'ятаємо про правила роботи",
        "пам'ятаємо про правила безпеки та поведінки",
        "дотримуємося правил безпеки",
        "етап уроку",
        "яка вона",
        "привітайтеся з гостями",
    )

    def __init__(self, api_key: str):
        self._api_key = str(api_key or "").strip()
        self.client = genai.Client(api_key=self._api_key) if self._api_key else None
        self.loader = BaseDocsLoader()
        self.model_name = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
        self.debug_mode = os.getenv("GENERATOR_DEBUG", "").strip().lower() in {"1", "true", "yes", "on"}
        self._core_generator = None

    def _ensure_client(self):
        if self.client is not None:
            return self.client
        self._api_key = str(self._api_key or "").strip()
        if not self._api_key:
            raise ValueError(
                "Gemini API key is not configured. Set GEMINI_API_KEY "
                "(or GOOGLE_API_KEY / GOOGLE_GENAI_API_KEY)."
            )
        self.client = genai.Client(api_key=self._api_key)
        return self.client

    @staticmethod
    def _to_log_value(value):
        if isinstance(value, Path):
            return str(value)
        if isinstance(value, (list, tuple)):
            return [LessonGenerator._to_log_value(item) for item in value]
        if isinstance(value, dict):
            return {str(key): LessonGenerator._to_log_value(item) for key, item in value.items()}
        return value

    def _log_generation_event(self, request_id, event, **payload):
        record = {"request_id": request_id, "event": event}
        for key, value in payload.items():
            record[key] = self._to_log_value(value)
        logger.info("generation_trace %s", json.dumps(record, ensure_ascii=False))

    def _capture_debug_snapshot(self, debug_payload, key, value):
        if not isinstance(debug_payload, dict):
            return
        debug_payload[key] = self._to_log_value(value)

    @staticmethod
    def _utc_now_iso():
        return datetime.now(timezone.utc).isoformat()

    def _record_model_call(
        self,
        runtime_trace,
        *,
        name,
        started_at_iso,
        ended_at_iso,
        duration_ms,
        prompt,
        status,
        timeout_sec=None,
        error=None,
        extra=None,
    ):
        if not isinstance(runtime_trace, dict):
            return
        model_calls = runtime_trace.setdefault("model_calls", [])
        prompt_text = str(prompt or "")
        call_info = {
            "name": name,
            "started_at": started_at_iso,
            "ended_at": ended_at_iso,
            "duration_ms": int(duration_ms),
            "duration_sec": round(duration_ms / 1000.0, 3),
            "status": status,
            "timeout_sec": timeout_sec,
            "prompt_chars": len(prompt_text),
            "prompt_lines": prompt_text.count("\n") + 1 if prompt_text else 0,
            "error": clean_text(str(error or "")),
        }
        if isinstance(extra, dict):
            for key, value in extra.items():
                call_info[key] = value
        model_calls.append(call_info)

    @staticmethod
    def _is_retryable_model_error(exc) -> bool:
        text = str(exc or "").lower()
        retryable_markers = (
            "503",
            "unavailable",
            "429",
            "resource_exhausted",
            "rate limit",
            "rate_limit",
            "quota exceeded",
            "deadline",
            "timed out",
            "timeout",
            "connection",
            "connectionerror",
            "connection reset",
            "temporarily",
            "try again later",
        )
        non_retryable_markers = (
            "api key not valid",
            "api_key_invalid",
            "invalid api key",
            "permission_denied",
            "unauthenticated",
            "invalid_argument",
            "bad request",
            "400",
            "safety",
            "blocked",
            "policy",
        )
        if any(marker in text for marker in non_retryable_markers):
            return False
        return any(marker in text for marker in retryable_markers)

    async def _execute_model_call(
        self,
        *,
        call_name,
        prompt,
        system_instruction,
        temperature,
        runtime_trace=None,
        timeout_sec=None,
        response_mime_type="application/json",
        extra=None,
    ):
        if timeout_sec is None:
            timeout_raw = os.getenv("GENERATOR_MODEL_TIMEOUT_SEC", "35")
            try:
                timeout_sec = float(timeout_raw)
            except ValueError:
                timeout_sec = 35.0
        if timeout_sec <= 0:
            timeout_sec = None

        started_at_iso = self._utc_now_iso()
        started_perf = time.perf_counter()
        status = "ok"
        error_text = ""
        response = None
        attempts_info = []
        try:
            max_attempts = int(float(os.getenv("GENERATOR_MODEL_RETRY_ATTEMPTS", "3")))
        except ValueError:
            max_attempts = 3
        max_attempts = max(1, min(5, max_attempts))
        try:
            backoff_base_sec = float(os.getenv("GENERATOR_MODEL_RETRY_BACKOFF_SEC", "1.2"))
        except ValueError:
            backoff_base_sec = 1.2
        backoff_base_sec = max(0.1, min(10.0, backoff_base_sec))
        try:
            client = self._ensure_client()

            def call():
                return client.models.generate_content(
                    model=self.model_name,
                    contents=prompt,
                    config=genai_types.GenerateContentConfig(
                        system_instruction=system_instruction,
                        response_mime_type=response_mime_type,
                        temperature=temperature,
                    ),
                )

            for attempt in range(1, max_attempts + 1):
                attempt_started = time.perf_counter()
                try:
                    if timeout_sec:
                        response = await asyncio.wait_for(asyncio.to_thread(call), timeout=timeout_sec)
                    else:
                        response = await asyncio.to_thread(call)
                    attempts_info.append(
                        {
                            "attempt": attempt,
                            "status": "ok",
                            "duration_ms": int((time.perf_counter() - attempt_started) * 1000),
                            "error": "",
                        }
                    )
                    if attempt > 1:
                        logger.info("%s succeeded after retry attempt=%s", call_name, attempt)
                    break
                except asyncio.TimeoutError as exc:
                    attempt_error = f"model call timed out after {timeout_sec}s"
                    attempts_info.append(
                        {
                            "attempt": attempt,
                            "status": "timeout",
                            "duration_ms": int((time.perf_counter() - attempt_started) * 1000),
                            "error": attempt_error,
                            "retryable": True,
                        }
                    )
                    if attempt >= max_attempts:
                        status = "timeout"
                        error_text = attempt_error
                        logger.error("%s failed after %s attempts: %s", call_name, max_attempts, attempt_error)
                        return None
                    sleep_sec = backoff_base_sec * (2 ** (attempt - 1))
                    logger.warning(
                        "%s retryable timeout attempt=%s/%s sleep_sec=%.2f reason=%s",
                        call_name,
                        attempt,
                        max_attempts,
                        sleep_sec,
                        attempt_error,
                    )
                    await asyncio.sleep(sleep_sec)
                except Exception as exc:
                    retryable = self._is_retryable_model_error(exc)
                    attempt_error = str(exc)
                    attempts_info.append(
                        {
                            "attempt": attempt,
                            "status": "error",
                            "duration_ms": int((time.perf_counter() - attempt_started) * 1000),
                            "error": clean_text(attempt_error),
                            "retryable": retryable,
                        }
                    )
                    if not retryable:
                        status = "error"
                        error_text = attempt_error
                        logger.error("%s failed non-retryable: %s", call_name, exc)
                        return None
                    if attempt >= max_attempts:
                        status = "error"
                        error_text = attempt_error
                        logger.error("%s failed after %s attempts: %s", call_name, max_attempts, exc)
                        return None
                    sleep_sec = backoff_base_sec * (2 ** (attempt - 1))
                    logger.warning(
                        "%s retryable error attempt=%s/%s sleep_sec=%.2f reason=%s",
                        call_name,
                        attempt,
                        max_attempts,
                        sleep_sec,
                        clean_text(attempt_error)[:240],
                    )
                    await asyncio.sleep(sleep_sec)
            return response
        except Exception as exc:
            status = "error"
            error_text = str(exc)
            logger.error("%s failed: %s", call_name, exc)
            return None
        finally:
            ended_at_iso = self._utc_now_iso()
            duration_ms = int((time.perf_counter() - started_perf) * 1000)
            self._record_model_call(
                runtime_trace,
                name=call_name,
                started_at_iso=started_at_iso,
                ended_at_iso=ended_at_iso,
                duration_ms=duration_ms,
                prompt=prompt,
                status=status,
                timeout_sec=timeout_sec,
                error=error_text,
                extra={
                    **(extra if isinstance(extra, dict) else {}),
                    "retry_attempts": len(attempts_info) or 1,
                    "retry_max_attempts": max_attempts,
                    "retry_attempts_detail": attempts_info,
                    "failed_after_retries": bool(status != "ok" and len(attempts_info) >= max_attempts),
                },
            )

    def _persist_debug_payload(self, request_id, debug_payload):
        if not isinstance(debug_payload, dict) or not debug_payload:
            return ""
        try:
            diagnostics_dir = Path("storage") / "diagnostics"
            diagnostics_dir.mkdir(parents=True, exist_ok=True)
            path = diagnostics_dir / f"{request_id}_generator_debug.json"
            with path.open("w", encoding="utf-8") as handle:
                json.dump(debug_payload, handle, ensure_ascii=False, indent=2)
            return str(path)
        except Exception as exc:
            logger.warning("Failed to persist generator debug payload: %s", exc)
            return ""

    @staticmethod
    def _last_model_call(runtime_trace, call_name):
        if not isinstance(runtime_trace, dict):
            return {}
        calls = runtime_trace.get("model_calls")
        if not isinstance(calls, list):
            return {}
        for item in reversed(calls):
            if isinstance(item, dict) and item.get("name") == call_name:
                return item
        return {}






    async def generate_lesson_files(self, topic, grade, requirements, mode, subject, context="", source_file=None, source_files=None, request_id=None):
        if self._core_generator is None:
            from core.lesson_rewrite_generator import LessonRewriteGenerator

            self._core_generator = LessonRewriteGenerator(self)
        return await self._core_generator.generate_lesson_files(
            topic=topic,
            grade=grade,
            requirements=requirements,
            mode=mode,
            subject=subject,
            context=context,
            source_file=source_file,
            source_files=source_files,
            request_id=request_id,
        )

    async def _generate_lesson_json(self, topic, grade, requirements, context="", subject="", template_blueprint=None, runtime_trace=None):
        allowed_section_titles = [
            title
            for title in ((template_blueprint or {}).get("section_order") or [])
            if clean_text(title)
        ]
        allowed_optional_sections = [
            title for title in allowed_section_titles if self.OPTIONAL_SECTION_FIELD_MAP.get(title)
        ]
        allowed_optional_block = ", ".join(allowed_optional_sections) if allowed_optional_sections else "немає"
        sys_instr = (
            "Ти — методист НУШ. Заповни JSON-каркас конспекту ВИКЛЮЧНО короткими робочими тезами без жодної зайвої води. "
            "Кожен рядок — максимум одне конкретне речення або дія. Ніяких вступів, ніяких загальних фраз типа "
            "'Вчитель пропонує учням...', 'Відбувається обговорення...'. "
            "Стиль: дієслово дії + конкретна навчальна активність. Поверни ТІЛЬКИ валідний JSON."
        )
        prompt = (
            "ФОРМАТ JSON:\n"
            "{\n"
            '  "topic": "Тема уроку",\n'
            '  "grade": "Клас",\n'
            '  "subject": "Предмет",\n'
            '  "lesson_type": "Тип уроку",\n'
            '  "education_area": "Освітня галузь (за НУШ)",\n'
            '  "goal": ["Пункт мети — 1 рядок"],\n'
            '  "equipment": ["Конкретний ресурс"],\n'
            '  "lesson_flow": [\n'
            "    {\n"
            '      "stage": "Назва етапу — точно як у шаблоні",\n'
            '      "time_min": 3,\n'
            '      "teacher_actions": ["Запитання або команда вчителя — 1 коротке речення"],\n'
            '      "student_actions": ["Дія учнів — 1 коротке речення"],\n'
            '      "activities": ["Назва вправи або завдання — 1 рядок"]\n'
            "    }\n"
            "  ]\n"
            "}\n"
            "ЗАЛІЗНІ ПРАВИЛА:\n"
            "- 6–8 етапів у lesson_flow; назви ТОЧНО як у шаблонах з base_docs (не перейменовувати).\n"
            "- Кожен пункт списку — ОДНЕ коротке речення або фраза (не більше 15 слів). Ніяких абзаців.\n"
            "- Без символів Markdown (**, #, -, *, нумерації) на початку рядків.\n"
            "- ЗАБОРОНЕНІ загальні кліше: 'проводиться', 'здійснюється', 'проводить бесіду', 'звертає увагу', "
            "'пропонує учням', 'відбувається', 'використовуючи знання'. Заміни їх конкретними діями.\n"
            "- У teacher_actions — конкретна репліка або питання вчителя в дужках або дієсловом наказ. форми.\n"
            "- У student_actions — короткий опис дії учнів (що роблять, а не що 'відбувається').\n"
            "- У activities — назва конкретної вправи, гри, методу або завдання.\n"
            "- Не додавай жодних блоків, яких немає у шаблонах (Завдання, Очікувані результати тощо), "
            f"крім дозволених розділів: {allowed_optional_block}.\n"
            "- У 'Обладнання' — тільки конкретні, не банальні ресурси; без крейди, ручки, дошки.\n"
            "- Тільки українська мова.\n\n"
            f"Предмет: {subject}\nТема: {topic}\nКлас: {grade}\nВимоги: {requirements}\nКонтекст: {context}"
        )

        try:
            rich_timeout_raw = os.getenv("GENERATOR_RICH_TIMEOUT_SEC", "55")
            try:
                rich_timeout_sec = float(rich_timeout_raw)
            except ValueError:
                rich_timeout_sec = 55.0
            res = await self._execute_model_call(
                call_name="rich_fallback_json",
                prompt=prompt,
                system_instruction=sys_instr,
                temperature=0.3,
                runtime_trace=runtime_trace,
                timeout_sec=rich_timeout_sec,
                extra={
                    "subject": clean_text(subject),
                    "grade": clean_text(grade),
                    "topic_chars": len(clean_text(topic)),
                    "context_chars": len(clean_text(context)),
                },
            )
            if not res:
                raise RuntimeError("model_call_failed")
            data = json.loads(res.text)
            if not isinstance(data, dict):
                data = {}
        except Exception as e:
            logger.error(f"Помилка генерації JSON конспекту: {e}")
            data = {}

        topic_val = clean_text(data.get("topic") or topic)
        grade_val = clean_text(data.get("grade") or grade)
        subject = clean_text(data.get("subject") or data.get("course") or data.get("предмет"))
        lesson_type = clean_text(data.get("lesson_type") or data.get("type") or data.get("тип уроку"))
        education_area = clean_text(data.get("education_area") or data.get("area") or data.get("освітня галузь"))
        goal = normalize_list(data.get("goal") or data.get("goals"))
        tasks = normalize_list(data.get("tasks") or data.get("objectives"))
        expected_results = normalize_expected_results(data.get("expected_results") or data.get("results"))
        key_competencies = normalize_list(data.get("key_competencies") or data.get("competencies") or data.get("ключові компетентності"))
        cross_cutting_skills = normalize_list(data.get("cross_cutting_skills") or data.get("soft_skills") or data.get("наскрізні вміння"))
        values = normalize_list(data.get("values") or data.get("value_orientations") or data.get("ціннісні орієнтири"))
        integration = normalize_list(data.get("integration") or data.get("interdisciplinary_links") or data.get("інтеграція"))
        methods = normalize_list(data.get("methods") or data.get("methodologies") or data.get("методи"))
        forms = normalize_list(data.get("forms") or data.get("work_forms") or data.get("форми"))
        assessment = normalize_list(data.get("assessment") or data.get("evaluation") or data.get("formative_assessment"))
        differentiation = normalize_list(data.get("differentiation") or data.get("support"))
        equipment = self._prune_equipment_items(data.get("equipment"))
        resources = normalize_list(data.get("resources") or data.get("джерела") or data.get("матеріали"))
        homework = normalize_list(data.get("homework") or data.get("home_task") or data.get("домашнє завдання"))
        flow_raw = data.get("lesson_flow") if isinstance(data.get("lesson_flow"), list) else []

        lesson_flow = []
        for item in flow_raw:
            if not isinstance(item, dict):
                continue
            stage = clean_text(item.get("stage") or item.get("title") or "")
            time_min = item.get("time_min") or item.get("time") or item.get("duration")
            teacher_actions = normalize_list(item.get("teacher_actions") or item.get("teacher") or item.get("teacher_steps"))
            student_actions = normalize_list(item.get("student_actions") or item.get("students") or item.get("student_steps"))
            activities = normalize_list(item.get("activities") or item.get("tasks") or item.get("practice") or item.get("bullets"))

            if not stage and any([teacher_actions, student_actions, activities]):
                stage = "Етап уроку"

            if stage or any([teacher_actions, student_actions, activities]):
                lesson_flow.append(self._simplify_rich_stage({
                    "stage": stage,
                    "time_min": time_min,
                    "teacher_actions": teacher_actions,
                    "student_actions": student_actions,
                    "activities": activities,
                }, equipment=equipment, resources=resources))

        if not lesson_flow:
            lesson_flow = [
                {
                    "stage": "Організаційний момент",
                    "time_min": 2,
                    "teacher_actions": ["Вітає дітей, організовує увагу класу.", "Коротко налаштовує на роботу."],
                    "student_actions": ["Вітаються, налаштовуються на роботу"],
                    "activities": ["Коротка організаційна вправа"],
                },
                {
                    "stage": "Мотивація навчальної діяльності",
                    "time_min": 3,
                    "teacher_actions": ["Ставить коротке проблемне запитання, вводить у тему."],
                    "student_actions": ["Висловлюють припущення, формують очікування"],
                    "activities": ["Міні-дискусія"],
                },
                {
                    "stage": "Актуалізація опорних знань",
                    "time_min": 5,
                    "teacher_actions": ["Організовує коротке фронтальне опитування, пропонує вправу на повторення."],
                    "student_actions": ["Відповідають, виконують короткі завдання"],
                    "activities": ["Вправа на повторення"],
                },
                {
                    "stage": "Опрацювання нового матеріалу",
                    "time_min": 15,
                    "teacher_actions": ["Пояснює новий матеріал, демонструє приклади та ставить уточнювальні запитання."],
                    "student_actions": ["Слухають, ставлять запитання, фіксують ключові ідеї"],
                    "activities": ["Робота з прикладами, мікрообговорення"],
                },
                {
                    "stage": "Практична діяльність",
                    "time_min": 12,
                    "teacher_actions": ["Пояснює завдання, супроводжує виконання, за потреби коригує відповіді."],
                    "student_actions": ["Виконують вправи індивідуально/у парах"],
                    "activities": ["Практичні завдання різного рівня"],
                },
                {
                    "stage": "Узагальнення і систематизація",
                    "time_min": 5,
                    "teacher_actions": ["Підводить дітей до висновків, уточнює головні результати уроку."],
                    "student_actions": ["Формулюють підсумки"],
                    "activities": ["Складання короткої схеми/пам'ятки"],
                },
                {
                    "stage": "Рефлексія",
                    "time_min": 3,
                    "teacher_actions": ["Організовує коротку рефлексивну вправу, ставить підсумкові запитання."],
                    "student_actions": ["Висловлюють, що вдалося/що потребує повторення"],
                    "activities": ["Рефлексивна вправа"],
                },
                {
                    "stage": "Домашнє завдання",
                    "time_min": 2,
                    "teacher_actions": ["Пояснює домашнє завдання, уточнює спосіб виконання."],
                    "student_actions": ["Уточнюють умови виконання"],
                    "activities": ["Запис домашнього завдання"],
                }
            ]

        return {
            "topic": topic_val,
            "grade": grade_val,
            "subject": subject,
            "lesson_type": lesson_type,
            "education_area": education_area,
            "goal": goal,
            "tasks": tasks,
            "expected_results": expected_results,
            "key_competencies": key_competencies,
            "cross_cutting_skills": cross_cutting_skills,
            "values": values,
            "integration": integration,
            "methods": methods,
            "forms": forms,
            "assessment": assessment,
            "differentiation": differentiation,
            "equipment": equipment,
            "resources": resources,
            "homework": homework,
            "lesson_flow": lesson_flow
        }







    @staticmethod
    def _extract_json_payload(text):
        raw = str(text or '').strip()
        if not raw:
            return {}
        try:
            data = json.loads(raw)
            return data if isinstance(data, dict) else {}
        except Exception:
            start = raw.find('{')
            end = raw.rfind('}')
            if start == -1 or end == -1 or end <= start:
                return {}
            try:
                data = json.loads(raw[start:end + 1])
                return data if isinstance(data, dict) else {}
            except Exception:
                return {}

    @staticmethod
    def _get_header_value(parsed_doc, label):
        for item in parsed_doc.get('header_fields', []):
            if item.get('label') == label:
                return clean_text(item.get('value') or '')
        return ''

    @staticmethod
    def _topic_tokens(text):
        stopwords = {
            'та', 'і', 'й', 'у', 'в', 'на', 'до', 'для', 'про', 'за', 'з', 'із', 'зі', 'по', 'що', 'як', 'це',
            'урок', 'тема', 'клас', 'число', 'числа', 'чисел', 'частина', 'робота', 'вправи', 'вправа',
        }
        normalized = re.sub(r"[^\w']+", ' ', clean_text(text).lower(), flags=re.UNICODE)
        tokens = []
        for token in normalized.split():
            if len(token) < 3 or token in stopwords:
                continue
            tokens.append(token)
        return set(tokens)

    def _score_reference_doc(self, requested_topic, parsed_doc, path):
        reference_topic = self._get_header_value(parsed_doc, 'Тема') or clean_text(path.stem)
        requested_norm = self._normalize_template_name(requested_topic)
        reference_norm = self._normalize_template_name(reference_topic)
        requested_tokens = self._topic_tokens(requested_topic)
        reference_tokens = self._topic_tokens(reference_topic)
        overlap = len(requested_tokens & reference_tokens)
        ratio = SequenceMatcher(None, requested_norm, reference_norm).ratio() if requested_norm and reference_norm else 0.0
        substring_bonus = 2.0 if requested_norm and reference_norm and (requested_norm in reference_norm or reference_norm in requested_norm) else 0.0
        header_bonus = min(2.0, len(parsed_doc.get('header_fields', [])) * 0.25)
        structure_bonus = min(2.0, len(parsed_doc.get('sections', [])) * 0.25)
        master_bonus = 1.0 if self.loader._is_master_template(path) else 0.0
        score = overlap * 10.0 + ratio * 8.0 + substring_bonus + header_bonus + structure_bonus + master_bonus
        return score

    def _select_reference_doc(self, topic, doc_paths, parsed_docs):
        if not doc_paths or not parsed_docs:
            return None, None

        pairs = list(zip(doc_paths, parsed_docs))
        scored = sorted(
            ((self._score_reference_doc(topic, parsed_doc, path), path, parsed_doc) for path, parsed_doc in pairs),
            key=lambda item: item[0],
            reverse=True,
        )
        best_score, best_path, best_doc = scored[0]
        if best_score > 0:
            return best_path, best_doc

        for path, parsed_doc in pairs:
            if self.loader._is_master_template(path):
                return path, parsed_doc
        return pairs[0]

    @staticmethod
    def _count_structure_stages(structure):
        if not isinstance(structure, dict):
            return 0
        explicit = structure.get("stages")
        if isinstance(explicit, list) and explicit:
            return len(explicit)
        total = 0
        for section in structure.get("sections") or []:
            total += len((section or {}).get("stages") or [])
        return total

    @staticmethod
    def _score_structure_completeness(structure):
        if not isinstance(structure, dict):
            return (0, 0, 0)
        return (
            LessonGenerator._count_structure_stages(structure),
            len(structure.get("sections") or []),
            len(structure.get("header_fields") or []),
        )

    def _pick_richest_structure(self, parsed_docs):
        candidates = [doc for doc in (parsed_docs or []) if isinstance(doc, dict)]
        if not candidates:
            return None
        return max(candidates, key=self._score_structure_completeness)

    def _merge_reference_structure_stages(self, primary_structure, parsed_docs, max_stages=10):
        if not isinstance(primary_structure, dict):
            return None

        merged = copy.deepcopy(primary_structure)
        sections = merged.get("sections")
        if not isinstance(sections, list):
            sections = []
            merged["sections"] = sections

        flow_index = None
        best_stage_count = -1
        for index, section in enumerate(sections):
            section_title = clean_text_preserve_prefix(section.get("display_title") or section.get("title") or "")
            normalized = self._normalize_template_name(section_title)
            stage_count = len((section or {}).get("stages") or [])
            if "хід уроку" in normalized:
                flow_index = index
                break
            if stage_count > best_stage_count:
                best_stage_count = stage_count
                flow_index = index

        if flow_index is None:
            sections.append(
                {
                    "title": "Хід уроку",
                    "display_title": "Хід уроку",
                    "style": "Normal",
                    "items": [],
                    "item_styles": [],
                    "substeps": [],
                    "stages": [],
                    "children_order": [],
                }
            )
            flow_index = len(sections) - 1

        flow_section = sections[flow_index]
        flow_title = clean_text_preserve_prefix(flow_section.get("display_title") or flow_section.get("title") or "") or "Хід уроку"
        flow_section.setdefault("stages", [])
        flow_section.setdefault("children_order", [])

        existing_keys = set()
        for stage in flow_section.get("stages") or []:
            title = clean_text_preserve_prefix(stage.get("display_title") or stage.get("title") or "")
            key = self._normalize_template_name(title) or clean_text(title).lower()
            if key:
                existing_keys.add(key)

        for parsed_doc in parsed_docs or []:
            for stage in parsed_doc.get("stages") or []:
                if len(flow_section.get("stages") or []) >= max_stages:
                    break
                stage_title = clean_text_preserve_prefix(stage.get("display_title") or stage.get("title") or "")
                stage_key = self._normalize_template_name(stage_title) or clean_text(stage_title).lower()
                if not stage_title or not stage_key or stage_key in existing_keys:
                    continue
                stage_copy = copy.deepcopy(stage)
                stage_copy["section_title"] = flow_title
                flow_section["stages"].append(stage_copy)
                flow_section["children_order"].append(("stage", len(flow_section["stages"]) - 1))
                existing_keys.add(stage_key)

        merged_stages = []
        for section in sections:
            for stage in section.get("stages") or []:
                merged_stages.append(stage)
        merged["stages"] = merged_stages
        return merged

    def _select_blueprint_reference_structure(self, reference_structure, parsed_docs, min_stages=4):
        parsed_docs = [doc for doc in (parsed_docs or []) if isinstance(doc, dict)]
        if not reference_structure and not parsed_docs:
            return None, {"mode": "none", "reference_stages": 0, "richest_stages": 0, "selected_stages": 0}

        reference_stages = self._count_structure_stages(reference_structure)
        richest_structure = self._pick_richest_structure(parsed_docs)
        richest_stages = self._count_structure_stages(richest_structure)

        selected = copy.deepcopy(reference_structure) if isinstance(reference_structure, dict) else None
        mode = "reference"
        if selected is None and richest_structure is not None:
            selected = copy.deepcopy(richest_structure)
            mode = "richest"

        if selected is not None and self._count_structure_stages(selected) < min_stages and richest_structure is not None:
            if richest_stages > self._count_structure_stages(selected):
                selected = copy.deepcopy(richest_structure)
                mode = "richest"

        merged = self._merge_reference_structure_stages(selected, parsed_docs, max_stages=max(min_stages + 2, 10))
        if self._count_structure_stages(merged) > self._count_structure_stages(selected):
            selected = merged
            mode = f"{mode}+merged"

        return selected, {
            "mode": mode,
            "reference_stages": reference_stages,
            "richest_stages": richest_stages,
            "selected_stages": self._count_structure_stages(selected),
        }

    @staticmethod
    def _detect_stage_collapse_step(reference_stages, blueprint_stages, stages_after_apply, stages_after_compact, final_stages):
        reference_stages = int(reference_stages or 0)
        blueprint_stages = int(blueprint_stages or 0)
        stages_after_apply = int(stages_after_apply or 0)
        stages_after_compact = int(stages_after_compact or 0)
        final_stages = int(final_stages or 0)

        if reference_stages > 1 and blueprint_stages <= 1:
            return "blueprint_assembly"
        if blueprint_stages > 1 and stages_after_apply <= 1:
            return "apply_generated_content"
        if stages_after_apply > 1 and stages_after_compact <= 1:
            return "compact_blueprint_document"
        if stages_after_compact > 1 and final_stages <= 1:
            return "finalize_or_render"
        return None

    @staticmethod
    def _truncate_example(text, limit=220):
        text = clean_text(text)
        if len(text) <= limit:
            return text
        return f"{text[:limit - 1].rstrip()}…"

    @staticmethod
    def _split_inline_title_payload(title):
        normalized = clean_text_preserve_prefix(title)
        if not normalized or ":" not in normalized:
            return normalized, []

        left, right = normalized.split(":", 1)
        left = clean_text_preserve_prefix(left)
        right = clean_text_preserve_prefix(right)
        if not left or not right:
            return normalized, []

        left_words = len(left.split())
        right_words = len(right.split())
        if left_words > 7 or right_words < 3:
            return normalized, []

        fragments = [
            clean_text_preserve_prefix(fragment)
            for fragment in re.split(r"[;•]+", right)
            if clean_text_preserve_prefix(fragment)
        ]
        if not fragments:
            fragments = [right]

        return left, fragments[:4]

    @staticmethod
    def _extract_stage_prefix(title):
        normalized = clean_text_preserve_prefix(title)
        match = re.match(r"^\s*((?:\d+|[IVXІХ]+)[\.\)])\s*", normalized, flags=re.IGNORECASE)
        return clean_text_preserve_prefix(match.group(1)) if match else ""

    @staticmethod
    def _contains_foreign_person_name(text, topic=""):
        value = clean_text_preserve_prefix(text)
        if not value:
            return False
        topic_text = clean_text(topic).lower()

        word_pattern = r"[A-Za-zА-Яа-яІіЇїЄєҐґ][A-Za-zА-Яа-яІіЇїЄєҐґ''\-]*"
        capitalized_word_pattern = r"[A-ZА-ЯІЇЄҐ][a-zа-яіїєґ''\-]+"
        topic_tokens = {
            token.lower()
            for token in re.findall(word_pattern, topic_text)
            if len(token) >= 3
        }

        candidates = []
        candidates.extend(
            re.findall(
                rf"(?:\b[A-ZА-ЯІЇЄҐ]\.\s*){{1,3}}{word_pattern}",
                value,
            )
        )
        candidates.extend(
            re.findall(
                rf"\b{capitalized_word_pattern}\s+{capitalized_word_pattern}\b",
                value,
            )
        )

        for candidate in candidates:
            candidate_norm = clean_text(candidate).lower()
            if topic_text and candidate_norm and candidate_norm in topic_text:
                continue
            candidate_tokens = {
                token.lower()
                for token in re.findall(word_pattern, candidate_norm)
                if len(token) >= 3
            }
            if topic_tokens and candidate_tokens & topic_tokens:
                continue
            return True
        return False

    def _derive_safe_stage_title(self, original_title):
        normalized = clean_text_preserve_prefix(original_title)
        lower = clean_text(normalized).lower()
        prefix = self._extract_stage_prefix(normalized)
        if "домаш" in lower:
            core = "Домашнє завдання"
        elif "підсум" in lower or "рефлекс" in lower:
            core = "Підсумок уроку"
        elif "закріп" in lower:
            core = "Закріплення знань"
        elif "вивчен" in lower or "нов" in lower or "основн" in lower:
            core = "Вивчення нового матеріалу"
        elif "орган" in lower or "вступ" in lower:
            core = "Організаційний момент"
        elif "мотивац" in lower:
            core = "Мотивація навчальної діяльності"
        elif "актуал" in lower:
            core = "Актуалізація опорних знань"
        else:
            core = "Етап уроку"

        if prefix:
            return f"{prefix} {core}"
        return core

    def _sanitize_substep_title_phrase(self, title, topic=""):
        normalized = clean_text_preserve_prefix(title)
        if not normalized:
            return ""

        topic_tokens = self._topic_tokens(topic)
        has_kazka_in_topic = any(token.startswith("казк") for token in topic_tokens)
        if has_kazka_in_topic:
            return normalized

        replacements = {
            "казки": "твору",
            "казку": "твір",
            "казка": "твір",
            "казці": "творі",
            "казкою": "твором",
        }
        updated = normalized
        for source, target in replacements.items():
            updated = re.sub(rf"\b{source}\b", target, updated, flags=re.IGNORECASE)
        return clean_text_preserve_prefix(updated)

    @staticmethod
    def _normalize_stage_title_key(title):
        value = clean_text_preserve_prefix(title).lower()
        value = re.sub(r"\s+", " ", value).strip()
        return value

    @staticmethod
    def _is_content_like_substep_title(title):
        normalized = clean_text_preserve_prefix(title)
        if not normalized:
            return False

        lower = normalized.lower()
        words = normalized.split()
        words_count = len(words)
        short_heading_prefixes = (
            "гра ",
            "вправа ",
            "завдання ",
            "робота в парах",
            "робота в групах",
            "фізкультхвилинка",
            "рефлексія",
        )
        short_heading_hit = any(lower.startswith(prefix) for prefix in short_heading_prefixes)
        if short_heading_hit and words_count <= 4 and not any(mark in normalized for mark in (":", ";", "(", ")", ",")):
            return False

        if words_count >= 8:
            return True
        if re.match(r"^\d+\.\s+", normalized):
            if words_count >= 4:
                return True
            if any(mark in normalized for mark in (":", ";", ".", ",", "(", ")")):
                return True
        if normalized.count("-") >= 3:
            return True
        if ":" in normalized and words_count >= 4:
            return True
        if ";" in normalized and words_count >= 4:
            return True
        if "." in normalized and words_count >= 3:
            return True
        if "," in normalized and words_count >= 6:
            return True
        if re.search(r"[!?][\"'Р'В»\)\]]*\s*$", normalized):
            return words_count >= 4
        if "(" in normalized and ")" in normalized and words_count >= 4:
            return True
        if "Р'В«" in normalized and "Р'В»" in normalized and words_count >= 6:
            return True
        return False

    def _normalize_example_node_title(self, kind, title, topic=""):
        normalized = clean_text_preserve_prefix(title)
        if not normalized:
            return "", []

        inline_items = []
        if kind in {"section", "stage"}:
            normalized, inline_items = self._split_inline_title_payload(normalized)

        if kind == "substep":
            normalized = self._sanitize_substep_title_phrase(normalized, topic=topic)
            if self._is_content_like_substep_title(normalized) or self._contains_foreign_person_name(normalized, topic):
                inline_items = [normalized, *inline_items]
                normalized = ""
        elif kind == "stage":
            if self._is_content_like_substep_title(normalized) or self._contains_foreign_person_name(normalized, topic):
                inline_items = [normalized, *inline_items]
                normalized = self._derive_safe_stage_title(normalized)

        return normalized, inline_items

    def _build_example_node(self, node, kind="section", topic=""):
        raw_sample_items = [
            item
            for item in normalize_block_list(node.get("items") or [])
            if not self._is_metadata_source_line(item)
        ]
        sample_item_styles = [clean_text(style) or "Normal" for style in (node.get("item_styles") or [])][: len(raw_sample_items)]
        raw_title = clean_text_preserve_prefix(node.get("title") or "")
        raw_display_title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
        normalized_title, title_inline_items = self._normalize_example_node_title(kind, raw_title, topic=topic)
        normalized_display_title, display_inline_items = self._normalize_example_node_title(kind, raw_display_title, topic=topic)

        if not normalized_display_title and normalized_title:
            normalized_display_title = normalized_title
        if not normalized_title and normalized_display_title:
            normalized_title = normalized_display_title

        seed_items = self._merge_sample_items([], [display_inline_items, title_inline_items], limit=6)
        seed_items = [
            item for item in seed_items
            if not self._is_metadata_source_line(item)
        ]
        sample_items = self._merge_sample_items(seed_items, [raw_sample_items], limit=18)
        content_style = next((style for style in sample_item_styles if style), "Normal")
        built_substeps = [self._build_example_node(substep, "substep", topic=topic) for substep in (node.get("substeps") or [])]
        substep_index_map = {}
        normalized_substeps = []
        untitled_stage_substeps = 0
        for index, substep in enumerate(built_substeps):
            substep_title = clean_text_preserve_prefix(substep.get("display_title") or substep.get("title") or "")
            substep_has_payload = bool(
                (substep.get("sample_items") or [])
                or (substep.get("substeps") or [])
                or (substep.get("stages") or [])
            )
            if not substep_title and not substep_has_payload:
                continue
            if kind == "stage":
                if not substep_title:
                    if untitled_stage_substeps >= 1:
                        continue
                    untitled_stage_substeps += 1
                if len(normalized_substeps) >= 5:
                    continue
            elif kind == "section":
                if not substep_title:
                    continue
                if self._contains_foreign_person_name(substep_title, topic):
                    continue
                if self._is_content_like_substep_title(substep_title) and len(substep_title.split()) >= 4:
                    continue
                if len(normalized_substeps) >= 3:
                    continue
            substep_index_map[index] = len(normalized_substeps)
            normalized_substeps.append(substep)

        normalized_children_order = []
        for child_kind, child_index in (node.get("children_order") or []):
            if child_kind == "substep":
                mapped_index = substep_index_map.get(child_index)
                if mapped_index is None:
                    continue
                normalized_children_order.append(("substep", mapped_index))
            else:
                normalized_children_order.append((child_kind, child_index))

        result = {
            "title": normalized_title,
            "display_title": normalized_display_title,
            "style": clean_text(node.get("style") or "Normal") or "Normal",
            "content_style": content_style,
            "sample_items": sample_items,
            "sample_item_styles": sample_item_styles,
            "source_items": [],
            "items": [],
            "substeps": normalized_substeps,
            "children_order": normalized_children_order,
        }
        if "stages" in node:
            raw_stages = [self._build_example_node(stage, "stage", topic=topic) for stage in (node.get("stages") or [])]
            deduped_stages = []
            stage_index_map = {}
            seen_stage_keys = {}
            for index, stage in enumerate(raw_stages):
                stage_title = clean_text_preserve_prefix(stage.get("display_title") or stage.get("title") or "")
                stage_key = self._normalize_stage_title_key(stage_title) or stage_title.lower()
                payload_size = (
                    len(stage.get("sample_items") or [])
                    + len(stage.get("substeps") or [])
                    + len(stage.get("items") or [])
                )
                if stage_key and stage_key in seen_stage_keys:
                    existing_idx, existing_payload = seen_stage_keys[stage_key]
                    if payload_size > existing_payload:
                        deduped_stages[existing_idx] = stage
                        seen_stage_keys[stage_key] = (existing_idx, payload_size)
                    stage_index_map[index] = existing_idx
                    continue

                stage_index_map[index] = len(deduped_stages)
                deduped_stages.append(stage)
                if stage_key:
                    seen_stage_keys[stage_key] = (len(deduped_stages) - 1, payload_size)

            result["stages"] = deduped_stages
            normalized_children_order = []
            for child_kind, child_index in (result.get("children_order") or []):
                if child_kind == "stage":
                    mapped_stage_index = stage_index_map.get(child_index)
                    if mapped_stage_index is None:
                        continue
                    normalized_children_order.append(("stage", mapped_stage_index))
                else:
                    normalized_children_order.append((child_kind, child_index))
            result["children_order"] = normalized_children_order

        deduped_children_order = []
        seen_children = set()
        for child_kind, child_index in (result.get("children_order") or []):
            key = (child_kind, child_index)
            if key in seen_children:
                continue
            seen_children.add(key)
            deduped_children_order.append((child_kind, child_index))
        result["children_order"] = deduped_children_order
        return result

    def _build_example_blueprint(self, reference_structure, topic, grade, subject, parsed_docs=None):
        blueprint = {
            'mode': 'strict_example',
            'reference_file': clean_text(reference_structure.get('file_name') or ''),
            'topic': clean_text(topic),
            'grade': clean_text(grade),
            'subject': clean_text(subject),
            'header_fields': [],
            'sections': [],
        }

        for item in reference_structure.get('header_fields', []):
            label = clean_text(item.get('label') or '')
            if not label:
                continue
            sample_value = clean_text(item.get('value') or '')
            default_value = ''
            if label == 'Тема':
                default_value = clean_text(topic)
            elif label == 'Клас':
                default_value = clean_text(grade)
            elif label == 'Предмет':
                default_value = clean_text(subject)
            elif label == 'Дата':
                default_value = ''
            blueprint['header_fields'].append({
                'label': label,
                'style': clean_text(item.get('style') or 'Normal') or 'Normal',
                'sample_value': sample_value,
                'value': default_value,
            })

        for section in reference_structure.get('sections', []):
            blueprint['sections'].append(self._build_example_node(section, topic=topic))

        if parsed_docs:
            example_index, grouped_stage_examples = self._collect_example_index(parsed_docs)
            for section in blueprint['sections']:
                self._enrich_blueprint_node(section, "section", example_index, grouped_stage_examples)

        return blueprint

    @staticmethod
    def _build_node_style_profile(node):
        sample_items = normalize_block_list((node or {}).get("sample_items") or [])
        has_children = bool((node or {}).get("stages") or (node or {}).get("substeps"))
        if not sample_items:
            return {
                "target_words_min": 6,
                "target_words_max": 26,
                "target_items_min": 1 if has_children else 3,
                "target_items_max": 3 if has_children else 6,
            }

        word_counts = []
        has_question = False
        for item in sample_items[:14]:
            text = clean_text(item)
            if not text:
                continue
            word_counts.append(len(text.split()))
            if "?" in text:
                has_question = True

        if not word_counts:
            return {
                "target_words_min": 6,
                "target_words_max": 26,
                "target_items_min": 1 if has_children else 3,
                "target_items_max": 3 if has_children else 6,
            }

        avg_words = sum(word_counts) / len(word_counts)
        min_words = max(3, int(round(avg_words * 0.65)))
        max_words = max(min_words + 6, min(48, int(round(avg_words * 1.85))))
        return {
            "target_words_min": min_words,
            "target_words_max": max_words,
            "allow_questions": has_question,
            "target_items_min": 1 if has_children else 3,
            "target_items_max": 3 if has_children else 6,
        }

    def _collect_prompt_source_items(self, node, limit=8):
        node = node if isinstance(node, dict) else {}
        buckets = [
            node.get("source_items") or [],
            node.get("sample_items") or [],
        ]
        for substep in (node.get("substeps") or [])[:6]:
            buckets.append(substep.get("source_items") or [])
        for stage in (node.get("stages") or [])[:8]:
            buckets.append(stage.get("source_items") or [])
        merged = self._merge_sample_items([], buckets, limit=max(limit, 24))
        return [self._truncate_example(item, 420) for item in merged[:limit]]

    def _build_prompt_node(self, node, kind="section"):
        source_limit_map = {
            "section": 2,
            "stage": 4,
            "substep": 0,
        }
        source_limit = source_limit_map.get(kind, 5)
        payload = {
            "display_title": node.get("display_title") or node.get("title") or "",
            "style_profile": self._build_node_style_profile(node),
            "substeps": [self._build_prompt_node(substep, kind="substep") for substep in (node.get("substeps") or [])],
        }
        if source_limit > 0:
            payload["source_items"] = self._collect_prompt_source_items(node, limit=source_limit)
        if "stages" in node:
            payload["stages"] = [self._build_prompt_node(stage, kind="stage") for stage in (node.get("stages") or [])]
        return payload

    def _iter_reference_nodes(self, node, kind):
        yield kind, node
        for substep in node.get("substeps") or []:
            yield from self._iter_reference_nodes(substep, "substep")
        for stage in node.get("stages") or []:
            yield from self._iter_reference_nodes(stage, "stage")

    def _collect_example_index(self, parsed_docs):
        exact_index = {}
        stage_group_index = {"intro": [], "main": [], "closing": []}

        for parsed_doc in parsed_docs or []:
            for section in parsed_doc.get("sections", []):
                for kind, node in self._iter_reference_nodes(section, "section"):
                    title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
                    if not title:
                        continue
                    items = normalize_block_list(node.get("items") or [])
                    if items:
                        key = (kind, self._normalize_template_name(title))
                        exact_index.setdefault(key, []).append(items)
                        if kind == "stage":
                            stage_group_index.setdefault(self._classify_stage_name(title), []).append(items)

        return exact_index, stage_group_index

    def _merge_sample_items(self, base_items, example_sets, limit=18):
        merged = []
        seen = set()
        for value in normalize_block_list(base_items):
            normalized = clean_text_preserve_prefix(value)
            if normalized and normalized not in seen:
                merged.append(normalized)
                seen.add(normalized)

        for example_items in example_sets:
            for value in normalize_block_list(example_items):
                normalized = clean_text_preserve_prefix(value)
                if not normalized or normalized in seen:
                    continue
                merged.append(normalized)
                seen.add(normalized)
                if len(merged) >= limit:
                    return merged
        return merged

    def _enrich_blueprint_node(self, node, kind, example_index, stage_group_index):
        title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
        if title:
            key = (kind, self._normalize_template_name(title))
            example_sets = list(example_index.get(key, []))
            if kind == "stage" and len(example_sets) < 2:
                example_sets.extend(stage_group_index.get(self._classify_stage_name(title), []))
            if example_sets:
                node["sample_items"] = self._merge_sample_items(node.get("sample_items") or [], example_sets)

        for substep in node.get("substeps") or []:
            self._enrich_blueprint_node(substep, "substep", example_index, stage_group_index)
        for stage in node.get("stages") or []:
            self._enrich_blueprint_node(stage, "stage", example_index, stage_group_index)

    def _build_reference_context(self, reference_structure, doc_paths, parsed_docs=None):
        if not reference_structure:
            return ""

        reference_topic = self._get_header_value(reference_structure, 'Тема')
        reference_goal = self._get_header_value(reference_structure, 'Мета')
        lines = [
            f"Референсний конспект: {reference_structure.get('file_name', 'невідомий файл')}",
        ]
        if reference_topic:
            lines.append(f"Тема прикладу: {reference_topic}")
        if reference_goal:
            lines.append(f"Мета прикладу: {self._truncate_example(reference_goal, 280)}")
        if doc_paths:
            lines.append(f"Доступні файли цього предмета/класу: {', '.join(path.name for path in doc_paths[:12])}")
        if parsed_docs:
            lines.append(f"Орієнтуйся не лише на референсний файл, а на всі {len(parsed_docs)} конспекти з цієї папки.")
            for parsed_doc in parsed_docs[:12]:
                doc_topic = self._get_header_value(parsed_doc, 'Тема') or clean_text(parsed_doc.get('file_name') or '')
                section_titles = [clean_text(section.get('display_title') or section.get('title') or '') for section in parsed_doc.get('sections', [])[:4]]
                stage_titles = [clean_text(stage.get('display_title') or stage.get('title') or '') for stage in parsed_doc.get('stages', [])[:4]]
                summary_parts = []
                if doc_topic:
                    summary_parts.append(f"тема: {doc_topic}")
                if section_titles:
                    summary_parts.append(f"розділи: {', '.join(section_titles)}")
                if stage_titles:
                    summary_parts.append(f"етапи: {', '.join(stage_titles)}")
                if summary_parts:
                    lines.append(f"- {parsed_doc.get('file_name', 'файл')}: {'; '.join(summary_parts)}")
        return '\n'.join(lines)

    def _collect_blueprint_stages(self, blueprint):
        titles = []
        if not blueprint:
            return titles

        for section in blueprint.get("sections", []):
            for stage in section.get("stages") or []:
                title = clean_text_preserve_prefix(stage.get("display_title") or stage.get("title") or "")
                if title:
                    titles.append(title)
        return titles

    @staticmethod
    def _format_slide_numbers(slide_numbers):
        normalized = sorted({int(number) for number in slide_numbers if str(number).isdigit()})
        if not normalized:
            return ""
        if len(normalized) == 1:
            return f"слайд {normalized[0]}"
        if len(normalized) == 2 and normalized[1] == normalized[0] + 1:
            return f"слайди {normalized[0]}–{normalized[1]}"
        return "слайди " + ", ".join(str(number) for number in normalized)

    def _assign_slides_to_stage_titles(self, slides, blueprint):
        if not slides or not blueprint:
            return {}

        slide_numbers = [
            slide.get("number")
            for slide in slides
            if not slide.get("is_title_slide") and (
                slide.get("text_blocks") or slide.get("bullet_points") or slide.get("notes") or slide.get("images")
            )
        ]
        stage_titles = self._collect_blueprint_stages(blueprint)
        if not slide_numbers or not stage_titles:
            return {}

        intro_titles = [title for title in stage_titles if self._classify_stage_name(title) == "intro"]
        closing_titles = [title for title in stage_titles if self._classify_stage_name(title) == "closing"]
        main_titles = [title for title in stage_titles if self._classify_stage_name(title) == "main"]
        assignments = {}

        first_slide = slide_numbers[0]
        last_slide = slide_numbers[-1]
        for title in intro_titles:
            assignments[title] = [first_slide]
        for title in closing_titles:
            assignments[title] = [last_slide]

        start_index = 1 if intro_titles and len(slide_numbers) > 1 else 0
        end_index = len(slide_numbers) - 1 if closing_titles and len(slide_numbers) - start_index > 1 else len(slide_numbers)
        middle_slides = slide_numbers[start_index:end_index]

        target_titles = main_titles or [title for title in stage_titles if title not in intro_titles + closing_titles]
        if not target_titles:
            target_titles = stage_titles

        if middle_slides and target_titles:
            chunk_size = max(1, (len(middle_slides) + len(target_titles) - 1) // len(target_titles))
            for index, title in enumerate(target_titles):
                chunk = middle_slides[index * chunk_size:(index + 1) * chunk_size]
                if not chunk:
                    chunk = [middle_slides[-1]]
                assignments[title] = chunk
        elif target_titles:
            for title in target_titles:
                assignments.setdefault(title, [first_slide])

        return assignments

    def _summarize_slide_for_source(self, slide, limit=320):
        parts = [f"Слайд {slide.get('number')}."]
        title = clean_text(slide.get("title") or "")
        if title:
            parts.append(f"Назва: {title}.")

        text_blocks = normalize_block_list(slide.get("text_blocks"))[:3]
        bullet_points = normalize_block_list(slide.get("bullet_points"))[:4]
        notes = clean_text(slide.get("notes") or "")
        image_labels = normalize_block_list(slide.get("image_labels"))[:2]
        images_count = int(slide.get("images") or 0)

        if text_blocks:
            parts.append("Основне: " + "; ".join(text_blocks) + ".")
        if bullet_points:
            parts.append("Ключові пункти: " + "; ".join(bullet_points) + ".")
        if notes:
            parts.append("Нотатки: " + self._truncate_example(notes, 180) + ".")
        if image_labels:
            parts.append("Візуальні опори: " + "; ".join(image_labels) + ".")
        elif images_count:
            parts.append(f"На слайді є {images_count} візуальних матеріали.")

        return self._truncate_example(" ".join(part for part in parts if part), limit)

    @staticmethod
    def _is_template_noise_line(text):
        normalized = clean_text(text).lower()
        if not normalized:
            return True
        if normalized.startswith((
            "робота над розділом",
            "опрацювання розділу",
            "обговорення теми",
            "робота з матеріалом",
            "обговорюють приклад із теми",
        )):
            return True
        if "робота над розділом" in normalized and "за темою" in normalized:
            return True
        if "обговорюють приклад із теми" in normalized and "на етапі" in normalized:
            return True
        return False

    @staticmethod
    def _looks_like_structural_dump(text):
        normalized = clean_text(text).lower()
        if not normalized:
            return True

        marker_score = 0
        for marker in (
            "хід уроку",
            "очікувані результати",
            "С–.",
            "іі.",
            "ііі.",
            "iv.",
            "v.",
            "vi.",
        ):
            if marker in normalized:
                marker_score += 1
        if marker_score >= 3 and len(normalized) > 120:
            return True

        roman_stage_hits = len(re.findall(r"\b(?:[ivx]{1,4})\.", normalized))
        if roman_stage_hits >= 3:
            return True
        if re.match(r"^(?:[ivx]{1,4}|\d+)\.\s+", normalized):
            if len(normalized.split()) <= 9:
                return True

        if normalized.count("Р'В«") >= 3 and len(normalized) > 180:
            return True
        return False

    @staticmethod
    def _is_metadata_source_line(text):
        normalized = clean_text(text).lower()
        if not normalized:
            return True
        if LessonGenerator._is_template_noise_line(normalized):
            return True
        metadata_prefixes = (
            "тема:",
            "мета:",
            "клас:",
            "предмет:",
            "дата:",
            "обладнання:",
            "тип уроку:",
            "освітня галузь:",
        )
        return normalized.startswith(metadata_prefixes)

    def _sanitize_source_items(self, items, limit=20):
        sanitized = []
        seen = set()

        for raw_item in normalize_block_list(items):
            candidate = clean_text_preserve_prefix(raw_item)
            if not candidate:
                continue
            fragments = [candidate]
            if len(candidate) > 420:
                fragments = [
                    clean_text_preserve_prefix(part)
                    for part in re.split(r"(?<=[\.\!\?;:])\s+", candidate)
                    if clean_text_preserve_prefix(part)
                ] or [candidate]

            for fragment in fragments:
                if (
                    not fragment
                    or self._is_metadata_source_line(fragment)
                    or self._looks_like_structural_dump(fragment)
                ):
                    continue
                normalized = fragment.lower()
                if normalized in seen:
                    continue
                seen.add(normalized)
                sanitized.append(self._truncate_example(fragment, 380))
                if len(sanitized) >= limit:
                    return sanitized

        return sanitized

    def _chunk_source_text(self, raw_text, limit=24):
        paragraphs = [clean_text(line) for line in str(raw_text or "").splitlines()]
        paragraphs = [paragraph for paragraph in paragraphs if len(paragraph) >= 3]
        if not paragraphs:
            return []

        normalized_paragraphs = []
        for paragraph in paragraphs:
            if len(paragraph) > 420:
                fragments = [
                    clean_text_preserve_prefix(part)
                    for part in re.split(r"(?<=[\.\!\?;:])\s+", paragraph)
                    if clean_text_preserve_prefix(part)
                ]
                normalized_paragraphs.extend(fragments or [paragraph])
            else:
                normalized_paragraphs.append(paragraph)

        chunks = []
        current_chunk = []
        current_length = 0
        for paragraph in normalized_paragraphs:
            if self._is_metadata_source_line(paragraph):
                continue
            if current_chunk and current_length + len(paragraph) > 420:
                chunks.append(" ".join(current_chunk))
                current_chunk = [paragraph]
                current_length = len(paragraph)
            else:
                current_chunk.append(paragraph)
                current_length += len(paragraph)
        if current_chunk:
            chunks.append(" ".join(current_chunk))

        return self._sanitize_source_items(chunks[:limit], limit=limit)

    def _build_plain_text_source_bundle(self, raw_text, blueprint):
        global_hints = self._chunk_source_text(raw_text, limit=24)
        exact_index = {}
        stage_group_index = {"intro": [], "main": [], "closing": []}

        stage_titles = self._collect_blueprint_stages(blueprint) if blueprint else []
        if global_hints and stage_titles:
            chunk_size = max(1, (len(global_hints) + len(stage_titles) - 1) // len(stage_titles))
            for index, title in enumerate(stage_titles):
                chunk = global_hints[index * chunk_size:(index + 1) * chunk_size]
                if not chunk:
                    chunk = [global_hints[min(index, len(global_hints) - 1)]]
                exact_index[("stage", self._normalize_template_name(title))] = [chunk]
                stage_group_index.setdefault(self._classify_stage_name(title), []).append(chunk)

        return exact_index, stage_group_index, global_hints

    def _build_parsed_doc_source_bundle(self, parsed_doc, raw_text, blueprint):
        exact_index = {}
        stage_group_index = {"intro": [], "main": [], "closing": []}

        for section in parsed_doc.get("sections", []):
            for kind, node in self._iter_reference_nodes(section, "section"):
                title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
                if not title:
                    continue
                items = self._sanitize_source_items(node.get("items") or [], limit=10)
                if not items:
                    continue
                exact_index.setdefault((kind, self._normalize_template_name(title)), []).append(items)
                if kind == "stage":
                    stage_group_index.setdefault(self._classify_stage_name(title), []).append(items)

        global_hints = self._chunk_source_text(raw_text, limit=24)
        if not global_hints:
            global_hints = [
                clean_text(paragraph.get("text") or "")
                for paragraph in parsed_doc.get("paragraphs", [])
                if paragraph.get("kind") == "content" and clean_text(paragraph.get("text") or "")
            ][:24]

        if blueprint:
            stage_titles = self._collect_blueprint_stages(blueprint)
            missing_titles = [
                title for title in stage_titles
                if ("stage", self._normalize_template_name(title)) not in exact_index
            ]
            if global_hints and missing_titles:
                chunk_size = max(1, (len(global_hints) + len(missing_titles) - 1) // len(missing_titles))
                for index, title in enumerate(missing_titles):
                    chunk = global_hints[index * chunk_size:(index + 1) * chunk_size]
                    if not chunk:
                        chunk = [global_hints[min(index, len(global_hints) - 1)]]
                    exact_index[("stage", self._normalize_template_name(title))] = [chunk]
                    stage_group_index.setdefault(self._classify_stage_name(title), []).append(chunk)

        return exact_index, stage_group_index, global_hints

    def _build_presentation_source_bundle(self, slides, blueprint):
        content_slides = [
            slide for slide in slides
            if not slide.get("is_title_slide") and (
                slide.get("text_blocks") or slide.get("bullet_points") or slide.get("notes") or slide.get("images")
            )
        ]
        if not content_slides:
            content_slides = slides[1:] if len(slides) > 1 else slides

        slide_summaries = {
            slide.get("number"): self._summarize_slide_for_source(slide)
            for slide in content_slides
            if slide.get("number")
        }
        global_hints = self._sanitize_source_items(
            [summary for summary in slide_summaries.values() if summary],
            limit=24,
        )
        exact_index = {}
        stage_group_index = {"intro": [], "main": [], "closing": []}
        assignments = self._assign_slides_to_stage_titles(content_slides, blueprint)

        for title, slide_numbers in assignments.items():
            chunk = [slide_summaries[number] for number in slide_numbers if slide_summaries.get(number)]
            if not chunk:
                continue
            exact_index.setdefault(("stage", self._normalize_template_name(title)), []).append(chunk)
            stage_group_index.setdefault(self._classify_stage_name(title), []).append(chunk)

        lines = []
        for title in self._collect_blueprint_stages(blueprint):
            slide_set = assignments.get(title)
            if not slide_set:
                continue
            lines.append(f'- Етап "{title}" -> {self._format_slide_numbers(slide_set)}')

        return exact_index, stage_group_index, global_hints, "\n".join(lines)

    @staticmethod
    def _empty_source_bundle():
        return {
            "context": "",
            "exact_index": {},
            "stage_group_index": {},
            "global_hints": [],
            "slide_reference_plan": "",
        }

    @staticmethod
    def _merge_source_indexes(target, source):
        for key, groups in (source or {}).items():
            target.setdefault(key, [])
            target[key].extend(groups)
        return target

    def _build_sources_bundle(self, source_files, blueprint=None):
        if not source_files:
            return self._empty_source_bundle()

        if isinstance(source_files, (str, Path)):
            source_files = [source_files]

        merged_context_parts = []
        merged_exact_index = {}
        merged_stage_group_index = {"intro": [], "main": [], "closing": []}
        merged_global_hints = []
        merged_slide_plans = []
        seen_hints = set()

        for source_file in source_files:
            bundle = self._build_source_bundle(source_file, blueprint)
            if not bundle:
                continue

            context = str(bundle.get("context") or "").strip()
            if context:
                merged_context_parts.append(context)

            self._merge_source_indexes(merged_exact_index, bundle.get("exact_index") or {})
            for group_name, groups in (bundle.get("stage_group_index") or {}).items():
                merged_stage_group_index.setdefault(group_name, [])
                merged_stage_group_index[group_name].extend(groups)

            for hint in bundle.get("global_hints") or []:
                normalized = clean_text_preserve_prefix(hint)
                if not normalized or normalized in seen_hints:
                    continue
                merged_global_hints.append(normalized)
                seen_hints.add(normalized)

            slide_plan = str(bundle.get("slide_reference_plan") or "").strip()
            if slide_plan:
                file_label = Path(source_file).name
                merged_slide_plans.append(f'Файл "{file_label}":\n{slide_plan}')

        context_header = ""
        if len(source_files) > 1:
            context_header = (
                f"ПРИКРІПЛЕНО {len(source_files)} МАТЕРІАЛИ.\n"
                "Поєднуй їхній зміст в один цілісний конспект. Якщо файли доповнюють один одного, "
                "використовуй усі. Якщо частина інформації дублюється, не повторюй її дослівно.\n\n"
            )

        merged_context = "\n\n".join(part for part in merged_context_parts if part).strip()
        if context_header and merged_context:
            merged_context = f"{context_header}{merged_context}"
        elif context_header:
            merged_context = context_header.strip()

        return {
            "context": merged_context,
            "exact_index": merged_exact_index,
            "stage_group_index": merged_stage_group_index,
            "global_hints": merged_global_hints[:24],
            "slide_reference_plan": "\n\n".join(merged_slide_plans).strip(),
        }

    def _build_source_bundle(self, source_file, blueprint=None):
        empty_bundle = self._empty_source_bundle()
        if not source_file:
            return empty_bundle

        source_path = Path(source_file)
        ext = source_path.suffix.lower()

        if ext == ".pptx":
            try:
                slides = PresentationAnalyzer.analyze(source_file)
            except Exception as exc:
                logger.warning("Failed to analyze presentation source %s: %s", source_file, exc)
                return empty_bundle
            if not slides:
                return empty_bundle

            exact_index, stage_group_index, global_hints, slide_reference_plan = self._build_presentation_source_bundle(slides, blueprint or {})
            slide_chunks = "\n".join(f"- {hint}" for hint in global_hints[:18])
            context = (
                "\n\nГОЛОВНЕ ДЖЕРЕЛО ЗМІСТУ: ЗАВАНТАЖЕНА ПРЕЗЕНТАЦІЯ.\n"
                "Фактичний зміст конспекту, приклади, вправи, питання до дітей і візуальні опори бери насамперед зі слайдів. "
                "Конспекти з base_docs використовуй тільки як зразок структури, стилю та послідовності блоків.\n\n"
                "КЛЮЧОВІ ОПОРИ ПРЕЗЕНТАЦІЇ:\n"
                f"{slide_chunks}"
            ).strip()
            return {
                "context": context,
                "exact_index": exact_index,
                "stage_group_index": stage_group_index,
                "global_hints": global_hints,
                "slide_reference_plan": slide_reference_plan,
            }

        if ext == ".docx":
            raw_text = self._read_docx_text(source_path)
            source_label = "DOCX-документ"
            parsed_source = None
            try:
                parsed_source = DocxParser.extract_structure(source_path)
            except Exception as exc:
                logger.warning("Failed to parse source DOCX %s: %s", source_file, exc)
            if parsed_source:
                exact_index, stage_group_index, global_hints = self._build_parsed_doc_source_bundle(parsed_source, raw_text, blueprint or {})
            else:
                exact_index, stage_group_index, global_hints = self._build_plain_text_source_bundle(raw_text, blueprint or {})
        elif ext == ".pdf":
            raw_text = self._read_pdf_text(source_path)
            source_label = "PDF-файл"
            exact_index, stage_group_index, global_hints = self._build_plain_text_source_bundle(raw_text, blueprint or {})
        elif ext in {".txt", ".md"}:
            raw_text = self._read_text_file(source_path)
            source_label = "текстовий файл"
            exact_index, stage_group_index, global_hints = self._build_plain_text_source_bundle(raw_text, blueprint or {})
        else:
            return empty_bundle

        if not raw_text and not global_hints:
            return empty_bundle

        outline = "\n".join(f"- {hint}" for hint in global_hints[:18])
        context = (
            f"\n\nГОЛОВНЕ ДЖЕРЕЛО ЗМІСТУ: ЗАВАНТАЖЕНИЙ {source_label.upper()}.\n"
            "Факти, визначення, приклади, вправи та формулювання бери з прикріпленого матеріалу. "
            "Конспекти з base_docs використовуй тільки як зразок структури, стилю та послідовності блоків.\n\n"
            "КЛЮЧОВІ ФРАГМЕНТИ МАТЕРІАЛУ:\n"
            f"{outline}"
        ).strip()
        return {
            "context": context,
            "exact_index": exact_index,
            "stage_group_index": stage_group_index,
            "global_hints": global_hints,
            "slide_reference_plan": "",
        }

    def _build_node_source_profile_tokens(self, node):
        parts = [
            clean_text_preserve_prefix(node.get("display_title") or node.get("title") or ""),
        ]
        parts.extend(normalize_block_list(node.get("sample_items") or [])[:6])
        parts.extend(
            clean_text_preserve_prefix(substep.get("display_title") or substep.get("title") or "")
            for substep in (node.get("substeps") or [])[:4]
        )
        parts.extend(
            clean_text_preserve_prefix(stage.get("display_title") or stage.get("title") or "")
            for stage in (node.get("stages") or [])[:4]
        )

        tokens = set()
        for part in parts:
            tokens.update(self._topic_tokens(part))
        return tokens

    def _score_global_hint_for_node(self, hint, node, kind):
        hint_text = clean_text_preserve_prefix(hint)
        if not hint_text:
            return 0.0

        hint_tokens = self._topic_tokens(hint_text)
        profile_tokens = self._build_node_source_profile_tokens(node)
        overlap = len(hint_tokens & profile_tokens)

        title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
        title_norm = self._normalize_template_name(title)
        hint_norm = self._normalize_template_name(hint_text)

        score = float(overlap) * 3.0
        if title_norm and title_norm in hint_norm:
            score += 1.5
        if kind == "stage" and "слайд" in hint_text.lower():
            score += 0.75
        if kind == "substep" and ":" in hint_text:
            score += 0.35
        return score

    def _inject_source_items_into_blueprint(self, blueprint, exact_index, stage_group_index, global_hints):
        if not blueprint:
            return

        global_pool = [
            clean_text_preserve_prefix(candidate)
            for candidate in (global_hints or [])
            if clean_text_preserve_prefix(candidate)
        ]
        used_global_indexes = set()

        def take_scored(node, kind, count):
            candidates = []
            for index, hint in enumerate(global_pool):
                if index in used_global_indexes:
                    continue
                score = self._score_global_hint_for_node(hint, node, kind)
                if score <= 0:
                    continue
                candidates.append((score, -index, index, hint))

            candidates.sort(reverse=True)
            taken = []
            for _, _, index, hint in candidates[:count]:
                used_global_indexes.add(index)
                taken.append(hint)
            return taken

        def take_global(count):
            taken = []
            for index, candidate in enumerate(global_pool):
                if index in used_global_indexes:
                    continue
                if candidate:
                    used_global_indexes.add(index)
                    taken.append(candidate)
                if len(taken) >= count:
                    break
            return taken

        def visit(node, kind, parent_source=None):
            title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
            example_sets = []
            if title:
                example_sets.extend(exact_index.get((kind, self._normalize_template_name(title)), []))
                if kind == "stage" and len(example_sets) < 2:
                        example_sets.extend(stage_group_index.get(self._classify_stage_name(title), []))

            source_items = self._merge_sample_items([], example_sets, limit=10)
            if not source_items:
                source_items = take_scored(node, kind, 3 if kind == "section" else 3)
            if not source_items and kind == "stage":
                source_items = take_global(3)
            if not source_items and parent_source:
                source_items = list(parent_source[:4])
            if kind in {"stage", "substep"} and len(source_items) < 3:
                extra_items = take_scored(node, kind, 3 - len(source_items))
                source_items = self._merge_sample_items(source_items, [extra_items], limit=6)
            if kind == "section" and len(source_items) < 3:
                extra_items = take_scored(node, kind, 3 - len(source_items))
                source_items = self._merge_sample_items(source_items, [extra_items], limit=6)
            node["source_items"] = source_items

            for stage in node.get("stages") or []:
                visit(stage, "stage", node.get("source_items") or parent_source)

            for substep in node.get("substeps") or []:
                visit(substep, "substep", node.get("source_items") or parent_source)

            if kind == "section" and not node.get("source_items"):
                aggregated = []
                for stage in node.get("stages") or []:
                    aggregated.extend(stage.get("source_items") or [])
                for substep in node.get("substeps") or []:
                    aggregated.extend(substep.get("source_items") or [])
                node["source_items"] = self._merge_sample_items([], [aggregated], limit=8) or take_global(3)

        for section in blueprint.get("sections") or []:
            visit(section, "section")

    def _build_generation_plan_template(self, blueprint):
        blueprint = blueprint if isinstance(blueprint, dict) else {}
        template = {
            "lesson_goal": "",
            "expected_results": [],
            "stages": [],
        }

        for section in blueprint.get("sections") or []:
            section_title = clean_text_preserve_prefix(section.get("display_title") or section.get("title") or "")
            for stage in section.get("stages") or []:
                stage_title = clean_text_preserve_prefix(stage.get("display_title") or stage.get("title") or "")
                if not stage_title:
                    continue
                template["stages"].append(
                    {
                        "section_title": section_title,
                        "stage_title": stage_title,
                        "stage_goal": "",
                        "stage_draft": [],
                        "teacher_actions": [],
                        "student_actions": [],
                        "activities": [],
                        "materials": [],
                        "assessment": [],
                        "source_items": self._collect_prompt_source_items(stage, limit=4),
                    }
                )
        return template

    def _coerce_plan_items(self, values, limit=8):
        cleaned = []
        seen = set()
        for item in normalize_block_list(values or []):
            normalized = normalize_sentence_punctuation(item)
            if not normalized or is_placeholder_text(normalized):
                continue
            key = self._normalize_template_name(normalized)
            if not key or key in seen:
                continue
            seen.add(key)
            cleaned.append(normalized)
            if len(cleaned) >= limit:
                break
        return cleaned

    def _is_math_subject(self, subject):
        normalized = self._normalize_for_matching(subject)
        return any(token in normalized for token in ("матем", "math", "алгебр", "геометр"))

    @staticmethod
    def _grade_to_int(grade):
        match = re.search(r"\d+", clean_text(grade))
        return int(match.group(0)) if match else 0

    def _looks_like_mojibake_text(self, value):
        text = clean_text_preserve_prefix(value)
        if not text:
            return False
        mojibake_blocks = (
            r"(?:Р .|РЎ.){3,}",
            r"(?:Гђ.|Г'.){3,}",
        )
        if any(re.search(pattern, text) for pattern in mojibake_blocks):
            return True
        markers = ("Ð", "Ñ", "Ò", "Ó", "«", "»")
        if any(marker in text for marker in markers):
            return True
        return False

    def _safe_display_text(self, value):
        text = clean_text_preserve_prefix(value)
        if not text:
            return ""
        repaired = self._repair_mojibake_utf8(text)
        if repaired and repaired != text and re.search(r"[А-Яа-яІіЇїЄєҐґ]", repaired):
            text = normalize_sentence_punctuation(repaired)
        else:
            text = normalize_sentence_punctuation(text)
        if self._looks_like_mojibake_text(text):
            return ""
        return text

    def _build_math_micro_tasks(self, grade, limit=2):
        grade_num = self._grade_to_int(grade)
        if grade_num == 3:
            tasks = [
                "Міні-вправа: обчисліть усно 6×4, 6×7, 36:6 та поясніть, як міркували.",
                "Задача: у 6 коробках по 5 олівців. Скільки олівців разом? Якою дією перевіримо відповідь?",
                "Гра «Я задумала число»: якщо число поділити на 6, отримаємо 4. Яке це число?",
            ]
        else:
            tasks = [
                "Міні-вправа: обчисліть 8×3, 24:3, 7×6 і перевірте одну відповідь оберненою дією.",
                "Коротка задача: у 4 кошиках по 9 яблук. Скільки яблук разом?",
            ]
        return tasks[: max(1, limit)]

    @staticmethod
    def _deterministic_pick(values, key, count=1):
        values = list(values or [])
        if not values:
            return []
        if len(values) <= count:
            return values
        seed = abs(hash(clean_text(str(key)))) % len(values)
        picked = []
        idx = seed
        while len(picked) < count:
            candidate = values[idx % len(values)]
            if candidate not in picked:
                picked.append(candidate)
            idx += 1
        return picked

    def _is_meta_report_line(self, value):
        text = clean_text_preserve_prefix(value)
        normalized = self._normalize_template_name(text)
        if not normalized:
            return True
        if any(marker in normalized for marker in self.META_REPORT_MARKERS):
            return True
        if normalized.startswith(("є ", "наявн", "присутн", "включає", "містить")):
            return True
        if normalized.startswith(("учитель ", "вчитель ")) and "?" not in text and len(normalized.split()) <= 6:
            return True
        return False

    def _rewrite_meta_line_to_direct(self, value, *, stage_title="", subject="", grade=""):
        normalized = self._normalize_template_name(value)
        stage_group = self._classify_stage_name(stage_title or "")
        is_math3 = self._is_math_subject(subject) and self._grade_to_int(grade) == 3
        if is_math3:
            if "уточнювальне запитання" in normalized or "запитання на розуміння" in normalized:
                return "Запитання: «Що спільного в прикладах 6×4 і 24:6?»"
            if "коротко озвучує мету" in normalized:
                return "Учитель: «Сьогодні закріплюємо таблицю множення і ділення на 6 та розв'язуємо задачі.»"
            if "швидкої перевірки" in normalized or "самооцін" in normalized:
                return "Покажіть сигналом руки: легко / було складно, і назвіть один приклад, який поясните класу."
            if "математичний сигнал" in normalized:
                return "Обчисліть усно: 6×5, 6×8, 42:6. Поясніть, як перевірили відповідь."
            if stage_group == "intro":
                return "Учитель: «Доброго дня! Приготуйте зошити, починаємо усний рахунок на 6.»"
        return ""

    def _is_topic_leak_line(self, value, *, topic="", subject="", grade=""):
        text = clean_text_preserve_prefix(value)
        normalized = self._normalize_template_name(text)
        if not normalized:
            return True
        if self._is_math_subject(subject) and self._grade_to_int(grade) == 3:
            if any(phrase in normalized for phrase in self.MATH3_GENERIC_PHRASES):
                return True
            if any(token in normalized for token in self.MATH_OFFTOPIC_TOKENS):
                return True
            if "на етапі «" in normalized and not re.search(r"\d", text):
                return True
            topic_tokens = self._topic_tokens(topic)
            line_tokens = self._topic_tokens(text)
            if topic_tokens and not (topic_tokens & line_tokens):
                if not any(token in normalized for token in ("6", "множ", "ділен", "задач", "приклад", "обчис")):
                    return True
        return False

    def _rewrite_generic_math3_line(self, line, *, stage_group=""):
        normalized = self._normalize_template_name(line)
        if not normalized:
            return ""
        rewrites = [
            ("подорож світом математики", "Давайте швидко: 6×2? 6×6? 24:6?"),
            ("подорож світом чисел", "Давайте швидко: 6×2? 6×6? 24:6?"),
            ("пам'ятаємо про правила роботи", "Працюємо по черзі: називаємо дію, пояснюємо відповідь коротко."),
            ("пам'ятаємо про правила безпеки та поведінки", "Працюємо по черзі: називаємо дію, пояснюємо відповідь коротко."),
            ("дотримуємося правил безпеки", "Тримаємо темп уроку: слухаємо умову, відповідаємо чітко."),
            ("подаруйте гарний настрій", "Налаштуймося на рахунок: усно порахуйте 6×2 і 24:6."),
            ("привітайтеся з гостями", "Відкрийте зошити, запишіть дату і підготуйтеся до усного рахунку."),
            ("виконайте коротке завдання й поясніть відповідь однокласнику", "Розв'яжіть 2 приклади на множення і 2 на ділення на 6 та поясніть один у парі."),
            ("запишіть каліграфічно число 66", "Запишіть дату і виконайте усно: 6×6, 42:6."),
        ]
        for marker, replacement in rewrites:
            if marker in normalized:
                return replacement
        if stage_group == "intro" and ("організацій" in normalized or "психолог" in normalized):
            return "Давайте швидко: 6×4? 30:6? 6×7?"
        return ""

    def _simplify_math3_line(self, line):
        text = normalize_sentence_punctuation(line)
        if not text:
            return ""
        simple_replacements = (
            ("Швидкий старт:", "Давайте швидко:"),
            ("Будемо застосовувати", "Застосуємо"),
            ("Будемо вчитися", "Навчимося"),
            ("Сьогодні ми поринемо у світ", "Сьогодні працюємо з"),
            ("Сьогодні ми зануримося у світ", "Сьогодні працюємо з"),
            ("Отже, тема нашого уроку:", "Тема уроку:"),
            ("Приготуйтеся відповідати швидко і правильно", "Працюємо усно: відповідаємо швидко й точно"),
        )
        for source, target in simple_replacements:
            if source in text:
                text = text.replace(source, target)
        return normalize_sentence_punctuation(text)

    def _polish_math3_stage_entry(self, entry):
        entry = dict(entry or {})
        stage_title = clean_text(entry.get("stage") or "")
        stage_group = self._classify_stage_name(stage_title)
        stage_key = self._normalize_template_name(stage_title)

        for field in ("teacher_actions", "student_actions", "activities"):
            polished = []
            for line in normalize_list(entry.get(field) or []):
                simplified = self._simplify_math3_line(line)
                if not simplified:
                    continue
                normalized = self._normalize_template_name(simplified)
                if (
                    stage_group == "intro"
                    and any(token in normalized for token in ("очікую", "настр", "поведінк"))
                    and not re.search(r"\d", simplified)
                ):
                    continue
                polished.append(simplified)
            entry[field] = self._unique_clean_items(
                polished,
                expected_field=field if field in {"teacher_actions", "student_actions", "activities"} else None,
                limit=12 if field == "activities" else 10,
            )

        if stage_group == "intro":
            simple_intro_questions = []
            for line in entry.get("teacher_actions") or []:
                normalized = self._normalize_template_name(line)
                if "?" in line and any(token in normalized for token in ("спільного", "чому", "найважч", "порівняйте", "доведіть")):
                    continue
                simple_intro_questions.append(line)
            if not any("?" in line for line in simple_intro_questions):
                starter_pool = [
                    "Учитель: «Давайте швидко: 6×2? 6×6? 24:6?»",
                    "Учитель: «Почнемо з трьох прикладів: 6×3? 18:6? 30:6?»",
                ]
                simple_intro_questions.extend(self._deterministic_pick(starter_pool, key=f"{stage_key}:intro_q", count=1))
            entry["teacher_actions"] = self._unique_clean_items(simple_intro_questions, expected_field="teacher_actions", limit=10)

        has_question = any("?" in line for line in [*(entry.get("teacher_actions") or []), *(entry.get("activities") or [])])
        has_student_quote = any("учні:" in self._normalize_template_name(line) or "«" in line for line in (entry.get("student_actions") or []))
        if has_question and not has_student_quote:
            student_dialogue_pool = [
                "Учні: «7.»",
                "Учні: «Я поділив 42 на 6.»",
                "Учні: «48:6=8, тому відповідь 8.»",
            ]
            teacher_feedback_pool = [
                "Учитель: «Добре, пояснив правильно.»",
                "Учитель: «Так, саме так міркуємо.»",
            ]
            entry["student_actions"] = self._unique_clean_items(
                [*(entry.get("student_actions") or []), *self._deterministic_pick(student_dialogue_pool, key=f"{stage_key}:stud", count=2)],
                expected_field="student_actions",
                limit=6,
            )
            entry["teacher_actions"] = self._unique_clean_items(
                [*(entry.get("teacher_actions") or []), *self._deterministic_pick(teacher_feedback_pool, key=f"{stage_key}:teach", count=1)],
                expected_field="teacher_actions",
                limit=6,
            )

        if any("«" in line or "учні:" in self._normalize_template_name(line) for line in (entry.get("student_actions") or [])):
            short_reactions = ["Добре.", "Правильно.", "Перевірмо разом.", "Хто думає інакше?"]
            if not any(self._normalize_template_name(line) in {self._normalize_template_name(x) for x in short_reactions} for line in (entry.get("teacher_actions") or [])):
                entry["teacher_actions"] = self._unique_clean_items(
                    [*(entry.get("teacher_actions") or []), *self._deterministic_pick(short_reactions, key=f"{stage_key}:react", count=1)],
                    expected_field="teacher_actions",
                    limit=7,
                )
        return entry

    def _postprocess_math3_lesson_flow(self, lesson_flow):
        stages = [dict(item) for item in (lesson_flow or []) if isinstance(item, dict)]
        if not stages:
            return stages

        varied_main_activities = [
            "Давайте швидко: 6×7? 42:6? 6×9?",
            "Порівняйте: 6×5 і 30:6. Що спільного?",
            "Питання-пастка: 6×8=46 чи 48? Поясніть.",
            "У парі складіть коротку задачу до виразу 6×4.",
        ]
        prev_sig = ""
        prev_student_sigs = set()
        student_variants = [
            "Учні: «7.»",
            "Учні: «Я поділив 42 на 6.»",
            "Учні: «48:6=8, тому відповідь 8.»",
            "Учні: «Спочатку 6×6=36, потім перевірили діленням.»",
        ]
        for idx, stage in enumerate(stages):
            stage = self._polish_math3_stage_entry(stage)
            student_lines = []
            for line in normalize_list(stage.get("student_actions") or []):
                sig = self._normalize_template_name(line)
                if sig and sig in prev_student_sigs:
                    continue
                student_lines.append(line)
            if not student_lines:
                student_lines = self._deterministic_pick(student_variants, key=f"{stage.get('stage') or ''}:{idx}:students", count=1)
            stage["student_actions"] = self._unique_clean_items(student_lines, expected_field="student_actions", limit=8)
            reaction_pool = ["Добре.", "Правильно.", "Перевірмо разом.", "Хто думає інакше?"]
            reaction_norms = {self._normalize_template_name(item) for item in reaction_pool}
            teacher_lines = normalize_list(stage.get("teacher_actions") or [])
            has_student_quotes = any("В«" in line for line in (stage.get("student_actions") or []))
            has_reaction = any(self._normalize_template_name(line) in reaction_norms for line in teacher_lines)
            if has_student_quotes and not has_reaction:
                teacher_lines.extend(self._deterministic_pick(reaction_pool, key=f"{stage.get('stage') or ''}:{idx}:reaction", count=1))
                stage["teacher_actions"] = self._unique_clean_items(teacher_lines, expected_field="teacher_actions", limit=10)

            activities = normalize_list(stage.get("activities") or [])
            stage_group = self._classify_stage_name(stage.get("stage") or "")
            if activities:
                first_sig = self._normalize_template_name(activities[0])
                if first_sig and first_sig == prev_sig:
                    replacement = self._deterministic_pick(varied_main_activities, key=f"{stage.get('stage') or ''}:{idx}", count=1)[0]
                    if stage_group == "intro":
                        replacement = "Давайте швидко: 6×2? 6×6? 24:6?"
                    elif stage_group == "closing":
                        replacement = "Коротко підсумуйте: який приклад на 6 вийшов найкраще?"
                    activities[0] = replacement
                    stage["activities"] = self._unique_clean_items(activities, expected_field="activities", limit=12)
                    first_sig = self._normalize_template_name(replacement)
                prev_sig = first_sig

            prev_student_sigs = {self._normalize_template_name(line) for line in (stage.get("student_actions") or []) if self._normalize_template_name(line)}
            stages[idx] = stage
        return stages

    def _rewrite_weak_math3_expression(self, line):
        text = normalize_sentence_punctuation(line)
        normalized = self._normalize_template_name(text)
        if not text:
            return ""
        weak_signals = (
            "5 * 6 + 10 / 4 - 7",
            "5Г—6+10/4-7",
            "ланцюжок",
        )
        has_mixed_ops = len(re.findall(r"[\+\-\*/Г—С…x]", text)) >= 3 and ("/" in text)
        if any(signal in normalized for signal in weak_signals) or has_mixed_ops:
            return "Обчисліть усно: 6×5, 36:6, 24:6, 6×7. Поясніть один обраний приклад."
        return text

    def _sanitize_stage_script_lines(self, items, *, stage_title="", topic="", subject="", grade="", allow_assessment=True):
        cleaned = []
        actions = []
        seen = set()
        stage_group = self._classify_stage_name(stage_title or "")
        is_math3 = self._is_math_subject(subject) and self._grade_to_int(grade) == 3
        for raw in normalize_block_list(items or []):
            line = normalize_sentence_punctuation(raw)
            if not line:
                continue
            if is_math3:
                rewritten_generic = self._rewrite_generic_math3_line(line, stage_group=stage_group)
                if rewritten_generic and rewritten_generic != line:
                    actions.append({"action": "rewritten_generic", "from": line, "to": rewritten_generic})
                    line = rewritten_generic
                rewritten_expr = self._rewrite_weak_math3_expression(line)
                if rewritten_expr and rewritten_expr != line:
                    actions.append({"action": "rewritten_weak_task", "from": line, "to": rewritten_expr})
                    line = rewritten_expr
            if self._is_meta_report_line(line):
                rewritten = self._rewrite_meta_line_to_direct(
                    line,
                    stage_title=stage_title,
                    subject=subject,
                    grade=grade,
                )
                if rewritten:
                    actions.append({"action": "rewritten", "from": line, "to": rewritten})
                    line = rewritten
                else:
                    actions.append({"action": "removed_meta", "from": line})
                    continue
            if self._is_topic_leak_line(line, topic=topic, subject=subject, grade=grade):
                actions.append({"action": "removed_leak", "from": line})
                continue
            if not allow_assessment and (
                self._looks_like_assessment_line(line)
                or self._looks_like_reflection_line(line)
                or "самооцін" in self._normalize_template_name(line)
            ):
                actions.append({"action": "removed_overused_assessment", "from": line})
                continue
            key = self._normalize_template_name(line)
            if not key or key in seen:
                continue
            seen.add(key)
            cleaned.append(line)
        if is_math3 and stage_group == "closing":
            normalized_join = " ".join(self._normalize_template_name(item) for item in cleaned)
            if "сьогодні на уроці я навчився" in normalized_join:
                cleaned = [
                    "Учитель: «Закінчіть речення: сьогодні я навчився(лася)...»" if "сьогодні на уроці я навчився" in self._normalize_template_name(item) else item
                    for item in cleaned
                ]
            if "я ще хочу попрацювати" in normalized_join:
                cleaned = [
                    "Учитель: «Що ще варто потренувати вдома: множення на 6 чи ділення на 6?»" if "я ще хочу попрацювати" in self._normalize_template_name(item) else item
                    for item in cleaned
                ]
            if not any(self._looks_like_teacher_line(item) for item in cleaned):
                cleaned.insert(0, "Учитель: «Назвіть приклад на 6, який сьогодні вдався найкраще.»")
            if not any(self._looks_like_student_line(item) for item in cleaned):
                cleaned.extend(
                    [
                        "Учні: «Мені вдалося швидко знайти 42:6=7.»",
                        "Учні: «Я ще потреную 6×8=48.»",
                    ]
                )
            if not any("смайлик" in self._normalize_template_name(item) or "сигналом" in self._normalize_template_name(item) for item in cleaned):
                cleaned.append("Покажіть смайликом: 🙂 — впевнено, 😐 — треба ще тренування, 🙁 — потрібна допомога.")
            if not any("дякую" in self._normalize_template_name(item) or "молодці" in self._normalize_template_name(item) for item in cleaned):
                cleaned.append("Учитель: «Молодці, дякую за роботу. Наступного уроку перевіримо себе ще раз на прикладах із 6.»")
        return cleaned, actions

    def _build_stage_draft_from_fields(self, stage, subject="", grade=""):
        stage = stage if isinstance(stage, dict) else {}
        draft = self._coerce_plan_items(stage.get("stage_draft") or [], limit=10)
        if draft:
            return draft

        lines = []
        goal = clean_text(stage.get("stage_goal") or "")
        if goal:
            lines.append(f"Мета етапу: {goal}.")
        for item in self._coerce_plan_items(stage.get("teacher_actions") or [], limit=3):
            lines.append(item if ":" in item else f"Учитель: {item}")
        for item in self._coerce_plan_items(stage.get("student_actions") or [], limit=3):
            lines.append(item if ":" in item else f"Учні: {item}")
        for item in self._coerce_plan_items(stage.get("activities") or [], limit=4):
            lines.append(item if ":" in item else f"Завдання: {item}")
        for item in self._coerce_plan_items(stage.get("assessment") or [], limit=2):
            lines.append(f"Міні-перевірка: {item}")

        if self._is_math_subject(subject):
            lines.extend(self._build_math_micro_tasks(grade, limit=1))
        return self._coerce_plan_items(lines, limit=10)

    def _enrich_stage_draft_locally(self, stage_draft, *, stage_title="", subject="", grade=""):
        lines = self._coerce_plan_items(stage_draft or [], limit=12)
        if not lines:
            lines = []
        stage_group = self._classify_stage_name(stage_title or "")
        is_math3 = self._is_math_subject(subject) and self._grade_to_int(grade) == 3

        has_teacher = any(self._looks_like_teacher_line(line) for line in lines)
        has_student = any(self._looks_like_student_line(line) for line in lines)
        has_task = any(self._has_task_signal(line) for line in lines)
        has_question = any("?" in line for line in lines)
        has_assessment = any(self._looks_like_assessment_line(line) or self._looks_like_reflection_line(line) for line in lines)
        title_hint = self._safe_display_text(stage_title) or "етап уроку"

        if not has_teacher:
            teacher_variants = [
                f"Учитель: «Працюємо на етапі «{title_hint}». Слухаємо інструкцію і виконуємо по кроках.»",
                "Учитель: «Спочатку рахуємо усно, потім перевіряємо відповідь у парі.»",
                "Учитель: «Поясніть хід міркування повним реченням.»",
            ]
            if is_math3 and stage_group == "intro":
                teacher_variants = [
                    "Учитель: «Починаємо з усного рахунку на 6. Готові?»",
                    "Учитель: «Давайте швидко: називаю приклад — ви даєте відповідь і коротке пояснення.»",
                    "Учитель: «Працюємо коротко й точно: приклад, відповідь, перевірка.»",
                ]
            lines.extend(self._deterministic_pick(teacher_variants, key=title_hint, count=1))
        if not has_student:
            student_variants = [
                "Учні в парах: один називає приклад, інший відповідає та пояснює спосіб.",
                "Учні обмінюються відповідями в парі та коригують помилки одне одного.",
                "Типова відповідь учня: «6×7=42, бо це сім разів по шість».",
            ]
            lines.extend(self._deterministic_pick(student_variants, key=f"{title_hint}:students", count=1))
        if not has_task:
            if self._is_math_subject(subject):
                if is_math3 and stage_group == "intro":
                    lines.append("Усний рахунок: 6×2, 6×6, 24:6, 30:6.")
                else:
                    lines.extend(self._build_math_micro_tasks(grade, limit=1))
            else:
                lines.append("Виконайте коротке завдання в парі та сформулюйте запитання до відповіді однокласника.")
        if not has_question:
            if self._is_math_subject(subject):
                question_variants = [
                    "Що спільного в прикладах 6×4 і 24:6?",
                    "Який крок у розв'язанні задачі був найважчим і чому?",
                    "Як перевірити, що відповідь до 36:6 правильна?",
                ]
                lines.extend(self._deterministic_pick(question_variants, key=f"{title_hint}:q", count=1))
            else:
                lines.append("Що допомогло вам виконати завдання точніше?")
        if not has_assessment and stage_group in {"main", "closing"}:
            assess_variants = [
                "Перевірка: обміняйтесь зошитами в парі та звірте одну відповідь за зразком на дошці.",
                "Покажіть сигналом руки: легко / було складно, і назвіть приклад, який поясните класу.",
            ]
            lines.extend(self._deterministic_pick(assess_variants, key=f"{title_hint}:a", count=1))
        if is_math3 and stage_group == "closing":
            lines.extend(
                [
                    "Учитель: «Що сьогодні вийшло найкраще: множення на 6 чи ділення на 6?»",
                    "Учні: «Найкраще вийшло ділити 36:6 і 48:6.»",
                    "Покажіть смайликом свій стан: 🙂 / 😐 / 🙁.",
                ]
            )
        if is_math3 and not any(re.search(r"\d", line) for line in lines):
            lines.extend(self._build_math_micro_tasks(grade, limit=1))

        return self._coerce_plan_items(lines, limit=12)

    def _coerce_generation_plan(self, plan_payload, blueprint, subject="", grade=""):
        template = self._build_generation_plan_template(blueprint)
        if not isinstance(plan_payload, dict):
            return template

        normalized_goal = clean_text(plan_payload.get("lesson_goal") or plan_payload.get("goal") or "")
        if normalized_goal:
            template["lesson_goal"] = normalized_goal

        template["expected_results"] = self._coerce_plan_items(
            plan_payload.get("expected_results") or plan_payload.get("results") or [],
            limit=10,
        )

        source_stages = plan_payload.get("stages") if isinstance(plan_payload.get("stages"), list) else []
        exact_index = {}
        for stage in source_stages:
            if not isinstance(stage, dict):
                continue
            title = clean_text_preserve_prefix(stage.get("stage_title") or stage.get("title") or stage.get("stage") or "")
            key = self._normalize_template_name(title)
            if key and key not in exact_index:
                exact_index[key] = stage

        for index, target_stage in enumerate(template["stages"]):
            target_key = self._normalize_template_name(target_stage.get("stage_title") or "")
            source_stage = exact_index.get(target_key)
            if not source_stage and index < len(source_stages) and isinstance(source_stages[index], dict):
                source_stage = source_stages[index]
            if not source_stage:
                continue

            target_stage["stage_goal"] = clean_text(
                source_stage.get("stage_goal")
                or source_stage.get("goal")
                or ""
            )
            target_stage["stage_draft"] = self._coerce_plan_items(
                source_stage.get("stage_draft")
                or source_stage.get("draft")
                or source_stage.get("narrative")
                or source_stage.get("content_block")
                or [],
                limit=12,
            )
            target_stage["teacher_actions"] = self._coerce_plan_items(
                source_stage.get("teacher_actions") or source_stage.get("teacher") or [],
                limit=8,
            )
            target_stage["student_actions"] = self._coerce_plan_items(
                source_stage.get("student_actions") or source_stage.get("students") or [],
                limit=8,
            )
            target_stage["activities"] = self._coerce_plan_items(
                source_stage.get("activities") or source_stage.get("tasks") or [],
                limit=8,
            )
            target_stage["materials"] = self._coerce_plan_items(
                source_stage.get("materials") or source_stage.get("resources") or [],
                limit=4,
            )
            target_stage["assessment"] = self._coerce_plan_items(
                source_stage.get("assessment") or source_stage.get("reflection") or [],
                limit=4,
            )
            target_stage["source_items"] = self._coerce_plan_items(
                [*(target_stage.get("source_items") or []), *(source_stage.get("source_items") or [])],
                limit=12,
            )

            fallback_source = self._coerce_plan_items(target_stage.get("source_items") or [], limit=6)
            if not target_stage["teacher_actions"] and fallback_source:
                target_stage["teacher_actions"] = fallback_source[:2]
            if not target_stage["student_actions"] and len(fallback_source) > 1:
                target_stage["student_actions"] = fallback_source[1:3]
            if not target_stage["activities"] and fallback_source:
                target_stage["activities"] = fallback_source[:3]
            if not target_stage["stage_draft"]:
                target_stage["stage_draft"] = self._build_stage_draft_from_fields(
                    target_stage,
                    subject=subject,
                    grade=grade,
                )
            target_stage["stage_draft"] = self._enrich_stage_draft_locally(
                target_stage.get("stage_draft") or [],
                stage_title=target_stage.get("stage_title") or "",
                subject=subject,
                grade=grade,
            )
            inferred_from_draft = self._infer_stage_fields_from_items(target_stage.get("stage_draft") or [])
            if not target_stage["teacher_actions"]:
                target_stage["teacher_actions"] = self._coerce_plan_items(
                    inferred_from_draft.get("teacher_actions") or [],
                    limit=6,
                )
            if not target_stage["student_actions"]:
                target_stage["student_actions"] = self._coerce_plan_items(
                    inferred_from_draft.get("student_actions") or [],
                    limit=6,
                )
            if not target_stage["activities"]:
                target_stage["activities"] = self._coerce_plan_items(
                    inferred_from_draft.get("activities") or [],
                    limit=8,
                )
            if not target_stage["assessment"]:
                target_stage["assessment"] = self._coerce_plan_items(
                    [
                        *(inferred_from_draft.get("assessment") or []),
                        *(inferred_from_draft.get("reflection") or []),
                    ],
                    limit=4,
                )

        return template

    def _compress_context_for_prompt(self, text, *, max_chars=2600, max_lines=70):
        lines = [
            clean_text_preserve_prefix(line)
            for line in str(text or "").splitlines()
            if clean_text_preserve_prefix(line)
        ]
        if len(lines) > max_lines:
            lines = lines[:max_lines]
            lines.append("...[context trimmed]")

        compact = "\n".join(lines).strip()
        if len(compact) > max_chars:
            compact = compact[:max_chars].rstrip()
            compact = f"{compact}\n...[context trimmed]"
        return compact

    def _build_generation_plan_prompt_outline(self, generation_plan, stage_limit=14):
        generation_plan = generation_plan if isinstance(generation_plan, dict) else {}
        stages = generation_plan.get("stages") if isinstance(generation_plan.get("stages"), list) else []
        outline = {
            "lesson_goal": clean_text(generation_plan.get("lesson_goal") or ""),
            "expected_results": self._coerce_plan_items(generation_plan.get("expected_results") or [], limit=6),
            "stages": [],
        }

        for stage in stages[: max(1, stage_limit)]:
            if not isinstance(stage, dict):
                continue
            outline["stages"].append(
                {
                    "section_title": clean_text(stage.get("section_title") or ""),
                    "stage_title": clean_text(stage.get("stage_title") or stage.get("title") or stage.get("stage") or ""),
                    "stage_goal": clean_text(stage.get("stage_goal") or ""),
                    "stage_draft": self._coerce_plan_items(stage.get("stage_draft") or [], limit=4),
                    "teacher_actions": self._coerce_plan_items(stage.get("teacher_actions") or [], limit=2),
                    "student_actions": self._coerce_plan_items(stage.get("student_actions") or [], limit=2),
                    "activities": self._coerce_plan_items(stage.get("activities") or [], limit=3),
                    "materials": self._coerce_plan_items(stage.get("materials") or [], limit=2),
                    "assessment": self._coerce_plan_items(stage.get("assessment") or [], limit=1),
                    "source_items": self._coerce_plan_items(stage.get("source_items") or [], limit=1),
                }
            )
        return outline

    def _extract_style_cues(self, blueprint, limit=8):
        cues = []
        for section in (blueprint or {}).get("sections") or []:
            for stage in section.get("stages") or []:
                stage_title = self._safe_display_text(stage.get("display_title") or stage.get("title") or "")
                stage_group = self._classify_stage_name(stage_title)
                source_pool = [
                    *(stage.get("source_items") or []),
                    *(stage.get("sample_items") or []),
                ]
                for raw in source_pool[:8]:
                    line = self._safe_display_text(raw)
                    if not line or self._is_template_noise_line(self._normalize_template_name(line)):
                        continue
                    if self._is_generic_lesson_item(line):
                        continue
                    if stage_group == "intro" and not any(token in self._normalize_template_name(line) for token in ("привіт", "настр", "мотивац", "гра", "вправа", "вірш", "усміх")):
                        continue
                    cues.append(line)
                    if len(cues) >= limit:
                        return cues
        return cues[:limit]

    def _build_script_prompt_outline(self, blueprint, section_limit=8, block_limit=6):
        outline = []
        for section in (blueprint or {}).get("sections") or []:
            section_title = self._safe_display_text(section.get("display_title") or section.get("title") or "")
            if not section_title:
                continue
            block_titles = []
            for stage in section.get("stages") or []:
                stage_title = self._safe_display_text(stage.get("display_title") or stage.get("title") or "")
                if stage_title:
                    block_titles.append(stage_title)
                if len(block_titles) >= block_limit:
                    break
            if not block_titles:
                for substep in section.get("substeps") or []:
                    substep_title = self._safe_display_text(substep.get("display_title") or substep.get("title") or "")
                    if substep_title:
                        block_titles.append(substep_title)
                    if len(block_titles) >= block_limit:
                        break
            outline.append({"title": section_title, "blocks": block_titles})
            if len(outline) >= section_limit:
                break
        return outline

    def _coerce_script_lines(self, lines, *, topic="", subject="", grade="", line_limit=8):
        raw_lines = normalize_block_list(lines or [])
        prepared = []
        for raw in raw_lines:
            line = self._safe_display_text(raw) or normalize_sentence_punctuation(raw)
            if not line:
                continue
            low = self._normalize_template_name(line)
            for prefix in ("вчитель:", "учитель:", "учні:", "діяльність:", "teacher:", "students:", "activities:"):
                if low.startswith(prefix):
                    line = normalize_sentence_punctuation(line[len(prefix):].strip())
                    low = self._normalize_template_name(line)
                    break
            if not line:
                continue
            if "типова відповідь" in low:
                line = "(відповіді учнів)"
            if self._is_topic_leak_line(line, topic=topic, subject=subject, grade=grade):
                continue
            if not line.startswith(("–", "(")):
                line = f"– {line}"
            prepared.append(line)

        cleaned_lines, _ = self._sanitize_stage_script_lines(
            prepared,
            stage_title="",
            topic=topic,
            subject=subject,
            grade=grade,
            allow_assessment=True,
        )
        result = []
        seen = set()
        for line in cleaned_lines:
            normalized = self._normalize_template_name(line)
            if not normalized or normalized in seen:
                continue
            seen.add(normalized)
            result.append(line)
            if len(result) >= max(3, line_limit):
                break
        return result

    def _coerce_lesson_script(self, payload, *, topic="", grade="", subject="", blueprint=None):
        payload = payload if isinstance(payload, dict) else {}
        header = payload.get("header") if isinstance(payload.get("header"), dict) else {}
        lesson_sections = payload.get("lesson_sections")
        if not isinstance(lesson_sections, list):
            lesson_sections = payload.get("sections")
        lesson_sections = lesson_sections if isinstance(lesson_sections, list) else []

        outline = self._build_script_prompt_outline(blueprint or {}, section_limit=8, block_limit=6)
        default_topic = clean_text(topic)
        default_goal = "Сформувати предметні вміння за темою уроку."
        default_equipment = "Підручник, зошит, картки із завданнями."
        script = {
            "header": {
                "topic": clean_text(header.get("topic") or payload.get("topic") or default_topic),
                "goal": clean_text(header.get("goal") or header.get("meta") or payload.get("goal") or default_goal),
                "equipment": clean_text(header.get("equipment") or payload.get("equipment") or default_equipment),
                "lesson_type": clean_text(header.get("lesson_type") or payload.get("lesson_type") or ""),
                "expected_results": self._coerce_plan_items(
                    header.get("expected_results") or payload.get("expected_results") or [],
                    limit=6,
                ),
            },
            "lesson_sections": [],
        }

        for section_index, section in enumerate(lesson_sections[:10]):
            if not isinstance(section, dict):
                continue
            fallback_title = (outline[section_index]["title"] if section_index < len(outline) else f"Р РѕР·РґС–Р» {section_index + 1}")
            section_title = self._safe_display_text(section.get("title") or fallback_title) or fallback_title
            blocks_raw = section.get("blocks")
            if not isinstance(blocks_raw, list):
                blocks_raw = section.get("items") if isinstance(section.get("items"), list) else []
            normalized_blocks = []
            for block_index, block in enumerate(blocks_raw[:8]):
                block = block if isinstance(block, dict) else {}
                fallback_block = ""
                if section_index < len(outline):
                    section_blocks = outline[section_index].get("blocks") or []
                    if block_index < len(section_blocks):
                        fallback_block = section_blocks[block_index]
                block_title = self._safe_display_text(
                    block.get("title") or block.get("name") or fallback_block or f"РџС–РґРїСѓРЅРєС' {block_index + 1}"
                ) or f"РџС–РґРїСѓРЅРєС' {block_index + 1}"
                block_lines = self._coerce_script_lines(
                    block.get("lines") or block.get("items") or [],
                    topic=topic,
                    subject=subject,
                    grade=grade,
                    line_limit=8,
                )
                if len(block_lines) < 2:
                    continue
                normalized_blocks.append({"title": block_title, "lines": block_lines})
            if not normalized_blocks:
                continue
            script["lesson_sections"].append({"title": section_title, "blocks": normalized_blocks})

        if not script["lesson_sections"]:
            fallback_sections = outline[:3] if outline else [{"title": "Хід уроку", "blocks": ["Робота за темою"]}]
            for section in fallback_sections:
                lines = self._coerce_script_lines(
                    [
                        "Працюємо за темою уроку.",
                        "Дайте відповідь на запитання вчителя.",
                        "(відповіді учнів)",
                        "Перевірмо разом.",
                    ],
                    topic=topic,
                    subject=subject,
                    grade=grade,
                    line_limit=6,
                )
                script["lesson_sections"].append(
                    {
                        "title": section.get("title") or "Хід уроку",
                        "blocks": [{"title": (section.get("blocks") or ["Робота за темою"])[0], "lines": lines}],
                    }
                )
        return script

    def _assess_lesson_script_alignment(self, script, *, topic="", subject="", grade=""):
        sections = script.get("lesson_sections") if isinstance(script, dict) else []
        sections = sections if isinstance(sections, list) else []
        total_lines = 0
        offtopic_lines = 0
        blocks_count = 0
        math_number_lines = 0
        issues = []

        for section in sections:
            for block in (section or {}).get("blocks") or []:
                blocks_count += 1
                for line in normalize_block_list((block or {}).get("lines") or []):
                    total_lines += 1
                    if self._is_topic_leak_line(line, topic=topic, subject=subject, grade=grade):
                        offtopic_lines += 1
                    if self._is_math_subject(subject) and re.search(r"\d", line):
                        math_number_lines += 1

        offtopic_ratio = (offtopic_lines / total_lines) if total_lines else 1.0
        score = (
            max(0.0, 1.0 - offtopic_ratio) * 0.65
            + min(1.0, total_lines / 24.0) * 0.2
            + min(1.0, blocks_count / 8.0) * 0.15
        )
        if self._is_math_subject(subject) and self._grade_to_int(grade) >= 2 and math_number_lines < 3:
            issues.append("для математики замало конкретних чисел/виразів")
            score *= 0.8
        if offtopic_ratio > 0.2:
            issues.append("виявлено інородні фрагменти, не пов'язані з темою")
        if total_lines < 12:
            issues.append("замало змістових реплік для проведення уроку")
        if len(sections) < 2:
            issues.append("замало секцій у структурі ходу уроку")

        acceptable = (
            len(sections) >= 2
            and total_lines >= 12
            and offtopic_ratio <= 0.2
            and score >= 0.62
        )
        return {
            "score": round(score, 4),
            "acceptable": acceptable,
            "section_count": len(sections),
            "blocks_count": blocks_count,
            "total_lines": total_lines,
            "offtopic_lines": offtopic_lines,
            "offtopic_ratio": round(offtopic_ratio, 4),
            "issues": issues,
        }

    def _lesson_script_to_document_model(self, script, *, topic="", grade="", subject=""):
        script = script if isinstance(script, dict) else {}
        header = script.get("header") if isinstance(script.get("header"), dict) else {}
        sections_raw = script.get("lesson_sections") if isinstance(script.get("lesson_sections"), list) else []

        header_fields = [
            {"label": "Тема", "value": clean_text(header.get("topic") or topic) or "—", "style": "Metodist Body"},
            {"label": "Мета", "value": clean_text(header.get("goal") or "") or "—", "style": "Metodist Body"},
            {"label": "Обладнання", "value": clean_text(header.get("equipment") or "") or "—", "style": "Metodist Body"},
        ]
        lesson_type = clean_text(header.get("lesson_type") or "")
        expected_results = self._coerce_plan_items(header.get("expected_results") or [], limit=6)
        if lesson_type:
            header_fields.append({"label": "Тип уроку", "value": lesson_type, "style": "Metodist Body"})
        if expected_results:
            header_fields.append(
                {"label": "Очікувані результати", "value": "; ".join(expected_results), "style": "Metodist Body"}
            )
        if clean_text(grade):
            header_fields.append({"label": "Клас", "value": clean_text(grade), "style": "Metodist Body"})
        if clean_text(subject):
            header_fields.append({"label": "Предмет", "value": clean_text(subject), "style": "Metodist Body"})

        sections = []
        lesson_flow = []
        for section in sections_raw:
            section_title = self._safe_display_text((section or {}).get("title") or "") or "Хід уроку"
            substeps = []
            for index, block in enumerate((section or {}).get("blocks") or [], start=1):
                block_title = self._safe_display_text((block or {}).get("title") or "") or f"РџС–РґРїСѓРЅРєС' {index}"
                block_lines = self._coerce_script_lines(
                    (block or {}).get("lines") or [],
                    topic=topic,
                    subject=subject,
                    grade=grade,
                    line_limit=8,
                )
                if len(block_lines) < 2:
                    continue
                substeps.append(
                    {
                        "title": f"{index}. {block_title}",
                        "display_title": f"{index}. {block_title}",
                        "style": "Metodist Substep",
                        "content_style": "Metodist Body",
                        "items": block_lines,
                        "sample_item_styles": [],
                        "substeps": [],
                        "children_order": [],
                    }
                )
                lesson_flow.append(
                    {
                        "stage": f"{section_title} — {block_title}",
                        "teacher_actions": [],
                        "student_actions": [],
                        "activities": block_lines,
                        "assessment": [],
                        "reflection": [],
                    }
                )
            if not substeps:
                continue
            sections.append(
                {
                    "title": section_title,
                    "display_title": section_title,
                    "style": "Metodist Section",
                    "content_style": "Metodist Body",
                    "items": [],
                    "sample_item_styles": [],
                    "substeps": substeps,
                    "stages": [],
                    "children_order": [("substep", idx) for idx in range(len(substeps))],
                }
            )

        if not sections:
            sections = [
                {
                    "title": "Хід уроку",
                    "display_title": "Хід уроку",
                    "style": "Metodist Section",
                    "content_style": "Metodist Body",
                    "items": [],
                    "sample_item_styles": [],
                    "substeps": [
                        {
                            "title": "1. Робота за темою",
                            "display_title": "1. Робота за темою",
                            "style": "Metodist Substep",
                            "content_style": "Metodist Body",
                            "items": ["– Працюємо за темою уроку.", "(відповіді учнів)", "– Перевірмо разом."],
                            "sample_item_styles": [],
                            "substeps": [],
                            "children_order": [],
                        }
                    ],
                    "stages": [],
                    "children_order": [("substep", 0)],
                }
            ]

        return {
            "header_fields": header_fields,
            "sections": sections,
            "lesson_flow": lesson_flow,
            "topic": clean_text(header.get("topic") or topic),
            "grade": clean_text(grade),
            "subject": clean_text(subject),
        }

    async def _generate_lesson_script(
        self,
        *,
        topic,
        grade,
        requirements,
        subject,
        context,
        source_context,
        blueprint,
        reference_context,
        slide_reference_plan="",
        runtime_trace=None,
    ):
        outline = self._build_script_prompt_outline(blueprint, section_limit=8, block_limit=6)
        outline_json = json.dumps(outline, ensure_ascii=False, indent=2)
        style_cues_json = json.dumps(self._extract_style_cues(blueprint, limit=10), ensure_ascii=False, indent=2)
        prompt_reference_context = self._compress_context_for_prompt(reference_context, max_chars=1800, max_lines=40)
        prompt_source_context = self._compress_context_for_prompt(source_context, max_chars=2200, max_lines=55)
        slide_plan_block = slide_reference_plan or "Окрема прив'язка етапів до слайдів не задана."

        sys_instr = (
            "Ти досвідчений вчитель НУШ. Згенеруй готовий lesson-script у форматі JSON для проведення уроку. "
            "Пиши одразу у фінальному стилі конспекту: короткі репліки через '–', питання, конкретні завдання, "
            "місця для відповідей учнів у форматі '(відповіді учнів)'. "
            "Не використовуй ролі 'Вчитель/Учні/Діяльність', не пиши meta-описи етапу."
        )

        prompt = f"""Побудуй JSON у форматі:
{{
  "header": {{
    "topic": "...",
    "goal": "...",
    "equipment": "...",
    "lesson_type": "...",
    "expected_results": ["..."]
  }},
  "lesson_sections": [
    {{
      "title": "І. ...",
      "blocks": [
        {{
          "title": "Назва підпункту",
          "lines": [
            "– ...",
            "– ...?",
            "(відповіді учнів)",
            "– Перевірмо разом."
          ]
        }}
      ]
    }}
  ]
}}

ПРАВИЛА:
- не копіюй дослівно фрази з шаблонів;
- не перенось чужі назви ігор/вправ, якщо вони не підходять темі;
- кожний block має бути тематично прив'язаний до теми уроку;
- прибери інородні фрагменти з інших предметів;
- для математики обов'язково додай конкретні числа/вирази/задачі;
- без markdown, тільки валідний JSON.

ПРЕДМЕТ: {subject}
КЛАС: {grade}
ТЕМА: {topic}
ВИМОГИ: {requirements}
ДОДАТКОВИЙ КОНТЕКСТ: {context}

ОЧІКУВАНА СТРУКТУРА РОЗДІЛІВ З ШАБЛОНІВ:
{outline_json}

STYLE_CUES_FROM_TEMPLATES:
{style_cues_json}

{prompt_reference_context}

{prompt_source_context}

ПРИВ'ЯЗКА ЕТАПІВ ДО СЛАЙДІВ:
{slide_plan_block}
"""
        script_timeout_raw = os.getenv("GENERATOR_SCRIPT_TIMEOUT_SEC", "40")
        try:
            script_timeout_sec = float(script_timeout_raw)
        except ValueError:
            script_timeout_sec = 40.0
        response = await self._execute_model_call(
            call_name="script_first_payload",
            prompt=prompt,
            system_instruction=sys_instr,
            temperature=0.25,
            runtime_trace=runtime_trace,
            timeout_sec=script_timeout_sec,
            extra={
                "subject": clean_text(subject),
                "grade": clean_text(grade),
                "topic_chars": len(clean_text(topic)),
                "outline_sections": len(outline),
                "style_cues_count": len(json.loads(style_cues_json)),
                "compressed_source_context_chars": len(clean_text(prompt_source_context)),
                "compressed_reference_context_chars": len(clean_text(prompt_reference_context)),
            },
        )
        if not response:
            raise RuntimeError("script_model_call_failed")
        payload = self._extract_json_payload(getattr(response, "text", ""))
        raw_sections = payload.get("lesson_sections")
        if not isinstance(raw_sections, list):
            raw_sections = payload.get("sections")
        raw_payload_empty = not isinstance(raw_sections, list) or len(raw_sections) == 0
        script = self._coerce_lesson_script(payload, topic=topic, grade=grade, subject=subject, blueprint=blueprint)
        return script, {"raw_payload_empty": raw_payload_empty}

    async def _generate_pedagogical_plan(
        self,
        topic,
        grade,
        requirements,
        subject,
        context,
        source_context,
        blueprint,
        reference_context,
        slide_reference_plan="",
        runtime_trace=None,
    ):
        plan_template = self._build_generation_plan_template(blueprint)
        if not plan_template.get("stages"):
            return plan_template
        plan_reference_context = self._compress_context_for_prompt(reference_context, max_chars=2000, max_lines=40)
        plan_source_context = self._compress_context_for_prompt(source_context, max_chars=2600, max_lines=60)
        plan_template_json = json.dumps(plan_template, ensure_ascii=False, indent=2)
        style_cues_json = json.dumps(self._extract_style_cues(blueprint, limit=8), ensure_ascii=False, indent=2)

        slide_plan_block = slide_reference_plan or "Окрема прив'язка етапів до слайдів не задана."
        sys_instr = (
            "Ти методист НУШ. Побудуй pedagogical_plan у підході content-first: "
            "спершу створи живий stage_draft кожного етапу, потім розклади його по ролях. "
            "Пиши природно, з конкретними репліками, запитаннями, вправами та короткими переходами між етапами. "
            "Поверни тільки валідний JSON українською."
        )
        prompt = f"""Сформуй pedagogical_plan за шаблоном JSON:
{{
  "lesson_goal": "конкретна мета уроку",
  "expected_results": ["3-6 конкретних очікуваних результатів"],
  "stages": [
    {{
      "section_title": "назва розділу",
      "stage_title": "назва етапу",
      "stage_goal": "мета етапу",
      "stage_draft": [
        "4-8 живих рядків етапу: репліки вчителя, дії дітей, мікродіалоги, питання, вправа"
      ],
      "teacher_actions": ["2-5 конкретних дій/реплік учителя, виведених із stage_draft"],
      "student_actions": ["2-5 конкретних дій учнів, виведених із stage_draft"],
      "activities": ["2-5 конкретних вправ/завдань/прикладів/ігор, виведених із stage_draft"],
      "materials": ["1-3 потрібні матеріали"],
      "assessment": ["1-2 кроки перевірки/самооцінювання/рефлексії"]
    }}
  ]
}}

ЖОРСТКІ ПРАВИЛА:
- порядок і кількість stages мають бути ТОЧНО як у PLAN_TEMPLATE;
- stage_draft має бути живим, а не канцелярським: додавай звернення до дітей, мікродіалоги, конкретні інструкції;
- не пиши meta-описи формату: "є конкретні дії", "є запитання", "учитель коротко озвучує мету етапу";
- уникай порожніх формулювань: "учні виконують завдання", "вчитель пояснює", "проводиться рефлексія" без конкретики;
- обов'язково дай щонайменше одне конкретне запитання й одне конкретне завдання в кожному stage;
- для математики 3 класу додавай конкретні числа/вирази/короткі задачі;
- не копіюй факти з чужих тем із шаблонів; бери лише стиль і методичний тон;
- без markdown і без коментарів поза JSON.

ПРЕДМЕТ: {subject}
КЛАС: {grade}
ТЕМА: {topic}
ВИМОГИ: {requirements}
ДОДАТКОВИЙ КОНТЕКСТ: {context}

{plan_reference_context}

{plan_source_context}

ПРИВ'ЯЗКА ЕТАПІВ ДО СЛАЙДІВ:
{slide_plan_block}

STYLE_CUES_FROM_TEMPLATES (лише для тону й методики, без копіювання фактів):
{style_cues_json}

PLAN_TEMPLATE:
{plan_template_json}
"""

        try:
            plan_timeout_raw = os.getenv("GENERATOR_PLAN_TIMEOUT_SEC", "12")
            try:
                plan_timeout_sec = float(plan_timeout_raw)
            except ValueError:
                plan_timeout_sec = 12.0
            response = await self._execute_model_call(
                call_name="pedagogical_plan",
                prompt=prompt,
                system_instruction=sys_instr,
                temperature=0.3,
                runtime_trace=runtime_trace,
                timeout_sec=plan_timeout_sec,
                extra={
                    "subject": clean_text(subject),
                    "grade": clean_text(grade),
                    "topic_chars": len(clean_text(topic)),
                    "source_context_chars": len(clean_text(source_context)),
                    "reference_context_chars": len(clean_text(reference_context)),
                    "compressed_source_context_chars": len(clean_text(plan_source_context)),
                    "compressed_reference_context_chars": len(clean_text(plan_reference_context)),
                    "plan_template_chars": len(plan_template_json),
                    "style_cues_chars": len(style_cues_json),
                },
            )
            if not response:
                raise RuntimeError("model_call_failed")
            payload = self._extract_json_payload(getattr(response, "text", ""))
            return self._coerce_generation_plan(payload, blueprint, subject=subject, grade=grade)
        except Exception as exc:
            logger.error("Pedagogical plan generation failed: %s", exc)
            return plan_template

    def _build_payload_from_generation_plan(self, blueprint, generation_plan, topic, grade, subject):
        blueprint = blueprint if isinstance(blueprint, dict) else {}
        generation_plan = generation_plan if isinstance(generation_plan, dict) else {}
        stage_entries = generation_plan.get("stages") if isinstance(generation_plan.get("stages"), list) else []
        stage_index = {}
        for item in stage_entries:
            if not isinstance(item, dict):
                continue
            key = self._normalize_template_name(item.get("stage_title") or item.get("title") or item.get("stage") or "")
            if key and key not in stage_index:
                stage_index[key] = item

        def compose_stage_items(plan_stage):
            plan_stage = plan_stage if isinstance(plan_stage, dict) else {}
            items = []
            stage_draft = self._coerce_plan_items(plan_stage.get("stage_draft") or [], limit=10)
            stage_title = clean_text(plan_stage.get("stage_title") or "")
            if not stage_draft:
                stage_draft = self._build_stage_draft_from_fields(plan_stage, subject=subject, grade=grade)
            stage_draft = self._enrich_stage_draft_locally(
                stage_draft,
                stage_title=stage_title,
                subject=subject,
                grade=grade,
            )
            items.extend(stage_draft[:8])

            inferred = self._infer_stage_fields_from_items(items)
            teacher_actions = self._coerce_plan_items(
                [*(plan_stage.get("teacher_actions") or []), *(inferred.get("teacher_actions") or [])],
                limit=6,
            )
            student_actions = self._coerce_plan_items(
                [*(plan_stage.get("student_actions") or []), *(inferred.get("student_actions") or [])],
                limit=6,
            )
            activities = self._coerce_plan_items(
                [*(plan_stage.get("activities") or []), *(inferred.get("activities") or [])],
                limit=8,
            )
            assessment = self._coerce_plan_items(
                [
                    *(plan_stage.get("assessment") or []),
                    *(inferred.get("assessment") or []),
                    *(inferred.get("reflection") or []),
                ],
                limit=4,
            )
            goal = clean_text(plan_stage.get("stage_goal") or "")
            if goal:
                items.append(f"Мета: {goal}")
            for item in teacher_actions[:2]:
                if item not in items:
                    items.append(f"Вчитель: {item}")
            for item in student_actions[:2]:
                if item not in items:
                    items.append(f"Учні: {item}")
            for item in activities[:3]:
                if item not in items:
                    items.append(item if ":" in item else f"Завдання: {item}")
            for item in assessment[:2]:
                normalized_assessment = self._normalize_template_name(item)
                label = "Рефлексія" if "рефлекс" in normalized_assessment else "Оцінювання"
                items.append(f"{label}: {item}")
            return items

        def node_payload(node, kind):
            payload = {"items": [], "substeps": []}
            stage_items_cache = []
            stage_inferred_cache = {}
            plan_stage = {}

            if kind == "stage":
                title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
                key = self._normalize_template_name(title)
                plan_stage = stage_index.get(key, {})
                stage_items_cache = compose_stage_items(plan_stage)
                stage_inferred_cache = self._infer_stage_fields_from_items(stage_items_cache)
                payload["items"] = self._coerce_plan_items(stage_items_cache, limit=10)

            for substep in node.get("substeps") or []:
                sub_payload = node_payload(substep, "substep")
                if kind == "stage":
                    expected_field = self._expected_substep_field(substep)
                    mapped_items = []
                    if expected_field == "teacher_actions":
                        mapped_items = [*(plan_stage.get("teacher_actions") or []), *(stage_inferred_cache.get("teacher_actions") or [])]
                    elif expected_field == "student_actions":
                        mapped_items = [*(plan_stage.get("student_actions") or []), *(stage_inferred_cache.get("student_actions") or [])]
                    elif expected_field == "activities":
                        mapped_items = [*(plan_stage.get("activities") or []), *(stage_inferred_cache.get("activities") or [])]
                    elif expected_field == "assessment":
                        mapped_items = [*(plan_stage.get("assessment") or []), *(stage_inferred_cache.get("assessment") or [])]
                    elif expected_field == "reflection":
                        mapped_items = [*(plan_stage.get("assessment") or []), *(stage_inferred_cache.get("reflection") or [])]
                    elif expected_field == "materials":
                        mapped_items = plan_stage.get("materials") or []
                    if mapped_items:
                        sub_payload["items"] = self._coerce_plan_items(mapped_items, limit=8)
                payload["substeps"].append(sub_payload)

            if "stages" in node:
                payload["stages"] = []
                for stage in node.get("stages") or []:
                    payload["stages"].append(node_payload(stage, "stage"))

            return payload

        payload = {
            "header_values": {},
            "sections": [],
        }
        for section in blueprint.get("sections") or []:
            payload["sections"].append(node_payload(section, "section"))
        return payload

    async def _generate_blueprint_content(
        self,
        topic,
        grade,
        requirements,
        subject,
        context,
        source_context,
        blueprint,
        reference_context,
        slide_reference_plan="",
        generation_plan=None,
        runtime_trace=None,
    ):
        prompt_blueprint = {
            "header_fields": [
                {
                    "label": field.get("label"),
                    "sample_value": self._truncate_example(field.get("sample_value") or "", 220),
                }
                for field in blueprint.get("header_fields", [])
            ],
            "sections": [self._build_prompt_node(section) for section in blueprint.get("sections", [])],
        }
        normalized_plan = self._coerce_generation_plan(generation_plan or {}, blueprint, subject=subject, grade=grade)
        compact_plan_outline = self._build_generation_plan_prompt_outline(normalized_plan, stage_limit=14)
        prompt_blueprint_json = json.dumps(prompt_blueprint, ensure_ascii=False, indent=2)
        compact_plan_json = json.dumps(compact_plan_outline, ensure_ascii=False, indent=2)
        style_cues_json = json.dumps(self._extract_style_cues(blueprint, limit=8), ensure_ascii=False, indent=2)
        prompt_reference_context = self._compress_context_for_prompt(reference_context, max_chars=2200, max_lines=45)
        prompt_source_context = self._compress_context_for_prompt(source_context, max_chars=3200, max_lines=80)
        slide_plan_block = slide_reference_plan or "Окрема прив'язка етапів до слайдів не задана."

        sys_instr = (
            "Ти методист НУШ. Крок 1 вже виконано: є pedagogical_plan. "
            "Крок 2: перенеси зміст plan у blueprint JSON без зміни структури. "
            "Пиши змістовно, конкретно, без води. Кожен рядок items = окрема дія, репліка, запитання або вправа. "
            "Зберігай порядок sections/stages/substeps один в один. Поверни лише валідний JSON українською."
        )

        prompt = f"""ФОРМАТ ВІДПОВІДІ (JSON):
{{
  "header_values": {{"Назва поля": "нове значення"}},
  "sections": [
    {{
      "items": ["конкретний змістовний рядок"],
      "substeps": [{{"items": ["конкретний змістовний рядок"]}}],
      "stages": [{{"items": ["конкретний змістовний рядок"], "substeps": [{{"items": ["конкретний змістовний рядок"]}}]}}]
    }}
  ]
}}

ОБОВ'ЯЗКОВІ ВИМОГИ ДО ЗМІСТУ:
- не змінюй структуру blueprint, лише header_values і items;
- у кожному stage має бути щонайменше 4-8 змістовних рядків сумарно (stage items + substeps);
- у кожному stage додай:
  1) конкретну дію/репліку вчителя;
  2) конкретну дію учнів;
  3) конкретне завдання, вправу, приклад або запитання;
- stage.items пиши як живий фрагмент уроку: короткі репліки, переходи, міні-діалоги, запитання;
- заборонено meta-коментарі типу: "є конкретні дії", "є математичний сигнал", "є елемент перевірки";
- явні префікси ролей ("Вчитель:", "Учні:") використовуй лише коли це справді допомагає, не роби ними весь текст;
- уникай загальних фраз: "провести обговорення", "ознайомити учнів", "виконати завдання", "провести рефлексію", "відбувається", "проводиться", "здійснюється";
- додавай 1-2 елементи формувального оцінювання/рефлексії там, де це доречно;
- якщо предмет = математика (3 клас), додай конкретні числа/вирази/короткі задачі;
- якщо в node є source_items, використовуй їх як опорні факти, а не як шаблон довжини;
- без markdown і без коментарів поза JSON.

ПРЕДМЕТ: {subject}
КЛАС: {grade}
ТЕМА: {topic}
ВИМОГИ: {requirements}
ДОДАТКОВИЙ КОНТЕКСТ: {context}

{prompt_reference_context}

{prompt_source_context}

ПРИВ'ЯЗКА ЕТАПІВ ДО СЛАЙДІВ:
{slide_plan_block}

STYLE_CUES_FROM_TEMPLATES (лише для тону, без копіювання фактів):
{style_cues_json}

PEDAGOGICAL_PLAN (використай як головну смислову основу):
{compact_plan_json}

BLUEPRINT (заповни, не змінюючи структуру):
{prompt_blueprint_json}
"""

        try:
            payload_timeout_raw = os.getenv("GENERATOR_PAYLOAD_TIMEOUT_SEC", "55")
            try:
                payload_timeout_sec = float(payload_timeout_raw)
            except ValueError:
                payload_timeout_sec = 55.0
            response = await self._execute_model_call(
                call_name="strict_blueprint_payload",
                prompt=prompt,
                system_instruction=sys_instr,
                temperature=0.25,
                runtime_trace=runtime_trace,
                timeout_sec=payload_timeout_sec,
                extra={
                    "subject": clean_text(subject),
                    "grade": clean_text(grade),
                    "topic_chars": len(clean_text(topic)),
                    "source_context_chars": len(clean_text(source_context)),
                    "reference_context_chars": len(clean_text(reference_context)),
                    "compressed_source_context_chars": len(clean_text(prompt_source_context)),
                    "compressed_reference_context_chars": len(clean_text(prompt_reference_context)),
                    "blueprint_json_chars": len(prompt_blueprint_json),
                    "plan_outline_chars": len(compact_plan_json),
                    "style_cues_chars": len(style_cues_json),
                    "plan_stages_full": len((normalized_plan or {}).get("stages") or []),
                    "plan_stages_outline": len((compact_plan_outline or {}).get("stages") or []),
                    "stages_in_blueprint": len(self._collect_blueprint_stages(blueprint)),
                },
            )
            if not response:
                raise RuntimeError("model_call_failed")
            return self._extract_json_payload(getattr(response, "text", ""))
        except Exception as exc:
            logger.error("Strict example generation failed: %s", exc)
            return {}

    def _coerce_items(self, items, sample_items=None):
        cleaned = self._sanitize_source_items(items or [], limit=28)
        if cleaned:
            sample_count = len(sample_items or [])
            limit = max(10, min(28, sample_count + 14 if sample_count else 16))
            return cleaned[:limit]

        fallback_items = self._sanitize_source_items(sample_items or [], limit=12)
        if fallback_items:
            return fallback_items[: min(6, len(fallback_items))]
        return []

    def _expected_substep_field(self, node):
        label = self._normalize_template_name(
            (node or {}).get("title")
            or (node or {}).get("display_title")
            or ""
        )
        return self._resolve_stage_field_from_label(label)

    @staticmethod
    def _repair_mojibake_utf8(value):
        text = clean_text(value)
        if not text:
            return ""
        best = text
        best_score = -1
        for encoding in ("cp1251", "cp866", "cp1252", "latin1"):
            try:
                repaired = text.encode(encoding).decode("utf-8")
            except Exception:
                continue
            if not repaired:
                continue
            cyr_score = len(re.findall(r"[А-Яа-яІіЇїЄєҐґ]", repaired))
            mojibake_penalty = len(re.findall(r"[ÐÑÃÂ]", repaired))
            score = cyr_score - mojibake_penalty
            if score > best_score:
                best = repaired
                best_score = score
        return best

    def _normalize_for_matching(self, value):
        normalized = self._normalize_template_name(value)
        repaired = self._repair_mojibake_utf8(normalized)
        if repaired != normalized:
            candidate = self._normalize_template_name(repaired)
            if candidate:
                return candidate
        return normalized

    def _resolve_stage_field_from_label(self, label):
        normalized = self._normalize_for_matching(label)
        if not normalized:
            return None

        direct_match = self.EXPLICIT_STAGE_LABEL_MAP.get(normalized)
        if direct_match:
            return direct_match

        canonical_map = {
            "мета": "goal",
            "ціль": "goal",
            "goal": "goal",
            "вчитель": "teacher_actions",
            "учитель": "teacher_actions",
            "teacher": "teacher_actions",
            "учні": "student_actions",
            "учень": "student_actions",
            "students": "student_actions",
            "student": "student_actions",
            "діяльність": "activities",
            "активність": "activities",
            "activities": "activities",
            "activity": "activities",
            "оцінювання": "assessment",
            "самооцінювання": "assessment",
            "взаємооцінювання": "assessment",
            "assessment": "assessment",
            "рефлексія": "reflection",
            "reflection": "reflection",
            "матеріали": "materials",
            "ресурси": "materials",
            "materials": "materials",
            "диференціація": "differentiation",
            "differentiation": "differentiation",
        }
        if normalized in canonical_map:
            return canonical_map[normalized]

        if normalized.startswith(("вчител", "учител", "teacher")):
            return "teacher_actions"
        if normalized.startswith(("учн", "student")):
            return "student_actions"
        if normalized.startswith(("діяльн", "активн", "activity")):
            return "activities"
        if normalized.startswith(("оцін", "assessment", "перевір")):
            return "assessment"
        if normalized.startswith(("рефлекс", "reflection")):
            return "reflection"
        if normalized.startswith(("матеріал", "ресурс", "material")):
            return "materials"
        if normalized.startswith(("мета", "ціль", "goal")):
            return "goal"
        if normalized.startswith(("диферен", "differentiation")):
            return "differentiation"
        return None

    def _is_over_generic_item(self, text, expected_field=None):
        normalized = self._normalize_template_name(text)
        if not normalized:
            return True

        words_count = len(normalized.split())
        if expected_field in {"teacher_actions", "student_actions", "activities"} and words_count > 50:
            return True
        if words_count > 70:
            return True

        return False

    def _normalize_blueprint_item_text(self, value, expected_field=None):
        text = normalize_sentence_punctuation(value)
        if not text or is_placeholder_text(text):
            return ""

        field_name, explicit_value = self._split_explicit_stage_label(text)
        if field_name:
            if expected_field:
                if field_name != expected_field:
                    return ""
                text = normalize_sentence_punctuation(explicit_value)
            else:
                # Preserve explicit role label on stage-level items.
                text = normalize_sentence_punctuation(text)

        if not text or is_placeholder_text(text):
            return ""
        if self._is_over_generic_item(text, expected_field=expected_field):
            return ""
        if self._is_generic_lesson_item(text) and len(text.split()) <= 4:
            return ""
        return text

    def _unique_clean_items(self, items, expected_field=None, limit=12):
        cleaned = []
        seen = set()
        for item in normalize_block_list(items):
            normalized = self._normalize_blueprint_item_text(item, expected_field=expected_field)
            if not normalized:
                continue
            key = self._normalize_template_name(normalized)
            if not key or key in seen:
                continue
            seen.add(key)
            cleaned.append(normalized)
            if len(cleaned) >= limit:
                break
        return cleaned

    def _build_topic_anchor_items(self, topic, node_title, expected_field=None, limit=2, subject="", grade=""):
        topic_text = self._safe_display_text(topic)
        node_text = self._safe_display_text(node_title)
        if not topic_text:
            return []

        is_math = self._is_math_subject(subject)
        if expected_field == "teacher_actions":
            candidates = [
                f"Учитель: «Тема уроку — {topic_text}. Починаємо з короткого усного завдання.»",
                "Учитель дає інструкцію: «Виконуємо по кроках, після цього звіряємо в парі».",
            ]
        elif expected_field == "student_actions":
            candidates = [
                "Учні працюють у парі: один називає приклад, інший пояснює відповідь.",
                "Кілька учнів озвучують розв'язання, клас уточнює або виправляє.",
            ]
        elif expected_field == "activities":
            candidates = [
                f"Обчисліть коротко за темою «{topic_text}» і поясніть спосіб розв'язання.",
                "Попрацюйте в парах: взаємоперевірка однієї відповіді за зразком.",
            ]
            if is_math:
                candidates.extend(self._build_math_micro_tasks(grade, limit=2))
        else:
            title_part = node_text or "етап уроку"
            candidates = [
                f"Етап «{title_part}»: виконайте коротке завдання й поясніть відповідь однокласнику.",
                "Учитель ставить одне запитання на розуміння, діти відповідають усно по черзі.",
                "Наприкінці етапу одна коротка перевірка: «правильно / треба підказка».",
            ]
            if is_math:
                candidates.extend(self._build_math_micro_tasks(grade, limit=2))

        return self._unique_clean_items(candidates, expected_field=expected_field, limit=limit)

    def _stabilize_blueprint_node(self, blueprint_node, model_node, kind, topic=""):
        blueprint_node = blueprint_node if isinstance(blueprint_node, dict) else {}
        model_node = model_node if isinstance(model_node, dict) else {}

        expected_field = self._expected_substep_field(blueprint_node) if kind == "substep" else None
        source_items = normalize_block_list(blueprint_node.get("source_items") or [])
        item_limit = 14 if kind in {"section", "stage"} else 10
        min_items = 0 if (blueprint_node.get("stages") or blueprint_node.get("substeps")) else (2 if source_items else 1)
        if kind == "substep" and not source_items:
            substep_title = clean_text_preserve_prefix(
                blueprint_node.get("display_title") or blueprint_node.get("title") or ""
            )
            if not substep_title:
                min_items = 0

        items = self._unique_clean_items(model_node.get("items") or [], expected_field=expected_field, limit=item_limit)
        if len(items) < min_items:
            fallback_cleaned = self._unique_clean_items(
                source_items,
                expected_field=expected_field,
                limit=item_limit,
            )
            for candidate in fallback_cleaned:
                if candidate in items:
                    continue
                items.append(candidate)
                if len(items) >= min_items:
                    break
        if len(items) < min_items:
            anchor_items = self._build_topic_anchor_items(
                topic=topic,
                node_title=blueprint_node.get("display_title") or blueprint_node.get("title") or "",
                expected_field=expected_field,
                limit=item_limit,
            )
            for candidate in anchor_items:
                if candidate in items:
                    continue
                items.append(candidate)
                if len(items) >= min_items:
                    break
        items = items[:item_limit]

        model_substeps = model_node.get("substeps") if isinstance(model_node.get("substeps"), list) else []
        model_stages = model_node.get("stages") if isinstance(model_node.get("stages"), list) else []

        stabilized_substeps = []
        for index, substep_blueprint in enumerate(blueprint_node.get("substeps") or []):
            source_substep = model_substeps[index] if index < len(model_substeps) else {}
            stabilized_substeps.append(
                self._stabilize_blueprint_node(substep_blueprint, source_substep, "substep", topic=topic)
            )

        stabilized_stages = []
        for index, stage_blueprint in enumerate(blueprint_node.get("stages") or []):
            source_stage = model_stages[index] if index < len(model_stages) else {}
            stabilized_stages.append(
                self._stabilize_blueprint_node(stage_blueprint, source_stage, "stage", topic=topic)
            )

        children_order = list(blueprint_node.get("children_order") or [])
        if not children_order:
            children_order = [("substep", index) for index in range(len(stabilized_substeps))]
            children_order.extend(("stage", index) for index in range(len(stabilized_stages)))

        stabilized = {
            "title": blueprint_node.get("title") or model_node.get("title") or "",
            "display_title": blueprint_node.get("display_title") or blueprint_node.get("title") or model_node.get("display_title") or model_node.get("title") or "",
            "style": blueprint_node.get("style") or model_node.get("style") or "Normal",
            "content_style": blueprint_node.get("content_style") or model_node.get("content_style") or "Normal",
            "items": items,
            "sample_items": blueprint_node.get("sample_items") or [],
            "source_items": blueprint_node.get("source_items") or [],
            "sample_item_styles": blueprint_node.get("sample_item_styles") or [],
            "substeps": stabilized_substeps,
            "children_order": children_order,
        }
        if "stages" in blueprint_node:
            stabilized["stages"] = stabilized_stages
        return stabilized

    def _stabilize_blueprint_document_model(self, blueprint, document_model):
        blueprint = blueprint if isinstance(blueprint, dict) else {}
        document_model = document_model if isinstance(document_model, dict) else {}
        model_sections = document_model.get("sections") if isinstance(document_model.get("sections"), list) else []
        topic = clean_text(document_model.get("topic") or "")
        if not topic:
            for field in document_model.get("header_fields") or []:
                if clean_text(field.get("label") or "").lower() == "тема":
                    topic = clean_text(field.get("value") or "")
                    if topic:
                        break

        stabilized_sections = []
        for index, section_blueprint in enumerate(blueprint.get("sections") or []):
            section_model = model_sections[index] if index < len(model_sections) else {}
            stabilized_sections.append(
                self._stabilize_blueprint_node(section_blueprint, section_model, "section", topic=topic)
            )

        stabilized_headers = []
        for field in document_model.get("header_fields") or []:
            if not isinstance(field, dict):
                continue
            label = clean_text(field.get("label") or "")
            value = normalize_sentence_punctuation(field.get("value") or "")
            if not label:
                continue
            stabilized_headers.append(
                {
                    "label": label,
                    "style": clean_text(field.get("style") or "Normal") or "Normal",
                    "value": value,
                }
            )

        stabilized = {
            **document_model,
            "header_fields": stabilized_headers,
            "sections": stabilized_sections,
            "lesson_flow": [],
        }
        for section in stabilized_sections:
            for stage in section.get("stages") or []:
                stabilized["lesson_flow"].append(self._stage_node_to_flow_entry(stage))
        return stabilized

    @staticmethod
    def _count_document_items(document_model):
        def count_node(node):
            total = len(normalize_block_list(node.get("items") or []))
            for substep in node.get("substeps") or []:
                total += count_node(substep)
            for stage in node.get("stages") or []:
                total += count_node(stage)
            return total

        total = 0
        for section in document_model.get("sections") or []:
            total += count_node(section)
        return total

    @staticmethod
    def _count_rich_lesson_items(lesson_data):
        total = 0
        for key in (
            "goal",
            "tasks",
            "key_competencies",
            "cross_cutting_skills",
            "values",
            "integration",
            "methods",
            "forms",
            "assessment",
            "differentiation",
            "equipment",
            "resources",
            "homework",
        ):
            total += len(normalize_list((lesson_data or {}).get(key) or []))

        expected = normalize_expected_results((lesson_data or {}).get("expected_results") or {})
        total += len(expected.get("knowledge") or [])
        total += len(expected.get("skills") or [])
        total += len(expected.get("values") or [])

        for stage in (lesson_data or {}).get("lesson_flow") or []:
            total += len(normalize_list(stage.get("teacher_actions") or []))
            total += len(normalize_list(stage.get("student_actions") or []))
            total += len(normalize_list(stage.get("activities") or []))
        return total

    @staticmethod
    def _node_item_count(node):
        if not isinstance(node, dict):
            return 0
        total = len(normalize_block_list(node.get("items") or []))
        for substep in node.get("substeps") or []:
            total += LessonGenerator._node_item_count(substep)
        for stage in node.get("stages") or []:
            total += LessonGenerator._node_item_count(stage)
        return total

    def _is_generic_lesson_item(self, text):
        cleaned = clean_text_preserve_prefix(text)
        if not cleaned:
            return True
        normalized = self._normalize_template_name(cleaned)
        if not normalized:
            return True
        if self._is_template_noise_line(normalized):
            return True
        for phrase in self.WEAK_GENERIC_PATTERNS:
            if phrase in normalized:
                return True
        words = normalized.split()
        if len(words) <= 2:
            return True
        if len(words) <= 4 and "?" not in cleaned and not re.search(r"\d", cleaned):
            return True
        return False

    def _has_task_signal(self, text):
        cleaned = clean_text_preserve_prefix(text)
        normalized = self._normalize_template_name(cleaned)
        if "?" in cleaned:
            return True
        return any(token in normalized for token in self.TASK_SIGNAL_TOKENS)

    def _has_student_signal(self, text):
        normalized = self._normalize_for_matching(text)
        return self._looks_like_student_line(text) or "\u0443\u0447\u043d" in normalized or "student" in normalized

    def _assess_blueprint_quality(self, blueprint, document_model):
        metrics = {
            "total_items": self._count_document_items(document_model),
            "section_total": 0,
            "section_filled": 0,
            "stage_total": 0,
            "stage_filled": 0,
            "stage_with_source": 0,
            "stage_source_covered": 0,
            "content_density_score": 0.0,
            "stage_completeness_score": 0.0,
            "anti_generic_score": 0.0,
            "student_activity_score": 0.0,
            "teacher_activity_score": 0.0,
            "assessment_coverage_score": 0.0,
            "pedagogically_collapsed": False,
            "weak_nodes": [],
            "blocking_weak_nodes": [],
            "cosmetic_weak_nodes": [],
            "blocking_weak_nodes_count": 0,
            "cosmetic_weak_nodes_count": 0,
            "only_cosmetic_weakness": False,
            "acceptance_fail_reasons": [],
            "score": 0.0,
            "acceptable": False,
        }

        sections_blueprint = blueprint.get("sections") or []
        sections_model = document_model.get("sections") or []
        subject = clean_text((document_model or {}).get("subject") or "")
        grade = clean_text((document_model or {}).get("grade") or "")
        stage_details = []

        def push_weak(kind, title, path, section_index, trace, issues):
            issues = [clean_text_preserve_prefix(item) for item in (issues or []) if clean_text_preserve_prefix(item)]
            node = {
                "kind": kind,
                "title": clean_text_preserve_prefix(title) or path,
                "path": path,
                "section_index": section_index,
                "trace": list(trace or []),
                "issues": issues,
            }
            key = (node["kind"], node["path"])
            if key in {(item.get("kind"), item.get("path")) for item in metrics["weak_nodes"]}:
                return
            metrics["weak_nodes"].append(node)

        def is_cosmetic_weak_node(node):
            if not isinstance(node, dict):
                return False
            kind = clean_text(node.get("kind") or "").lower()
            title = clean_text_preserve_prefix(node.get("title") or "")
            path = clean_text_preserve_prefix(node.get("path") or "")
            issues = [
                self._normalize_template_name(item)
                for item in (node.get("issues") or [])
                if clean_text_preserve_prefix(item)
            ]
            mojibake_title = self._looks_like_mojibake_text(title) or self._looks_like_mojibake_text(path)
            if not mojibake_title:
                return False
            if not issues:
                return True
            if kind == "substep" and all("РїРѕСЂРѕР¶РЅС–Р№ РїС–РґРїСѓРЅРєС'" in issue for issue in issues):
                return True
            cosmetic_markers = ("назв", "заголов", "лейбл", "кодиров", "mojibake")
            return all(any(marker in issue for marker in cosmetic_markers) for issue in issues)

        def walk(bp_node, model_node, kind, path, section_index, trace):
            bp_node = bp_node if isinstance(bp_node, dict) else {}
            model_node = model_node if isinstance(model_node, dict) else {}

            item_count = self._node_item_count(model_node)
            direct_items = normalize_block_list(model_node.get("items") or [])
            has_source = bool(bp_node.get("source_items"))
            display_title = clean_text_preserve_prefix(
                bp_node.get("display_title") or bp_node.get("title") or path
            )

            if kind == "section":
                metrics["section_total"] += 1
                if item_count >= 2 or (bp_node.get("stages") or bp_node.get("substeps")):
                    metrics["section_filled"] += 1
                if item_count < 2 and not (bp_node.get("stages") or bp_node.get("substeps")):
                    push_weak(kind, display_title, path, section_index, trace, ["замало змісту"])

            if kind == "stage":
                metrics["stage_total"] += 1
                if item_count >= 3:
                    metrics["stage_filled"] += 1
                if has_source:
                    metrics["stage_with_source"] += 1
                    if item_count >= 3:
                        metrics["stage_source_covered"] += 1

                stage_entry = self._stage_node_to_flow_entry(model_node)
                teacher_actions = normalize_list(stage_entry.get("teacher_actions") or [])
                student_actions = normalize_list(stage_entry.get("student_actions") or [])
                activities = normalize_list(stage_entry.get("activities") or [])
                assessment_items = normalize_list(stage_entry.get("assessment") or [])
                assessment_items.extend(normalize_list(stage_entry.get("reflection") or []))
                stage_model_substeps = model_node.get("substeps") or []
                for sub_index, substep_bp in enumerate(bp_node.get("substeps") or []):
                    if sub_index >= len(stage_model_substeps):
                        continue
                    model_substep = stage_model_substeps[sub_index] if isinstance(stage_model_substeps[sub_index], dict) else {}
                    sub_items = normalize_block_list(model_substep.get("items") or [])
                    expected_field = self._expected_substep_field(substep_bp)
                    if not expected_field and len(bp_node.get("substeps") or []) >= 3:
                        if sub_index == 0:
                            expected_field = "teacher_actions"
                        elif sub_index == 1:
                            expected_field = "student_actions"
                        elif sub_index == 2:
                            expected_field = "activities"
                    if expected_field == "teacher_actions":
                        teacher_actions.extend(sub_items)
                    elif expected_field == "student_actions":
                        student_actions.extend(sub_items)
                    elif expected_field == "activities":
                        activities.extend(sub_items)
                    elif expected_field in {"assessment", "reflection"}:
                        assessment_items.extend(sub_items)
                teacher_actions = self._unique_clean_items(teacher_actions, expected_field="teacher_actions", limit=16)
                student_actions = self._unique_clean_items(student_actions, expected_field="student_actions", limit=16)
                activities = self._unique_clean_items(activities, expected_field="activities", limit=18)
                assessment_items = self._unique_clean_items(assessment_items, expected_field="assessment", limit=10)

                all_lines = []
                all_lines.extend(normalize_block_list(direct_items))
                for substep in model_node.get("substeps") or []:
                    all_lines.extend(normalize_block_list((substep or {}).get("items") or []))
                all_lines.extend(teacher_actions)
                all_lines.extend(student_actions)
                all_lines.extend(activities)
                all_lines = [clean_text_preserve_prefix(line) for line in all_lines if clean_text_preserve_prefix(line)]

                generic_count = sum(1 for line in all_lines if self._is_generic_lesson_item(line))
                report_like_count = sum(
                    1
                    for line in all_lines
                    if any(self._normalize_template_name(line).startswith(prefix) for prefix in self.REPORT_LABEL_PREFIXES)
                )
                question_count = sum(1 for line in all_lines if "?" in line)
                number_signal_count = sum(1 for line in all_lines if re.search(r"\d", line))
                task_signal = any(self._has_task_signal(line) for line in all_lines)
                student_signal = bool(student_actions) or any(self._has_student_signal(line) for line in all_lines)
                teacher_signal = bool(teacher_actions)
                assessment_signal = bool(assessment_items) or any(
                    self._looks_like_assessment_line(line) or self._looks_like_reflection_line(line)
                    for line in all_lines
                )
                assessment_relevant = self._classify_stage_name(display_title) in {"main", "closing"}
                complete_stage = bool(teacher_signal and student_signal and task_signal)
                avg_words = (
                    sum(len(clean_text(line).split()) for line in all_lines) / len(all_lines)
                    if all_lines
                    else 0.0
                )

                stage_details.append(
                    {
                        "title": display_title,
                        "item_count": item_count,
                        "line_count": len(all_lines),
                        "avg_words": avg_words,
                        "generic_count": generic_count,
                        "report_like_count": report_like_count,
                        "question_count": question_count,
                        "number_signal_count": number_signal_count,
                        "teacher_signal": teacher_signal,
                        "student_signal": student_signal,
                        "assessment_signal": assessment_signal,
                        "assessment_relevant": assessment_relevant,
                        "task_signal": task_signal,
                        "complete": complete_stage,
                    }
                )

                issues = []
                if item_count < 3:
                    issues.append("етап занадто короткий")
                if not teacher_signal:
                    issues.append("немає конкретних дій вчителя")
                if not student_signal:
                    issues.append("немає конкретних дій учнів")
                if not task_signal:
                    issues.append("немає конкретного завдання/питання")
                if assessment_relevant and not assessment_signal:
                    issues.append("немає формувального оцінювання/рефлексії")
                if all_lines and (generic_count / len(all_lines)) > 0.45:
                    issues.append("забагато загальних формулювань")
                if all_lines and (report_like_count / len(all_lines)) > 0.72 and question_count == 0:
                    issues.append("етап виглядає як формальний звіт, бракує живих реплік")
                if has_source and not direct_items and len(all_lines) < 3:
                    issues.append("source_items не розкрито в змісті")
                if self._is_math_subject(subject) and self._grade_to_int(grade) == 3 and number_signal_count == 0:
                    issues.append("для математики 3 класу бракує конкретних чисел/виразів")
                if issues:
                    push_weak("stage", display_title, path, section_index, trace, issues)

            bp_substeps = bp_node.get("substeps") or []
            model_substeps = model_node.get("substeps") or []
            for index, substep_bp in enumerate(bp_substeps):
                substep_model = model_substeps[index] if index < len(model_substeps) else {}
                substep_title = clean_text_preserve_prefix(substep_bp.get("display_title") or substep_bp.get("title") or "")
                substep_count = self._node_item_count(substep_model)
                substep_has_source = bool(substep_bp.get("source_items"))
                substep_path = f"{path} > {substep_title or index + 1}"
                substep_trace = [*trace, {"key": "substeps", "index": index}]
                if substep_count == 0 and (substep_has_source or len(bp_substeps) <= 3):
                    push_weak(
                        "substep",
                        substep_title or f"{display_title} / РїС–РґРїСѓРЅРєС' {index + 1}",
                        substep_path,
                        section_index,
                        substep_trace,
                        ["РїРѕСЂРѕР¶РЅС–Р№ РїС–РґРїСѓРЅРєС'"],
                    )
                walk(substep_bp, substep_model, "substep", substep_path, section_index, substep_trace)

            bp_stages = bp_node.get("stages") or []
            model_stages = model_node.get("stages") or []
            for index, stage_bp in enumerate(bp_stages):
                stage_model = model_stages[index] if index < len(model_stages) else {}
                stage_title = clean_text_preserve_prefix(stage_bp.get("display_title") or stage_bp.get("title") or "")
                stage_path = f"{path} > {stage_title or index + 1}"
                stage_trace = [*trace, {"key": "stages", "index": index}]
                walk(stage_bp, stage_model, "stage", stage_path, section_index, stage_trace)

        for index, section_bp in enumerate(sections_blueprint):
            section_model = sections_model[index] if index < len(sections_model) else {}
            section_title = clean_text_preserve_prefix(section_bp.get("display_title") or section_bp.get("title") or "")
            walk(section_bp, section_model, "section", section_title or f"section {index + 1}", index, [])

        stage_coverage = (
            metrics["stage_filled"] / metrics["stage_total"]
            if metrics["stage_total"]
            else 1.0
        )
        source_coverage = (
            metrics["stage_source_covered"] / metrics["stage_with_source"]
            if metrics["stage_with_source"]
            else 1.0
        )
        section_coverage = (
            metrics["section_filled"] / metrics["section_total"]
            if metrics["section_total"]
            else 1.0
        )

        if stage_details:
            avg_stage_items = sum(item.get("item_count", 0) for item in stage_details) / len(stage_details)
            avg_stage_words = sum(item.get("avg_words", 0.0) for item in stage_details) / len(stage_details)
            total_lines = sum(item.get("line_count", 0) for item in stage_details)
            total_generic = sum(item.get("generic_count", 0) for item in stage_details)
            total_report_like = sum(item.get("report_like_count", 0) for item in stage_details)
            total_questions = sum(item.get("question_count", 0) for item in stage_details)
            total_number_signals = sum(item.get("number_signal_count", 0) for item in stage_details)
            complete_stages = sum(1 for item in stage_details if item.get("complete"))
            student_stages = sum(1 for item in stage_details if item.get("student_signal"))
            teacher_stages = sum(1 for item in stage_details if item.get("teacher_signal"))
            assessment_relevant = sum(1 for item in stage_details if item.get("assessment_relevant"))
            assessment_stages = sum(
                1 for item in stage_details
                if item.get("assessment_relevant") and item.get("assessment_signal")
            )

            content_density_score = min(1.0, ((avg_stage_items / 5.0) * 0.65) + ((avg_stage_words / 9.0) * 0.35))
            stage_completeness_score = complete_stages / len(stage_details)
            anti_generic_score = max(0.0, 1.0 - ((total_generic / total_lines) if total_lines else 0.65))
            student_activity_score = student_stages / len(stage_details)
            teacher_activity_score = teacher_stages / len(stage_details)
            assessment_coverage_score = (
                assessment_stages / assessment_relevant
                if assessment_relevant
                else 1.0
            )
            report_penalty = max(0.0, 1.0 - ((total_report_like / total_lines) if total_lines else 0.0))
            question_bonus = min(1.0, total_questions / max(1, len(stage_details)))
            if self._is_math_subject(subject) and self._grade_to_int(grade) == 3:
                number_bonus = min(1.0, total_number_signals / max(1, len(stage_details)))
                anti_generic_score = max(0.0, anti_generic_score * 0.8 + number_bonus * 0.2)
            content_density_score = max(0.0, min(1.0, content_density_score * 0.9 + question_bonus * 0.1))
            stage_completeness_score = max(0.0, min(1.0, stage_completeness_score * 0.88 + report_penalty * 0.12))
        else:
            content_density_score = 0.0
            stage_completeness_score = 0.0
            anti_generic_score = 0.0
            student_activity_score = 0.0
            teacher_activity_score = 0.0
            assessment_coverage_score = 0.0

        metrics["content_density_score"] = round(content_density_score, 4)
        metrics["stage_completeness_score"] = round(stage_completeness_score, 4)
        metrics["anti_generic_score"] = round(anti_generic_score, 4)
        metrics["student_activity_score"] = round(student_activity_score, 4)
        metrics["teacher_activity_score"] = round(teacher_activity_score, 4)
        metrics["assessment_coverage_score"] = round(assessment_coverage_score, 4)

        minimum_expected_stages = 4
        teacher_missing_majority = teacher_activity_score < 0.5
        student_missing_majority = student_activity_score < 0.5
        assessment_missing_majority = assessment_coverage_score < 0.4
        if metrics["stage_total"] < minimum_expected_stages:
            metrics["pedagogically_collapsed"] = True
            push_weak(
                "document",
                "Конспект уроку",
                "lesson_flow",
                -1,
                [],
                [f"замало етапів уроку: {metrics['stage_total']} (< {minimum_expected_stages})"],
            )
        if teacher_missing_majority:
            metrics["pedagogically_collapsed"] = True
            push_weak(
                "document",
                "Конспект уроку",
                "lesson_flow",
                -1,
                [],
                ["у більшості етапів немає конкретних дій вчителя"],
            )
        if student_missing_majority:
            metrics["pedagogically_collapsed"] = True
            push_weak(
                "document",
                "Конспект уроку",
                "lesson_flow",
                -1,
                [],
                ["у більшості етапів немає конкретних дій учнів"],
            )
        if assessment_missing_majority:
            metrics["pedagogically_collapsed"] = True
            push_weak(
                "document",
                "Конспект уроку",
                "lesson_flow",
                -1,
                [],
                ["у більшості релевантних етапів відсутні оцінювання/рефлексія"],
            )

        score = (
            min(metrics["total_items"] / max(18, metrics["stage_total"] * 5), 1.0) * 0.22
            + stage_coverage * 0.18
            + source_coverage * 0.10
            + section_coverage * 0.03
            + content_density_score * 0.16
            + stage_completeness_score * 0.14
            + anti_generic_score * 0.09
            + student_activity_score * 0.08
        )
        collapse_penalty = 1.0
        if metrics["stage_total"] < minimum_expected_stages:
            collapse_penalty *= 0.45
        if teacher_missing_majority:
            collapse_penalty *= 0.7
        if student_missing_majority:
            collapse_penalty *= 0.75
        if assessment_missing_majority:
            collapse_penalty *= 0.8
        if metrics["pedagogically_collapsed"]:
            collapse_penalty *= 0.85
        score *= collapse_penalty
        metrics["score"] = round(score, 4)
        blocking_nodes = [node for node in (metrics.get("weak_nodes") or []) if not is_cosmetic_weak_node(node)]
        cosmetic_nodes = [node for node in (metrics.get("weak_nodes") or []) if is_cosmetic_weak_node(node)]
        metrics["blocking_weak_nodes"] = blocking_nodes
        metrics["cosmetic_weak_nodes"] = cosmetic_nodes
        metrics["blocking_weak_nodes_count"] = len(blocking_nodes)
        metrics["cosmetic_weak_nodes_count"] = len(cosmetic_nodes)
        metrics["only_cosmetic_weakness"] = bool(cosmetic_nodes) and not blocking_nodes

        weak_nodes_threshold = max(3, len(sections_blueprint) + 1)
        checks = {
            "stage_total": metrics["stage_total"] >= minimum_expected_stages,
            "total_items": metrics["total_items"] >= max(14, metrics["stage_total"] * 4),
            "stage_coverage": stage_coverage >= 0.72,
            "source_coverage": source_coverage >= 0.55,
            "content_density": content_density_score >= 0.58,
            "stage_completeness": stage_completeness_score >= 0.55,
            "anti_generic": anti_generic_score >= 0.62,
            "teacher_activity": teacher_activity_score >= 0.62,
            "student_activity": student_activity_score >= 0.62,
            "assessment_coverage": assessment_coverage_score >= 0.4,
            "pedagogically_collapsed": not metrics["pedagogically_collapsed"],
            "weak_nodes": len(blocking_nodes) <= weak_nodes_threshold,
        }
        metrics["acceptance_fail_reasons"] = [name for name, ok in checks.items() if not ok]
        metrics["acceptable"] = all(checks.values())
        return metrics

    @staticmethod
    def _summarize_quality_report(quality_report):
        weak_nodes = quality_report.get("weak_nodes") or []
        unique_titles = []
        seen = set()
        for node in weak_nodes:
            title = clean_text_preserve_prefix(node.get("title") or "")
            if not title or title in seen:
                continue
            seen.add(title)
            unique_titles.append(title)

        lines = [
            f"Загальна кількість змістових фрагментів: {quality_report.get('total_items', 0)}.",
            f"Заповненість етапів: {quality_report.get('stage_filled', 0)}/{quality_report.get('stage_total', 0)}.",
            f"Покриття етапів source-матеріалом: {quality_report.get('stage_source_covered', 0)}/{quality_report.get('stage_with_source', 0)}.",
            f"Щільність змісту: {quality_report.get('content_density_score', 0.0)}.",
            f"Повнота етапів: {quality_report.get('stage_completeness_score', 0.0)}.",
            f"Anti-generic score: {quality_report.get('anti_generic_score', 0.0)}.",
            f"Активність учителя: {quality_report.get('teacher_activity_score', 0.0)}.",
            f"Активність учнів: {quality_report.get('student_activity_score', 0.0)}.",
            f"Оцінювання/рефлексія: {quality_report.get('assessment_coverage_score', 0.0)}.",
        ]
        if unique_titles:
            lines.append("Найслабші вузли, які треба посилити:")
            lines.extend(f"- {title}" for title in unique_titles[:12])
        return "\n".join(lines)

    def _extract_repair_targets(self, quality_report, limit=12):
        weak_nodes = quality_report.get("weak_nodes") if isinstance(quality_report, dict) else []
        targets = []
        for node in weak_nodes or []:
            if not isinstance(node, dict):
                continue
            section_index = node.get("section_index")
            trace = node.get("trace") if isinstance(node.get("trace"), list) else []
            if section_index is None:
                continue
            targets.append(
                {
                    "kind": clean_text(node.get("kind") or ""),
                    "title": clean_text_preserve_prefix(node.get("title") or ""),
                    "path": clean_text_preserve_prefix(node.get("path") or ""),
                    "section_index": int(section_index),
                    "trace": trace,
                    "issues": [clean_text_preserve_prefix(item) for item in (node.get("issues") or []) if clean_text_preserve_prefix(item)],
                }
            )
            if len(targets) >= limit:
                break
        return targets

    @staticmethod
    def _resolve_payload_node(payload, section_index, trace):
        if not isinstance(payload, dict):
            return None
        sections = payload.get("sections")
        if not isinstance(sections, list) or section_index < 0 or section_index >= len(sections):
            return None
        node = sections[section_index]
        for step in trace or []:
            if not isinstance(node, dict):
                return None
            key = step.get("key")
            index = step.get("index")
            if not isinstance(index, int):
                return None
            bucket = node.get(key)
            if not isinstance(bucket, list) or index < 0 or index >= len(bucket):
                return None
            node = bucket[index]
        return node

    def _build_repair_target_snapshot(self, payload, repair_targets):
        if not isinstance(payload, dict):
            return {"header_values": {}, "targets": []}
        header_values = payload.get("header_values") if isinstance(payload.get("header_values"), dict) else {}
        snapshot = {"header_values": header_values, "targets": []}
        for target in repair_targets or []:
            section_index = target.get("section_index")
            trace = target.get("trace") or []
            if not isinstance(section_index, int):
                continue
            node = self._resolve_payload_node(payload, section_index, trace)
            if not isinstance(node, dict):
                continue
            snapshot["targets"].append(
                {
                    "section_index": section_index,
                    "trace": trace,
                    "title": clean_text(target.get("title") or ""),
                    "reason": clean_text(target.get("reason") or ""),
                    "items": normalize_block_list(node.get("items") or []),
                    "substeps_count": len(node.get("substeps") or []),
                    "stages_count": len(node.get("stages") or []),
                }
            )
        return snapshot

    def _merge_payload_by_repair_targets(self, previous_payload, candidate_payload, repair_targets):
        if not isinstance(previous_payload, dict):
            return candidate_payload if isinstance(candidate_payload, dict) else {}
        if not isinstance(candidate_payload, dict):
            return previous_payload

        merged = copy.deepcopy(previous_payload)
        if isinstance(candidate_payload.get("header_values"), dict):
            merged["header_values"] = candidate_payload.get("header_values")

        unique_targets = []
        seen = set()
        for target in repair_targets or []:
            key = (target.get("section_index"), json.dumps(target.get("trace") or [], ensure_ascii=False, sort_keys=True))
            if key in seen:
                continue
            seen.add(key)
            unique_targets.append(target)

        unique_targets.sort(key=lambda item: len(item.get("trace") or []))
        if not unique_targets:
            return candidate_payload

        for target in unique_targets:
            section_index = int(target.get("section_index"))
            trace = target.get("trace") or []
            source_node = self._resolve_payload_node(candidate_payload, section_index, trace)
            destination_node = self._resolve_payload_node(merged, section_index, trace)
            if not isinstance(source_node, dict) or not isinstance(destination_node, dict):
                continue
            for key in ("items", "substeps", "stages"):
                if key in source_node:
                    destination_node[key] = copy.deepcopy(source_node.get(key) or [])

        return merged

    async def _repair_blueprint_content(
        self,
        topic,
        grade,
        requirements,
        subject,
        context,
        source_context,
        blueprint,
        reference_context,
        slide_reference_plan,
        previous_payload,
        quality_report,
        generation_plan=None,
        runtime_trace=None,
    ):
        prompt_blueprint = {
            "header_fields": [
                {
                    "label": field.get("label"),
                    "sample_value": self._truncate_example(field.get("sample_value") or "", 220),
                }
                for field in blueprint.get("header_fields", [])
            ],
            "sections": [self._build_prompt_node(section) for section in blueprint.get("sections", [])],
        }
        normalized_plan = self._coerce_generation_plan(generation_plan or {}, blueprint, subject=subject, grade=grade)
        repair_targets = self._extract_repair_targets(quality_report, limit=14)
        compact_plan_outline = self._build_generation_plan_prompt_outline(normalized_plan, stage_limit=14)
        prompt_blueprint_json = json.dumps(prompt_blueprint, ensure_ascii=False, indent=2)
        compact_plan_json = json.dumps(compact_plan_outline, ensure_ascii=False, indent=2)
        style_cues_json = json.dumps(self._extract_style_cues(blueprint, limit=8), ensure_ascii=False, indent=2)
        repair_targets_json = json.dumps(repair_targets, ensure_ascii=False, indent=2)
        repair_snapshot = self._build_repair_target_snapshot(previous_payload or {}, repair_targets)
        repair_snapshot_json = json.dumps(repair_snapshot, ensure_ascii=False, indent=2)
        prompt_reference_context = self._compress_context_for_prompt(reference_context, max_chars=1800, max_lines=40)
        prompt_source_context = self._compress_context_for_prompt(source_context, max_chars=2600, max_lines=65)

        sys_instr = (
            "Ти методист НУШ. Потрібно посилити лише слабкі вузли JSON-конспекту без зміни структури blueprint. "
            "Переписуй тільки цільові nodes з REPAIR_TARGETS: додавай конкретні дії вчителя і учнів, вправи, питання, оцінювання. "
            "Зберігай порядок sections/stages/substeps. Поверни тільки валідний JSON українською."
        )
        quality_notes = self._summarize_quality_report(quality_report)

        prompt = f"""ПОПЕРЕДНІЙ JSON потребує repair-pass.

ЩО САМЕ ПОСИЛИТИ:
{quality_notes}

REPAIR_TARGETS (переписуй саме ці вузли):
{repair_targets_json}

ОБОВ'ЯЗКОВІ ПРАВИЛА REPAIR:
- не додавай і не видаляй вузли;
- не змінюй порядок вузлів;
- оновлюй тільки header_values та items;
- для кожного цільового stage забезпеч 4-8 змістовних рядків;
- додай конкретні дії вчителя, дії учнів, 1-2 конкретні вправи/питання/приклади;
- додай формувальне оцінювання або рефлексію там, де доречно;
- переписуй stage як живий фрагмент уроку, не як сухий звіт по ролях;
- прибери загальні формулювання і повтори;
- без markdown, лише JSON.

ПРЕДМЕТ: {subject}
КЛАС: {grade}
ТЕМА: {topic}
ВИМОГИ: {requirements}
ДОДАТКОВИЙ КОНТЕКСТ: {context}

{prompt_reference_context}

{prompt_source_context}

ПРИВ'ЯЗКА ЕТАПІВ ДО СЛАЙДІВ:
{slide_reference_plan or "Окрема прив'язка етапів до слайдів не задана."}

STYLE_CUES_FROM_TEMPLATES:
{style_cues_json}

PEDAGOGICAL_PLAN:
{compact_plan_json}

ПОПЕРЕДНІЙ JSON:
{repair_snapshot_json}

BLUEPRINT:
{prompt_blueprint_json}
"""

        try:
            repair_timeout_raw = os.getenv("GENERATOR_REPAIR_TIMEOUT_SEC", "25")
            try:
                repair_timeout_sec = float(repair_timeout_raw)
            except ValueError:
                repair_timeout_sec = 25.0
            response = await self._execute_model_call(
                call_name="strict_repair_payload",
                prompt=prompt,
                system_instruction=sys_instr,
                temperature=0.2,
                runtime_trace=runtime_trace,
                timeout_sec=repair_timeout_sec,
                extra={
                    "subject": clean_text(subject),
                    "grade": clean_text(grade),
                    "topic_chars": len(clean_text(topic)),
                    "source_context_chars": len(clean_text(source_context)),
                    "reference_context_chars": len(clean_text(reference_context)),
                    "compressed_source_context_chars": len(clean_text(prompt_source_context)),
                    "compressed_reference_context_chars": len(clean_text(prompt_reference_context)),
                    "plan_outline_chars": len(compact_plan_json),
                    "style_cues_chars": len(style_cues_json),
                    "blueprint_json_chars": len(prompt_blueprint_json),
                    "repair_targets_chars": len(repair_targets_json),
                    "repair_snapshot_chars": len(repair_snapshot_json),
                    "repair_targets": len(repair_targets),
                    "weak_nodes": len((quality_report or {}).get("weak_nodes") or []),
                },
            )
            if not response:
                raise RuntimeError("model_call_failed")
            candidate = self._extract_json_payload(getattr(response, "text", ""))
            if not self._payload_has_content(candidate):
                return previous_payload if isinstance(previous_payload, dict) else {}
            return self._merge_payload_by_repair_targets(previous_payload, candidate, repair_targets)
        except Exception as exc:
            logger.error("Strict example repair generation failed: %s", exc)
            return previous_payload if isinstance(previous_payload, dict) else {}

    @staticmethod
    def _payload_has_content(payload):
        if not isinstance(payload, dict):
            return False

        header_values = payload.get("header_values")
        if isinstance(header_values, dict):
            for value in header_values.values():
                if clean_text(value):
                    return True

        def node_has_items(node):
            if not isinstance(node, dict):
                return False
            if normalize_block_list(node.get("items") or []):
                return True
            for child in node.get("substeps") or []:
                if node_has_items(child):
                    return True
            for child in node.get("stages") or []:
                if node_has_items(child):
                    return True
            return False

        for section in payload.get("sections") or []:
            if node_has_items(section):
                return True
        return False

    def _build_local_blueprint_fallback_document(self, blueprint, topic, grade, subject):
        blueprint = blueprint if isinstance(blueprint, dict) else {}
        topic = clean_text(topic)
        grade = clean_text(grade)
        subject = clean_text(subject)

        header_fallbacks = {
            "Тема": topic,
            "Клас": grade,
            "Предмет": subject,
        }
        anti_meta_actions = []

        def pick_source_items(node, limit):
            candidates = self._sanitize_source_items(node.get("source_items") or [], limit=12)
            filtered = []
            topic_tokens = self._topic_tokens(topic)
            for candidate in candidates:
                normalized = clean_text_preserve_prefix(candidate)
                if not normalized:
                    continue
                if self._looks_like_structural_dump(normalized):
                    continue
                if self._is_meta_report_line(normalized):
                    continue
                if self._is_topic_leak_line(normalized, topic=topic, subject=subject, grade=grade):
                    continue
                filtered.append(normalized)
            if topic_tokens:
                filtered.sort(
                    key=lambda item: (
                        0 if topic_tokens & self._topic_tokens(item) else 1,
                        len(item),
                    )
                )
            return filtered[:limit]

        def visit(node, kind):
            title = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
            expected_field = self._expected_substep_field(node) if kind == "substep" else None
            item_limit = 4 if kind in {"section", "stage"} else 2
            items = pick_source_items(node, item_limit)
            if not items and title:
                items = self._build_topic_anchor_items(
                    topic=topic,
                    node_title=title,
                    expected_field=expected_field,
                    limit=1,
                    subject=subject,
                    grade=grade,
                )
            if kind == "stage":
                stage_seed = list(items)
                stage_seed.extend(
                    self._build_topic_anchor_items(
                        topic=topic,
                        node_title=title,
                        expected_field=None,
                        limit=2,
                        subject=subject,
                        grade=grade,
                    )
                )
                stage_enriched = self._enrich_stage_draft_locally(
                    stage_seed,
                    stage_title=title,
                    subject=subject,
                    grade=grade,
                )
                allow_assessment = self._classify_stage_name(title) in {"main", "closing"}
                stage_cleaned, stage_actions = self._sanitize_stage_script_lines(
                    stage_enriched,
                    stage_title=title,
                    topic=topic,
                    subject=subject,
                    grade=grade,
                    allow_assessment=allow_assessment,
                )
                items = self._coerce_plan_items(stage_cleaned or stage_enriched, limit=12)
                if stage_actions:
                    anti_meta_actions.append(
                        {
                            "stage": self._safe_display_text(title) or "Етап уроку",
                            "actions": stage_actions[:16],
                        }
                    )

            substeps = []
            substep_index_map = {}
            for index, substep in enumerate(node.get("substeps") or []):
                built = visit(substep, "substep")
                has_title = clean_text_preserve_prefix(built.get("display_title") or built.get("title") or "")
                has_payload = bool((built.get("items") or []) or (built.get("substeps") or []) or (built.get("stages") or []))
                if not has_title and not has_payload:
                    continue
                substep_index_map[index] = len(substeps)
                substeps.append(built)

            stages = []
            stage_index_map = {}
            for index, stage in enumerate(node.get("stages") or []):
                built = visit(stage, "stage")
                has_title = clean_text_preserve_prefix(built.get("display_title") or built.get("title") or "")
                has_payload = bool((built.get("items") or []) or (built.get("substeps") or []) or (built.get("stages") or []))
                if not has_title and not has_payload:
                    continue
                stage_index_map[index] = len(stages)
                stages.append(built)

            children_order = []
            for child_kind, child_index in (node.get("children_order") or []):
                if child_kind == "substep":
                    mapped = substep_index_map.get(child_index)
                    if mapped is not None:
                        children_order.append(("substep", mapped))
                elif child_kind == "stage":
                    mapped = stage_index_map.get(child_index)
                    if mapped is not None:
                        children_order.append(("stage", mapped))

            if not children_order:
                children_order = [("substep", idx) for idx in range(len(substeps))]
                children_order.extend(("stage", idx) for idx in range(len(stages)))

            built_node = {
                "title": node.get("title") or "",
                "display_title": self._safe_display_text(node.get("display_title") or node.get("title") or ""),
                "style": node.get("style") or "Normal",
                "content_style": node.get("content_style") or "Normal",
                "items": items,
                "sample_items": node.get("sample_items") or [],
                "source_items": node.get("source_items") or [],
                "sample_item_styles": node.get("sample_item_styles") or [],
                "substeps": substeps,
                "children_order": children_order,
            }
            if "stages" in node:
                built_node["stages"] = stages
            return built_node

        result = {
            "mode": "strict_example_local_fallback",
            "reference_file": blueprint.get("reference_file") or "",
            "topic": topic,
            "grade": grade,
            "subject": subject,
            "header_fields": [],
            "sections": [],
            "lesson_flow": [],
            "anti_meta_actions": [],
        }
        for field in blueprint.get("header_fields") or []:
            label = clean_text(field.get("label") or "")
            value = header_fallbacks.get(label) or clean_text(field.get("value") or "")
            result["header_fields"].append(
                {
                    "label": label,
                    "style": field.get("style") or "Normal",
                    "value": value,
                }
            )

        for section in blueprint.get("sections") or []:
            built_section = visit(section, "section")
            result["sections"].append(built_section)
            for stage in built_section.get("stages") or []:
                result["lesson_flow"].append(self._stage_node_to_flow_entry(stage))
        result["anti_meta_actions"] = anti_meta_actions[:80]
        return result

    def _apply_generated_content_to_blueprint(self, blueprint, payload, topic, grade, subject):
        payload = payload if isinstance(payload, dict) else {}
        header_values = payload.get('header_values') if isinstance(payload.get('header_values'), dict) else {}
        header_fallbacks = {
            'Тема': clean_text(topic),
            'Клас': clean_text(grade),
            'Предмет': clean_text(subject),
            'Тема': clean_text(topic),
            'Клас': clean_text(grade),
            'Предмет': clean_text(subject),
        }

        result = {
            'mode': 'strict_example',
            'reference_file': blueprint.get('reference_file') or '',
            'topic': clean_text(topic),
            'grade': clean_text(grade),
            'subject': clean_text(subject),
            'header_fields': [],
            'sections': [],
            'lesson_flow': [],
        }

        for field in blueprint.get('header_fields', []):
            label = field.get('label') or ''
            value = clean_text(header_values.get(label) or '')
            if not value:
                value = header_fallbacks.get(label) or clean_text(field.get('value') or '')
            if label in {'Обладнання', 'Обладнання'}:
                value = self._join_items(self._prune_equipment_items(value), separator="; ")
            result['header_fields'].append({
                'label': self._safe_display_text(label),
                'style': field.get('style') or 'Normal',
                'value': value,
            })
            if label in {'Тема', 'Тема'} and value:
                result['topic'] = value
            elif label in {'Клас', 'Клас'} and value:
                result['grade'] = value
            elif label in {'Предмет', 'Предмет'} and value:
                result['subject'] = value

        section_payloads = payload.get('sections') if isinstance(payload.get('sections'), list) else []

        def apply_node(node, node_payload, kind="section"):
            node_payload = node_payload if isinstance(node_payload, dict) else {}
            expected_field = self._expected_substep_field(node) if kind == "substep" else None
            generated_items = self._coerce_items(
                node_payload.get('items'),
                list(node.get('source_items') or []),
            )
            max_items_by_kind = {
                "section": 6,
                "stage": 10,
                "substep": 8,
            }
            generated_items = self._unique_clean_items(
                generated_items,
                expected_field=expected_field,
                limit=max_items_by_kind.get(kind, 6),
            )
            if not generated_items and not (node.get("stages") or node.get("substeps")):
                generated_items = self._build_topic_anchor_items(
                    topic=topic,
                    node_title=node.get("display_title") or node.get("title") or "",
                    expected_field=expected_field,
                    limit=1,
                    subject=subject,
                    grade=grade,
                )
            applied = {
                'title': self._safe_display_text(node.get('title') or '') or 'Етап уроку',
                'display_title': self._safe_display_text(node.get('display_title') or node.get('title') or '') or self._safe_display_text(node.get('title') or '') or 'Етап уроку',
                'style': node.get('style') or 'Normal',
                'content_style': node.get('content_style') or 'Normal',
                'items': generated_items,
                'sample_items': node.get('sample_items') or [],
                'source_items': node.get('source_items') or [],
                'sample_item_styles': node.get('sample_item_styles') or [],
                'substeps': [],
                'children_order': list(node.get('children_order') or []),
            }
            payload_substeps = node_payload.get('substeps') if isinstance(node_payload.get('substeps'), list) else []
            for index, substep in enumerate(node.get('substeps') or []):
                applied['substeps'].append(
                    apply_node(
                        substep,
                        payload_substeps[index] if index < len(payload_substeps) else {},
                        kind="substep",
                    )
                )
            if kind == "stage" and applied.get("substeps"):
                inferred_fields = self._infer_stage_fields_from_items(applied.get("items") or [])
                consumed_items = set()
                for substep_index, substep in enumerate(applied.get("substeps") or []):
                    blueprint_substep = (node.get("substeps") or [])[substep_index] if substep_index < len(node.get("substeps") or []) else {}
                    expected_field = self._expected_substep_field(blueprint_substep)
                    if not expected_field or normalize_block_list(substep.get("items") or []):
                        continue
                    mapped_items = self._unique_clean_items(
                        inferred_fields.get(expected_field) or [],
                        expected_field=expected_field if expected_field in {"teacher_actions", "student_actions", "activities"} else None,
                        limit=8,
                    )
                    if not mapped_items:
                        continue
                    substep["items"] = mapped_items
                    consumed_items.update(mapped_items)
                if consumed_items:
                    applied["items"] = [item for item in (applied.get("items") or []) if item not in consumed_items]
            if kind == "stage":
                stage_items = self._coerce_plan_items(applied.get("items") or [], limit=12)
                generic_ratio = (
                    sum(1 for item in stage_items if self._is_generic_lesson_item(item)) / len(stage_items)
                    if stage_items else 1.0
                )
                if len(stage_items) < 4 or generic_ratio > 0.55:
                    topic_tokens = self._topic_tokens(topic)
                    source_seed = []
                    for source_line in [
                        *normalize_block_list(node.get("source_items") or []),
                        *normalize_block_list(node.get("sample_items") or []),
                    ]:
                        cleaned_source = self._safe_display_text(source_line)
                        if not cleaned_source:
                            continue
                        if len(cleaned_source.split()) > 24:
                            continue
                        if topic_tokens and not (topic_tokens & self._topic_tokens(cleaned_source)):
                            if self._has_task_signal(cleaned_source) or re.search(r"\d", cleaned_source):
                                continue
                        source_seed.append(cleaned_source)
                    seed_lines = [
                        *stage_items,
                        *source_seed,
                    ]
                    enriched = self._enrich_stage_draft_locally(
                        seed_lines,
                        stage_title=applied.get("display_title") or applied.get("title") or "",
                        subject=subject,
                        grade=grade,
                    )
                    merged_stage_items = self._coerce_plan_items([*stage_items, *enriched], limit=12)
                    applied["items"] = merged_stage_items
                stage_group = self._classify_stage_name(applied.get("display_title") or applied.get("title") or "")
                allow_assessment = stage_group in {"main", "closing"}
                filtered_stage_items, anti_meta_actions = self._sanitize_stage_script_lines(
                    applied.get("items") or [],
                    stage_title=applied.get("display_title") or applied.get("title") or "",
                    topic=topic,
                    subject=subject,
                    grade=grade,
                    allow_assessment=allow_assessment,
                )
                if filtered_stage_items:
                    applied["items"] = self._coerce_plan_items(filtered_stage_items, limit=12)
                if anti_meta_actions:
                    result.setdefault("anti_meta_actions", []).append(
                        {
                            "stage": applied.get("display_title") or applied.get("title") or "",
                            "actions": anti_meta_actions[:20],
                        }
                    )
                inferred_after_enrich = self._infer_stage_fields_from_items(applied.get("items") or [])
                for substep_index, substep in enumerate(applied.get("substeps") or []):
                    blueprint_substep = (node.get("substeps") or [])[substep_index] if substep_index < len(node.get("substeps") or []) else {}
                    expected_field = self._expected_substep_field(blueprint_substep)
                    current_items = normalize_block_list(substep.get("items") or [])
                    if current_items or not expected_field:
                        continue
                    fallback_items = inferred_after_enrich.get(expected_field) or []
                    if expected_field == "reflection":
                        fallback_items = [*(inferred_after_enrich.get("reflection") or []), *(inferred_after_enrich.get("assessment") or [])]
                    substep["items"] = self._unique_clean_items(
                        fallback_items,
                        expected_field=expected_field if expected_field in {"teacher_actions", "student_actions", "activities"} else None,
                        limit=8,
                    )
                    filtered_sub_items, sub_actions = self._sanitize_stage_script_lines(
                        substep.get("items") or [],
                        stage_title=applied.get("display_title") or applied.get("title") or "",
                        topic=topic,
                        subject=subject,
                        grade=grade,
                        allow_assessment=allow_assessment,
                    )
                    if filtered_sub_items:
                        substep["items"] = self._coerce_plan_items(filtered_sub_items, limit=8)
                    if sub_actions:
                        result.setdefault("anti_meta_actions", []).append(
                            {
                                "stage": applied.get("display_title") or applied.get("title") or "",
                                "substep": substep.get("display_title") or substep.get("title") or "",
                                "actions": sub_actions[:12],
                            }
                        )
            if 'stages' in node:
                payload_stages = node_payload.get('stages') if isinstance(node_payload.get('stages'), list) else []
                applied['stages'] = []
                for index, stage in enumerate(node.get('stages') or []):
                    stage_applied = apply_node(
                        stage,
                        payload_stages[index] if index < len(payload_stages) else {},
                        kind="stage",
                    )
                    applied['stages'].append(stage_applied)
                    result['lesson_flow'].append(
                        self._stage_node_to_flow_entry(
                            stage_applied,
                            topic=topic,
                            subject=subject,
                            grade=grade,
                        )
                    )
            return applied

        for index, section in enumerate(blueprint.get('sections', [])):
            result['sections'].append(
                apply_node(
                    section,
                    section_payloads[index] if index < len(section_payloads) else {},
                    kind="section",
                )
            )

        return result

    def _compact_blueprint_document_model(self, document_model):
        if not isinstance(document_model, dict):
            return {}

        def compact_node(node):
            if not isinstance(node, dict):
                return {}

            compacted_substeps = []
            substep_index_map = {}
            for index, substep in enumerate(node.get("substeps") or []):
                compacted = compact_node(substep)
                has_title = clean_text_preserve_prefix(
                    compacted.get("display_title") or compacted.get("title") or ""
                )
                has_payload = bool(
                    (compacted.get("items") or [])
                    or (compacted.get("substeps") or [])
                    or (compacted.get("stages") or [])
                )
                if not has_title and not has_payload:
                    continue
                substep_index_map[index] = len(compacted_substeps)
                compacted_substeps.append(compacted)

            compacted_stages = []
            stage_index_map = {}
            for index, stage in enumerate(node.get("stages") or []):
                compacted = compact_node(stage)
                has_title = clean_text_preserve_prefix(
                    compacted.get("display_title") or compacted.get("title") or ""
                )
                has_payload = bool(
                    (compacted.get("items") or [])
                    or (compacted.get("substeps") or [])
                    or (compacted.get("stages") or [])
                )
                if not has_title and not has_payload:
                    continue
                stage_index_map[index] = len(compacted_stages)
                compacted_stages.append(compacted)

            children_order = []
            for child_kind, child_index in (node.get("children_order") or []):
                if child_kind == "substep":
                    mapped = substep_index_map.get(child_index)
                    if mapped is not None:
                        children_order.append(("substep", mapped))
                elif child_kind == "stage":
                    mapped = stage_index_map.get(child_index)
                    if mapped is not None:
                        children_order.append(("stage", mapped))

            if not children_order:
                children_order = [("substep", index) for index in range(len(compacted_substeps))]
                children_order.extend(("stage", index) for index in range(len(compacted_stages)))

            compacted_items = []
            local_seen_item_keys = set()
            for item in normalize_block_list(node.get("items") or []):
                normalized_item = clean_text_preserve_prefix(item)
                key = self._normalize_template_name(normalized_item)
                if not normalized_item or not key:
                    continue
                if key in local_seen_item_keys:
                    continue
                local_seen_item_keys.add(key)
                compacted_items.append(normalized_item)

            compacted_node = {
                **node,
                "items": compacted_items,
                "substeps": compacted_substeps,
                "children_order": children_order,
            }
            if "stages" in node:
                compacted_node["stages"] = compacted_stages
            return compacted_node

        compacted_sections = [compact_node(section) for section in (document_model.get("sections") or [])]
        compacted_lesson_flow = []
        doc_topic = clean_text(document_model.get("topic") or "")
        doc_subject = clean_text(document_model.get("subject") or "")
        doc_grade = clean_text(document_model.get("grade") or "")
        for section in compacted_sections:
            for stage in section.get("stages") or []:
                compacted_lesson_flow.append(
                    self._stage_node_to_flow_entry(
                        stage,
                        topic=doc_topic,
                        subject=doc_subject,
                        grade=doc_grade,
                    )
                )
        if self._is_math_subject(doc_subject) and self._grade_to_int(doc_grade) == 3:
            compacted_lesson_flow = self._postprocess_math3_lesson_flow(compacted_lesson_flow)

        return {
            **document_model,
            "sections": compacted_sections,
            "lesson_flow": compacted_lesson_flow,
        }

    @staticmethod
    def _normalize_template_name(value):
        text = clean_text(value).lower().replace(chr(8217), "'")
        text = re.sub(r"[^\w']+", " ", text, flags=re.UNICODE)
        return re.sub(r"\s+", " ", text).strip()

    @staticmethod
    def _build_template_blueprint(template_structure):
        if not template_structure:
            return None

        def unique(values):
            seen = set()
            result = []
            for value in values:
                if not value or value in seen:
                    continue
                seen.add(value)
                result.append(value)
            return result

        return {
            "header_order": unique(item.get("label") for item in template_structure.get("header_fields", [])),
            "section_order": unique(item.get("title") for item in template_structure.get("sections", [])),
            "stage_order": unique(item.get("title") for item in template_structure.get("stages", [])),
        }

    def _get_allowed_optional_sections(self, template_blueprint):
        if not template_blueprint:
            return set()

        normalized_allowed = {
            self._normalize_template_name(title)
            for title in (template_blueprint.get("section_order") or [])
            if clean_text(title)
        }
        return {
            field_name
            for title, field_name in self.OPTIONAL_SECTION_FIELD_MAP.items()
            if self._normalize_template_name(title) in normalized_allowed
        }

    def _prune_rich_lesson_data(self, lesson_data, template_blueprint=None):
        if not isinstance(lesson_data, dict):
            return lesson_data

        allowed_optional_sections = self._get_allowed_optional_sections(template_blueprint)
        for field_name in set(self.OPTIONAL_SECTION_FIELD_MAP.values()):
            if field_name in allowed_optional_sections:
                continue
            if field_name == "expected_results":
                lesson_data[field_name] = {"knowledge": [], "skills": [], "values": []}
            else:
                lesson_data[field_name] = []

        return lesson_data

    def _build_stage_substep_allowlist(self, template_structure):
        allowlist = {}
        if not isinstance(template_structure, dict):
            return allowlist

        for section in template_structure.get("sections", []) or []:
            for stage in section.get("stages", []) or []:
                stage_title = clean_text(stage.get("title") or "")
                if not stage_title:
                    continue
                stage_key = self._normalize_template_name(stage_title)
                labels = allowlist.setdefault(stage_key, set())
                for substep in stage.get("substeps", []) or []:
                    label = clean_text(substep.get("title") or "")
                    if label:
                        labels.add(label)

        return allowlist

    def _prune_stage_materials(self, materials, equipment=None, resources=None):
        cleaned_materials = normalize_list(materials)
        if not cleaned_materials:
            return []

        global_items = {
            self._normalize_template_name(item)
            for item in [*(equipment or []), *(resources or [])]
            if clean_text(item)
        }
        generic_stage_items = {
            self._normalize_template_name(item)
            for item in (
                "підручник",
                "зошит",
                "зошити",
                "дошка",
                "крейда",
                "презентація",
                "картки",
                "робочий аркуш",
                "робочі аркуші",
            )
        }

        filtered = []
        for item in cleaned_materials:
            normalized_item = self._normalize_template_name(item)
            if normalized_item in global_items and normalized_item in generic_stage_items:
                continue
            filtered.append(item)

        return filtered

    def _prune_equipment_items(self, equipment):
        cleaned_items = normalize_list(equipment)
        if not cleaned_items:
            return []

        trivial_names = {
            self._normalize_template_name(item)
            for item in self.TRIVIAL_EQUIPMENT_NAMES
        }
        filtered = [
            item for item in cleaned_items
            if self._normalize_template_name(item) not in trivial_names
        ]
        return filtered

    def _classify_stage_name(self, stage_name):
        normalized_stage = self._normalize_template_name(stage_name)
        if any(
            token in normalized_stage
            for token in (
                "вступ",
                "організац",
                "мотивац",
                "повідомлення теми",
                "початок уроку",
                "start",
                "intro",
                "вступ",
                "організац",
                "мотивац",
                "повідомлення теми",
            )
        ):
            return "intro"
        if any(
            token in normalized_stage
            for token in (
                "підсумок",
                "рефлекс",
                "домашнє",
                "завершення",
                "closing",
                "summary",
                "підсумок",
                "рефлекс",
                "домашнє",
            )
        ):
            return "closing"
        return "main"

    def _extract_stage_title(self, stage):
        if not isinstance(stage, dict):
            return ""
        return clean_text(
            stage.get("stage")
            or stage.get("title")
            or stage.get("display_title")
            or ""
        )

    def _stage_group_order(self, stage_title):
        group = self._classify_stage_name(stage_title)
        return {"intro": 0, "main": 1, "closing": 2}.get(group, 1)

    def _stage_suborder(self, stage_title):
        normalized = self._normalize_template_name(stage_title)
        if "організац" in normalized or "організац" in normalized:
            return 0
        if "рефлекс" in normalized or "рефлекс" in normalized:
            return 0
        if "підсумок" in normalized or "РїС–РґСЃСѓРјРѕРє" in normalized:
            return 1
        if "домашн" in normalized or "домашн" in normalized:
            return 2
        return 0

    def _normalize_stage_sequence(self, stages):
        prepared = []
        for index, stage in enumerate(stages or []):
            if not isinstance(stage, dict):
                continue
            prepared.append((index, dict(stage)))

        def sort_key(item):
            index, stage = item
            title = self._extract_stage_title(stage)
            return (
                self._stage_group_order(title),
                self._stage_suborder(title),
                index,
            )

        return [stage for _, stage in sorted(prepared, key=sort_key)]

    def _split_explicit_stage_label(self, value):
        text = clean_text_preserve_prefix(value)
        match = re.match(r"^([^:]{2,40}):\s*(.+)$", text)
        if not match:
            return None, text
        field_name = self._resolve_stage_field_from_label(match.group(1))
        if not field_name:
            return None, text
        return field_name, normalize_sentence_punctuation(match.group(2))

    def _looks_like_goal_line(self, value):
        normalized = self._normalize_for_matching(value)
        canonical_prefixes = (
            "ознайомити",
            "формувати",
            "закріпити",
            "розвивати",
            "виховувати",
            "перевірити",
            "узагальнити",
            "повторити",
            "створити",
            "навчити",
            "вчити",
            "відпрацювати",
            "налаштувати",
            "пояснити",
        )
        if any(normalized.startswith(prefix) for prefix in canonical_prefixes):
            return True
        for prefix in self.GOAL_LINE_PREFIXES:
            candidate = self._normalize_for_matching(prefix)
            if candidate and normalized.startswith(candidate):
                return True
        return False

    def _looks_like_student_line(self, value):
        normalized = self._normalize_for_matching(value)
        explicit_field, _ = self._split_explicit_stage_label(value)
        if explicit_field == "student_actions":
            return True
        canonical_prefixes = (
            "учні",
            "учень",
            "діти",
            "виконують",
            "відповідають",
            "слухають",
            "розглядають",
            "читають",
            "записують",
            "працюють",
            "називають",
            "будують",
            "висловлюють",
            "малюють",
            "обговорюють",
            "рухаються",
            "спостерігають",
            "складають",
            "повторюють",
            "розв'язують",
            "розв'язують",
            "запам'ятовують",
            "запам'ятовують",
            "викладають",
            "порівнюють",
            "характеризують",
            "вивчають",
            "визначають",
            "пишуть",
        )
        if any(normalized.startswith(prefix) for prefix in canonical_prefixes):
            return True
        for prefix in self.STUDENT_ACTION_PREFIXES:
            candidate = self._normalize_for_matching(prefix)
            if candidate and normalized.startswith(candidate):
                return True
        return False

    def _looks_like_teacher_line(self, value):
        normalized = self._normalize_for_matching(value)
        explicit_field, _ = self._split_explicit_stage_label(value)
        if explicit_field == "teacher_actions":
            return True
        stripped = clean_text_preserve_prefix(value)
        if "?" in stripped or stripped.startswith(("Р'В«", "\"")):
            return True
        canonical_prefixes = (
            "вчитель",
            "учитель",
            "подивіться",
            "послухайте",
            "пригадайте",
            "прочитайте",
            "відкрийте",
            "запишіть",
            "скажіть",
            "назвіть",
            "поміркуйте",
            "розгляньте",
            "погляньте",
            "спробуйте",
            "давайте",
            "зараз",
            "поясніть",
            "порівняйте",
            "знайдіть",
            "виконайте",
        )
        if any(normalized.startswith(prefix) for prefix in canonical_prefixes):
            return True
        for prefix in self.TEACHER_ACTION_PREFIXES:
            candidate = self._normalize_for_matching(prefix)
            if candidate and normalized.startswith(candidate):
                return True
        return False

    def _looks_like_material_line(self, value):
        normalized = self._normalize_for_matching(value)
        if normalized.startswith(("робота з", "гра ", "вправа", "бесіда", "читання", "слухання")):
            return False
        if any(token in normalized for token in ("підручник", "картк", "презентац", "обладнання", "матеріал", "ресурс")):
            return True
        return any(self._normalize_for_matching(token) in normalized for token in self.MATERIAL_HINT_TOKENS)

    def _looks_like_assessment_line(self, value):
        normalized = self._normalize_for_matching(value)
        explicit_field, _ = self._split_explicit_stage_label(value)
        if explicit_field == "assessment":
            return True
        assessment_tokens = (
            "оцін",
            "перевір",
            "самоперевір",
            "взаємоперевір",
            "критер",
            "правильність",
            "бал",
            "rubric",
            "чек лист",
            "формувальн",
        )
        return any(token in normalized for token in assessment_tokens)

    def _looks_like_reflection_line(self, value):
        normalized = self._normalize_for_matching(value)
        explicit_field, _ = self._split_explicit_stage_label(value)
        if explicit_field == "reflection":
            return True
        reflection_tokens = (
            "рефлекс",
            "підсум",
            "що вдалося",
            "що було складно",
            "настрій",
            "самооцін",
            "виснов",
            "дерево успіху",
            "мікрофон",
            "exit ticket",
        )
        return any(token in normalized for token in reflection_tokens)

    def _infer_stage_fields_from_items(self, items, equipment=None, resources=None):
        grouped = {
            "goal": [],
            "teacher_actions": [],
            "student_actions": [],
            "activities": [],
            "assessment": [],
            "reflection": [],
            "differentiation": [],
            "materials": [],
        }
        for raw_item in normalize_block_list(items):
            field_name, explicit_value = self._split_explicit_stage_label(raw_item)
            if field_name:
                grouped[field_name].append(explicit_value)
                continue

            item = normalize_sentence_punctuation(raw_item)
            if not item:
                continue

            if self._looks_like_goal_line(item):
                grouped["goal"].append(item)
            elif self._looks_like_reflection_line(item):
                grouped["reflection"].append(item)
            elif self._looks_like_student_line(item):
                grouped["student_actions"].append(item)
            elif self._looks_like_teacher_line(item):
                grouped["teacher_actions"].append(item)
            elif self._looks_like_assessment_line(item):
                grouped["assessment"].append(item)
            elif self._looks_like_material_line(item):
                grouped["materials"].append(item)
            else:
                grouped["activities"].append(item)

        grouped["materials"] = self._prune_stage_materials(
            grouped["materials"],
            equipment=equipment,
            resources=resources,
        )
        return {
            key: list(dict.fromkeys(normalize_list(value)))
            for key, value in grouped.items()
        }

    def _normalize_lesson_flow_stage(self, stage, equipment=None, resources=None):
        if not isinstance(stage, dict):
            return stage

        normalized_stage = dict(stage)
        inferred = self._infer_stage_fields_from_items(
            normalized_stage.get("activities") or [],
            equipment=equipment,
            resources=resources,
        )

        for field_name in (
            "goal",
            "teacher_actions",
            "student_actions",
            "activities",
            "assessment",
            "reflection",
            "differentiation",
            "materials",
        ):
            current_items = normalize_list(normalized_stage.get(field_name))
            merged = current_items + [item for item in inferred.get(field_name, []) if item not in current_items]
            if field_name == "materials":
                merged = self._prune_stage_materials(merged, equipment=equipment, resources=resources)
            normalized_stage[field_name] = merged

        if normalized_stage.get("teacher_actions") or normalized_stage.get("student_actions"):
            activity_items = []
            for item in normalize_list(normalized_stage.get("activities")):
                if item in normalized_stage.get("teacher_actions", []) or item in normalized_stage.get("student_actions", []):
                    continue
                activity_items.append(item)
            normalized_stage["activities"] = activity_items

        return normalized_stage

    def _simplify_rich_stage(self, stage, equipment=None, resources=None):
        normalized_stage = self._normalize_lesson_flow_stage(
            stage,
            equipment=equipment,
            resources=resources,
        )

        normalized_stage["goal"] = []
        normalized_stage["assessment"] = []
        normalized_stage["reflection"] = []
        normalized_stage["differentiation"] = []
        normalized_stage["materials"] = []

        teacher_actions = normalize_list(normalized_stage.get("teacher_actions"))
        student_actions = normalize_list(normalized_stage.get("student_actions"))
        activities = normalize_list(normalized_stage.get("activities"))

        if not activities:
            fallback_items = []
            fallback_items.extend(teacher_actions[:1])
            fallback_items.extend(student_actions[:1])
            normalized_stage["activities"] = list(dict.fromkeys(item for item in fallback_items if clean_text(item)))
        else:
            normalized_stage["activities"] = activities

        normalized_stage["teacher_actions"] = teacher_actions
        normalized_stage["student_actions"] = student_actions
        return normalized_stage

    def _stage_node_to_flow_entry(self, stage_node, *, topic="", subject="", grade=""):
        stage_node = stage_node if isinstance(stage_node, dict) else {}
        title = self._safe_display_text(
            stage_node.get("title")
            or stage_node.get("display_title")
            or stage_node.get("stage")
            or ""
        )
        if not title:
            title = "Етап уроку"
        entry = {
            "stage": title,
            "teacher_actions": [],
            "student_actions": [],
            "activities": [],
            "goal": [],
            "assessment": [],
            "reflection": [],
            "differentiation": [],
            "materials": [],
        }

        substeps = stage_node.get("substeps") if isinstance(stage_node.get("substeps"), list) else []
        if substeps:
            for substep in substeps:
                substep = substep if isinstance(substep, dict) else {}
                label = self._normalize_template_name(
                    substep.get("title")
                    or substep.get("display_title")
                    or ""
                )
                expected_field = self._resolve_stage_field_from_label(label)
                items = self._unique_clean_items(
                    substep.get("items") or [],
                    expected_field=expected_field,
                    limit=10,
                )
                if expected_field == "teacher_actions":
                    entry["teacher_actions"].extend(items)
                elif expected_field == "student_actions":
                    entry["student_actions"].extend(items)
                elif expected_field == "activities":
                    entry["activities"].extend(items)
                elif expected_field == "assessment":
                    entry["assessment"].extend(items)
                elif expected_field == "reflection":
                    entry["reflection"].extend(items)
                elif expected_field == "goal":
                    entry["goal"].extend(items)
                elif expected_field == "materials":
                    entry["materials"].extend(items)
                elif expected_field == "differentiation":
                    entry["differentiation"].extend(items)
                else:
                    entry["activities"].extend(items)

        direct_items_raw = self._unique_clean_items(stage_node.get("items") or [], limit=16)
        direct_items, _ = self._sanitize_stage_script_lines(
            direct_items_raw,
            stage_title=title,
            topic=topic,
            subject=subject,
            grade=grade,
            allow_assessment=True,
        )
        inferred_fields = self._infer_stage_fields_from_items(direct_items)
        entry["goal"].extend(inferred_fields.get("goal") or [])
        entry["teacher_actions"].extend(inferred_fields.get("teacher_actions") or [])
        entry["student_actions"].extend(inferred_fields.get("student_actions") or [])
        entry["activities"].extend(inferred_fields.get("activities") or [])
        entry["assessment"].extend(inferred_fields.get("assessment") or [])
        entry["reflection"].extend(inferred_fields.get("reflection") or [])
        entry["differentiation"].extend(inferred_fields.get("differentiation") or [])
        entry["materials"].extend(inferred_fields.get("materials") or [])

        entry["teacher_actions"] = self._unique_clean_items(
            entry["teacher_actions"],
            expected_field="teacher_actions",
            limit=10,
        )
        entry["student_actions"] = self._unique_clean_items(
            entry["student_actions"],
            expected_field="student_actions",
            limit=10,
        )
        entry["activities"] = self._unique_clean_items(
            entry["activities"],
            expected_field="activities",
            limit=12,
        )
        entry["assessment"] = self._unique_clean_items(
            entry["assessment"],
            expected_field="assessment",
            limit=8,
        )
        entry["reflection"] = self._unique_clean_items(
            entry["reflection"],
            expected_field="reflection",
            limit=6,
        )
        entry["goal"] = self._unique_clean_items(
            entry["goal"],
            expected_field="goal",
            limit=4,
        )
        entry["differentiation"] = self._unique_clean_items(
            entry["differentiation"],
            expected_field="differentiation",
            limit=6,
        )
        entry["materials"] = self._unique_clean_items(
            entry["materials"],
            expected_field="materials",
            limit=6,
        )
        stage_group = self._classify_stage_name(title)
        allow_assessment = stage_group in {"main", "closing"}
        for field_name in ("teacher_actions", "student_actions", "activities", "assessment", "reflection"):
            filtered_lines, _ = self._sanitize_stage_script_lines(
                entry.get(field_name) or [],
                stage_title=title,
                topic=topic,
                subject=subject,
                grade=grade,
                allow_assessment=allow_assessment,
            )
            entry[field_name] = self._unique_clean_items(
                filtered_lines,
                expected_field=field_name if field_name in {"teacher_actions", "student_actions", "activities"} else None,
                limit=12 if field_name == "activities" else 10,
            )
        if not entry["student_actions"]:
            inferred_students = []
            for item in entry["activities"]:
                normalized = self._normalize_for_matching(item)
                if self._has_student_signal(item) or any(token in normalized for token in ("учні", "діти", "у пар", "у груп", "працюють", "відповідають")):
                    inferred_students.append(item)
            if self._is_math_subject(subject) and self._grade_to_int(grade) == 3:
                if stage_group == "closing":
                    inferred_students = inferred_students or [
                        "Учні: «Я впевнено обчислюю 36:6=6.»",
                        "Учні: «Ще потреную приклади 6×8 і 48:6.»",
                    ]
                elif stage_group == "intro":
                    inferred_students = inferred_students or ["Учні відповідають усно на приклади в темпі та звіряють відповідь у парі."]
                else:
                    inferred_students = inferred_students or ["Учні обчислюють приклади на 6, пояснюють один хід розв'язання в парі."]
            entry["student_actions"] = self._unique_clean_items(
                inferred_students or ["Учні в парах виконують завдання етапу та коротко пояснюють свій спосіб розв'язання."],
                expected_field="student_actions",
                limit=4,
            )
        if not entry["teacher_actions"]:
            inferred_teacher = []
            for item in entry["activities"]:
                if "?" in item or self._looks_like_teacher_line(item):
                    inferred_teacher.append(item)
            if self._is_math_subject(subject) and self._grade_to_int(grade) == 3:
                if stage_group == "closing":
                    inferred_teacher = inferred_teacher or [
                        "Учитель: «Що сьогодні вийшло найкраще: множення на 6 чи ділення на 6?»",
                        "Учитель: «Назвіть один приклад, який можете пояснити класу.»",
                    ]
                elif stage_group == "intro":
                    inferred_teacher = inferred_teacher or ["Учитель: «Швидкий старт: 6×2, 6×6, 24:6, 30:6.»"]
                else:
                    inferred_teacher = inferred_teacher or ["Учитель: «Обчисліть усно 6×5, 42:6, 6×7 і поясніть один приклад.»"]
            entry["teacher_actions"] = self._unique_clean_items(
                inferred_teacher or ["Учитель ставить уточнювальні запитання і керує темпом виконання завдання."],
                expected_field="teacher_actions",
                limit=4,
            )
        if entry["teacher_actions"] or entry["student_actions"]:
            role_items = set(entry["teacher_actions"]) | set(entry["student_actions"])
            entry["activities"] = [item for item in entry["activities"] if item not in role_items]
        if not entry["activities"]:
            if self._is_math_subject(subject) and self._grade_to_int(grade) == 3:
                if stage_group == "intro":
                    fallback_activities = ["Усний рахунок: 6×3, 6×4, 18:6, 24:6."]
                elif stage_group == "closing":
                    fallback_activities = [
                        "Покажіть смайликом результат: 🙂 — впевнено, 😐 — ще тренуюсь, 🙁 — потрібна допомога.",
                        "Закінчіть речення: «Сьогодні я навчився(лася)...»",
                    ]
                else:
                    fallback_activities = [
                        "Розв'яжіть у парі: 6×8, 54:6, (12+24):6.",
                        "Звірте одну відповідь із сусідом і поясніть, де може бути помилка.",
                    ]
                entry["activities"] = self._unique_clean_items(
                    fallback_activities,
                    expected_field="activities",
                    limit=4,
                )
            else:
                entry["activities"] = self._unique_clean_items(
                    [*(entry["teacher_actions"][:1]), *(entry["student_actions"][:1])],
                    expected_field="activities",
                    limit=3,
                )
        if not entry["assessment"] and entry["reflection"]:
            entry["assessment"] = self._unique_clean_items(entry["reflection"], expected_field="assessment", limit=3)
        return entry

    def _align_lesson_flow_to_template(self, lesson_flow, template_stage_order):
        stages = [dict(item) for item in lesson_flow or [] if isinstance(item, dict)]
        template_stage_order = [clean_text(item) for item in template_stage_order or [] if clean_text(item)]
        if not stages:
            return []
        if not template_stage_order:
            return self._normalize_stage_sequence(stages)

        aligned = []
        remaining = [dict(item) for item in stages]

        def take_match(target_title):
            normalized_target = self._normalize_template_name(target_title)
            for index, stage in enumerate(remaining):
                source_title = clean_text(stage.get("stage") or "")
                if self._normalize_template_name(source_title) == normalized_target:
                    return remaining.pop(index)

            target_group = self._classify_stage_name(target_title)
            for index, stage in enumerate(remaining):
                source_title = clean_text(stage.get("stage") or "")
                if self._classify_stage_name(source_title) == target_group:
                    return remaining.pop(index)

            return None

        for template_title in template_stage_order:
            stage = take_match(template_title)
            if not stage:
                continue
            stage["stage"] = template_title
            aligned.append(stage)

        aligned.extend(remaining)
        return self._normalize_stage_sequence(aligned)

    @staticmethod
    def _join_items(items, separator=". "):
        cleaned = normalize_list(items)
        if not cleaned:
            return ""
        return separator.join(cleaned)

    def _read_docx_text(self, file_path):
        try:
            doc = Document(file_path)
        except Exception:
            return ""
        parts = [clean_text(paragraph.text) for paragraph in doc.paragraphs]
        return "\n".join(part for part in parts if part)

    def _read_pdf_text(self, file_path):
        try:
            reader = PdfReader(str(file_path))
        except Exception:
            return ""
        parts = []
        for page in reader.pages[:25]:
            parts.append(clean_text(page.extract_text() or ""))
        return "\n".join(part for part in parts if part)

    def _read_text_file(self, file_path):
        try:
            return clean_text(Path(file_path).read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            return ""

    def _lesson_to_slides(self, lesson_data):
        def chunk(items, size):
            for i in range(0, len(items), size):
                yield items[i:i + size]

        def labeled_list(label, items):
            return [f"{label}: {item}" for item in normalize_list(items)]

        slides = []
        topic = clean_text(lesson_data.get("topic") or "Урок")
        grade = clean_text(lesson_data.get("grade") or "")
        subject = clean_text(lesson_data.get("subject") or "")
        lesson_type = clean_text(lesson_data.get("lesson_type") or "")
        subtitle_parts = []
        if grade:
            subtitle_parts.append(grade if "клас" in grade.lower() else f"{grade} клас")
        if subject:
            subtitle_parts.append(subject)
        if lesson_type:
            subtitle_parts.append(lesson_type)
        slides.append({
            "type": "title",
            "title": topic,
            "subtitle": " • ".join(subtitle_parts)
        })

        goal = normalize_list(lesson_data.get("goal"))
        if goal:
            for part in chunk(goal, 6):
                slides.append({"type": "bullets", "emoji": "🎯", "title": "Мета уроку", "bullets": part})

        tasks = normalize_list(lesson_data.get("tasks"))
        if tasks:
            for part in chunk(tasks, 6):
                slides.append({"type": "bullets", "emoji": "✅", "title": "Завдання уроку", "bullets": part})

        expected = lesson_data.get("expected_results") or {}
        expected_bullets = []
        if isinstance(expected, dict):
            expected_bullets += labeled_list("Знаннєві", expected.get("knowledge"))
            expected_bullets += labeled_list("Діяльнісні", expected.get("skills"))
            expected_bullets += labeled_list("Ціннісні", expected.get("values"))
        else:
            expected_bullets += normalize_list(expected)
        if expected_bullets:
            for part in chunk(expected_bullets, 6):
                slides.append({"type": "bullets", "emoji": "📚", "title": "Очікувані результати", "bullets": part})

        competencies = normalize_list(lesson_data.get("key_competencies"))
        if competencies:
            for part in chunk(competencies, 6):
                slides.append({"type": "bullets", "emoji": "🧩", "title": "Ключові компетентності", "bullets": part})

        cross_skills = normalize_list(lesson_data.get("cross_cutting_skills"))
        if cross_skills:
            for part in chunk(cross_skills, 6):
                slides.append({"type": "bullets", "emoji": "🧠", "title": "Наскрізні вміння", "bullets": part})

        methods = normalize_list(lesson_data.get("methods"))
        forms = normalize_list(lesson_data.get("forms"))
        methods_forms = labeled_list("Методи", methods) + labeled_list("Форми", forms)
        if methods_forms:
            for part in chunk(methods_forms, 6):
                slides.append({"type": "bullets", "emoji": "🛠️", "title": "Методи та форми роботи", "bullets": part})

        equipment = normalize_list(lesson_data.get("equipment"))
        resources = normalize_list(lesson_data.get("resources"))
        equipment_resources = labeled_list("Обладнання", equipment) + labeled_list("Ресурси", resources)
        if equipment_resources:
            for part in chunk(equipment_resources, 6):
                slides.append({"type": "bullets", "emoji": "🧰", "title": "Ресурси та обладнання", "bullets": part})

        for stage in lesson_data.get("lesson_flow", []):
            title = clean_text(stage.get("stage") or "")
            time_txt = format_time(stage.get("time_min"))
            if time_txt:
                title = f"{title} ({time_txt})"

            bullets = []
            bullets += labeled_list("Вчитель", stage.get("teacher_actions"))
            bullets += labeled_list("Учні", stage.get("student_actions"))
            bullets += labeled_list("Діяльність", stage.get("activities"))

            if not title or not bullets:
                continue
            for part in chunk(bullets, 6):
                slides.append({"type": "bullets", "emoji": "??", "title": title, "bullets": part})

        return slides





    async def create_docx(self, lesson_data, filename, doc_paths=None, template_structure=None):
        path = f"storage/{filename}.docx"

        def build():
            data = lesson_data or {}
            selected_template = doc_paths[0] if doc_paths else None
            template_blueprint = self._build_template_blueprint(template_structure) or {}
            allowed_optional_sections = self._get_allowed_optional_sections(template_blueprint)
            allowed_header_labels = {
                self._normalize_template_name(label)
                for label in (template_blueprint.get("header_order") or self.DEFAULT_GENERIC_HEADER_ORDER)
                if clean_text(label)
            }
            stage_substep_allowlist = self._build_stage_substep_allowlist(template_structure)

            if selected_template:
                try:
                    doc = Document(selected_template)
                    body = doc._element.body
                    for element in list(body):
                        if element.tag.endswith("sectPr"):
                            continue
                        body.remove(element)
                except Exception as exc:
                    logger.warning("Failed to open template docx %s: %s. Falling back to blank document.", selected_template, exc)
                    doc = Document()
            else:
                doc = Document()

            style = doc.styles["Normal"]
            style.font.name = "Times New Roman"
            style.font.size = DocxPt(14)
            style._element.rPr.rFonts.set(qn("w:ascii"), "Times New Roman")
            style._element.rPr.rFonts.set(qn("w:hAnsi"), "Times New Roman")
            style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
            style._element.rPr.rFonts.set(qn("w:cs"), "Times New Roman")
            style.paragraph_format.space_before = DocxPt(0)
            style.paragraph_format.space_after = DocxPt(6)
            style.paragraph_format.line_spacing = 1.15
            style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

            def ensure_paragraph_style(
                style_name,
                *,
                base_name="Normal",
                font_size=14,
                bold=False,
                italic=False,
                alignment=WD_ALIGN_PARAGRAPH.LEFT,
                space_before=0,
                space_after=6,
                keep_with_next=False,
            ):
                try:
                    target_style = doc.styles[style_name]
                except KeyError:
                    target_style = doc.styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)

                try:
                    base_style = doc.styles[base_name]
                    target_style.base_style = base_style
                except KeyError:
                    pass

                target_style.font.name = "Times New Roman"
                target_style.font.size = DocxPt(font_size)
                target_style.font.bold = bold
                target_style.font.italic = italic
                target_style._element.rPr.rFonts.set(qn("w:ascii"), "Times New Roman")
                target_style._element.rPr.rFonts.set(qn("w:hAnsi"), "Times New Roman")
                target_style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
                target_style._element.rPr.rFonts.set(qn("w:cs"), "Times New Roman")
                target_style.paragraph_format.alignment = alignment
                target_style.paragraph_format.space_before = DocxPt(space_before)
                target_style.paragraph_format.space_after = DocxPt(space_after)
                target_style.paragraph_format.keep_with_next = keep_with_next
                return target_style

            body_style = ensure_paragraph_style(
                "Metodist Body",
                base_name="Normal",
                font_size=14,
                bold=False,
                alignment=WD_ALIGN_PARAGRAPH.JUSTIFY,
                space_before=0,
                space_after=6,
            )
            title_style = ensure_paragraph_style(
                "Metodist Title",
                base_name="Normal",
                font_size=18,
                bold=True,
                alignment=WD_ALIGN_PARAGRAPH.CENTER,
                space_before=0,
                space_after=12,
                keep_with_next=True,
            )
            section_style = ensure_paragraph_style(
                "Metodist Section",
                base_name="Normal",
                font_size=16,
                bold=True,
                alignment=WD_ALIGN_PARAGRAPH.LEFT,
                space_before=10,
                space_after=4,
                keep_with_next=True,
            )
            stage_style = ensure_paragraph_style(
                "Metodist Stage",
                base_name="Normal",
                font_size=14,
                bold=True,
                italic=True,
                alignment=WD_ALIGN_PARAGRAPH.LEFT,
                space_before=8,
                space_after=3,
                keep_with_next=True,
            )
            substep_style = ensure_paragraph_style(
                "Metodist Substep",
                base_name="Normal",
                font_size=14,
                bold=True,
                alignment=WD_ALIGN_PARAGRAPH.LEFT,
                space_before=4,
                space_after=2,
                keep_with_next=True,
            )

            def make_generic_model():
                topic = clean_text(data.get("topic") or "")
                grade = clean_text(data.get("grade") or "")
                subject = clean_text(data.get("subject") or "")
                goal = self._join_items(data.get("goal"), separator=". ")
                lesson_type = clean_text(data.get("lesson_type") or "")
                education_area = clean_text(data.get("education_area") or "")
                equipment = self._join_items(self._prune_equipment_items(data.get("equipment")), separator="; ")
                header_fields = []
                for label, value in [
                    ("Тема", topic),
                    ("Клас", grade),
                    ("Предмет", subject),
                    ("Мета", goal),
                    ("Тип уроку", lesson_type),
                    ("Освітня галузь", education_area),
                    ("Обладнання", equipment),
                ]:
                    if clean_text(value) and self._normalize_template_name(label) in allowed_header_labels:
                        header_fields.append({"label": label, "value": normalize_sentence_punctuation(value), "style": "Normal"})

                sections = []

                def make_content_node(title, items=None, style="Normal", content_style="Normal", display_title=None):
                    cleaned_items = [normalize_sentence_punctuation(item) for item in normalize_list(items)]
                    return {
                        "title": title,
                        "display_title": display_title or title,
                        "style": style,
                        "content_style": content_style,
                        "items": [item for item in cleaned_items if item],
                        "sample_item_styles": [],
                        "substeps": [],
                        "children_order": [],
                    }

                def attach_substeps(node, groups, allowed_labels=None):
                    normalized_allowed = None
                    if allowed_labels:
                        normalized_allowed = {
                            self._normalize_template_name(label)
                            for label in allowed_labels
                            if clean_text(label)
                        }
                    substeps = []
                    for label, items in groups:
                        if normalized_allowed is not None and self._normalize_template_name(label) not in normalized_allowed:
                            continue
                        cleaned_items = [normalize_sentence_punctuation(item) for item in normalize_list(items)]
                        cleaned_items = [item for item in cleaned_items if item]
                        if not cleaned_items:
                            continue
                        substeps.append(
                            {
                                "title": label,
                                "display_title": label,
                                "style": "Normal",
                                "content_style": "Normal",
                                "items": cleaned_items,
                                "sample_item_styles": [],
                                "substeps": [],
                                "children_order": [],
                            }
                        )
                    if substeps:
                        node["substeps"] = substeps
                        node["children_order"] = [("substep", index) for index in range(len(substeps))]
                    return node

                def add_section(title, items):
                    if self.OPTIONAL_SECTION_FIELD_MAP.get(title) not in allowed_optional_sections:
                        return
                    cleaned_items = normalize_list(items)
                    if not cleaned_items:
                        return
                    sections.append(
                        {
                            **make_content_node(title, cleaned_items, style="Heading 2"),
                            "stages": [],
                        }
                    )

                add_section("Завдання", data.get("tasks"))

                expected_results = normalize_expected_results(data.get("expected_results") or {})
                expected_section = {**make_content_node("Очікувані результати", [], style="Heading 2"), "stages": []}
                attach_substeps(
                    expected_section,
                    [
                        ("Знаннєві результати", expected_results.get("knowledge")),
                        ("Діяльнісні результати", expected_results.get("skills")),
                        ("Ціннісні результати", expected_results.get("values")),
                    ],
                )
                if expected_section.get("substeps") and "expected_results" in allowed_optional_sections:
                    sections.append(expected_section)
                add_section("Ключові компетентності", data.get("key_competencies"))
                add_section("Наскрізні вміння", data.get("cross_cutting_skills"))
                add_section("Ціннісні орієнтири", data.get("values"))
                add_section("Інтеграція", data.get("integration"))
                add_section("Методи", data.get("methods"))
                add_section("Форми роботи", data.get("forms"))
                add_section("Формувальне оцінювання", data.get("assessment"))
                add_section("Диференціація", data.get("differentiation"))
                add_section("Ресурси", data.get("resources"))
                add_section("Домашнє завдання", data.get("homework"))

                flow_stages = []
                for stage in data.get("lesson_flow") or []:
                    stage_title = clean_text(stage.get("stage") or "") or "Етап уроку"
                    time_text = format_time(stage.get("time_min"))
                    display_title = f"{stage_title} ({time_text})" if time_text else stage_title
                    stage_node = make_content_node(stage_title, [], style="Heading 3", display_title=display_title)
                    allowed_stage_labels = stage_substep_allowlist.get(self._normalize_template_name(stage_title))
                    attach_substeps(
                        stage_node,
                        [
                            ("Вчитель", stage.get("teacher_actions")),
                            ("Учні", stage.get("student_actions")),
                            ("Діяльність", stage.get("activities")),
                        ],
                        allowed_labels=allowed_stage_labels or self.DEFAULT_GENERIC_STAGE_SUBSTEPS,
                    )
                    if not stage_node.get("substeps"):
                        fallback_items = []
                        for source_items in (
                            stage.get("teacher_actions"),
                            stage.get("student_actions"),
                            stage.get("activities"),
                        ):
                            fallback_items.extend(normalize_list(source_items))
                        stage_node["items"] = [
                            normalize_sentence_punctuation(item)
                            for item in dict.fromkeys(item for item in fallback_items if clean_text(item))
                        ]
                    flow_stages.append(stage_node)
                if flow_stages:
                    sections.append({
                        "title": "Хід уроку",
                        "display_title": "Хід уроку",
                        "style": "Heading 2",
                        "content_style": "Normal",
                        "items": [],
                        "sample_item_styles": [],
                        "substeps": [],
                        "stages": flow_stages,
                        "children_order": [("stage", index) for index in range(len(flow_stages))],
                    })

                return {
                    "header_fields": header_fields,
                    "sections": sections,
                }

            def make_nush_script_model():
                def strip_role_prefix(value):
                    text = normalize_sentence_punctuation(value)
                    if not text:
                        return ""
                    for prefix in ("Вчитель:", "Учитель:", "Учні:", "Діяльність:", "Teacher:", "Students:", "Activities:"):
                        if text.lower().startswith(prefix.lower()):
                            text = text[len(prefix):].strip()
                    return normalize_sentence_punctuation(text)

                def normalize_student_response(value):
                    text = strip_role_prefix(value)
                    if not text:
                        return ""
                    normalized = self._normalize_template_name(text)
                    if "типова відповідь учня" in normalized or "типова відповідь" in normalized:
                        numbers = re.findall(r"\d+", text)
                        if len(numbers) == 1 and len(text.split()) <= 8:
                            return f"({numbers[0]})"
                        return "(кілька відповідей)"
                    if "відповіді учнів" in normalized:
                        return "(відповіді учнів)"
                    if "кілька відповідей" in normalized:
                        return "(кілька відповідей)"
                    if "В«" in text or re.search(r"\d", text):
                        compact = text.replace("В«", "").replace("В»", "").strip()
                        if re.fullmatch(r"\d+(?:[.,]\d+)?", compact):
                            return f"({compact})"
                        return f"({compact})"
                    return "(відповіді учнів)"

                def stage_group_from_title(title):
                    group = self._classify_stage_name(title)
                    if group == "intro":
                        return ["Усний рахунок", "Каліграфічна хвилинка", "Робота з прикладами", "Математична розминка"]
                    if group == "closing":
                        return ["Підбиття підсумків", "Рефлексія", "Самоперевірка", "Домашнє завдання"]
                    return ["Робота з прикладами", "Робота в парах", "Розв'язування задач", "Математична гра", "Перевірка"]

                def polish_action_line(text):
                    cleaned = strip_role_prefix(text)
                    if not cleaned:
                        return ""
                    normalized = self._normalize_template_name(cleaned)
                    replacements = (
                        ("сьогодні працюємо з табличного множення та ділення на 6", "Працюємо з табличним множенням і діленням на 6."),
                        ("сьогодні працюємо з", "Працюємо з"),
                    )
                    for src, dst in replacements:
                        if src in normalized:
                            cleaned = dst
                    if normalized.startswith("пригадаємо"):
                        return "– Обчисліть приклади з дошки і поясніть один обраний приклад."
                    if normalized.startswith("згадаємо"):
                        return "– Порівняйте два приклади та поясніть, чим вони відрізняються."
                    if normalized.startswith("поговоримо"):
                        return "– Запишіть відповідь і коротко поясніть свій спосіб."
                    if normalized.startswith("давайте"):
                        return f"– {cleaned}"
                    if "?" in cleaned:
                        return f"– {cleaned}"
                    return f"– {cleaned}"

                def classify_script_line(line):
                    normalized = self._normalize_template_name(line)
                    if line.startswith("("):
                        return "response"
                    if any(token in normalized for token in ("перевір", "звір", "самоперев", "оцініть", "сигналом", "смайлик")):
                        return "check"
                    if "?" in line:
                        return "question"
                    return "action"

                def split_into_blocks(lines):
                    lines = [line for line in lines if clean_text(line)]
                    if not lines:
                        return [["– Виконайте коротке завдання за темою уроку.", "(відповіді учнів)"]]
                    buckets = {"action": [], "question": [], "response": [], "check": []}
                    for line in lines:
                        buckets[classify_script_line(line)].append(line)

                    blocks = []
                    while len(blocks) < 6 and any(buckets.values()):
                        block = []
                        if buckets["action"]:
                            block.append(buckets["action"].pop(0))
                        if buckets["question"]:
                            block.append(buckets["question"].pop(0))
                        if buckets["response"]:
                            block.append(buckets["response"].pop(0))
                        elif any("?" in item for item in block):
                            block.append("(відповіді учнів)")
                        if buckets["check"]:
                            block.append(buckets["check"].pop(0))
                        elif block:
                            block.append("– Перевірмо разом.")
                        block = [item for item in block if clean_text(item)]
                        if block:
                            blocks.append(block)

                    if len(blocks) == 1:
                        blocks.append(["– Виконайте ще один приклад самостійно.", "(відповіді учнів)", "– Перевірмо разом."])
                    return blocks[:6]

                def stage_script_lines(stage):
                    lines = []
                    for item in normalize_list(stage.get("teacher_actions")):
                        cleaned = polish_action_line(item)
                        if cleaned:
                            lines.append(cleaned)
                    for item in normalize_list(stage.get("activities")):
                        cleaned = polish_action_line(item)
                        if cleaned:
                            lines.append(cleaned)
                    for item in normalize_list(stage.get("student_actions")):
                        response = normalize_student_response(item)
                        if response:
                            lines.append(response)
                    for item in normalize_list(stage.get("assessment")):
                        cleaned = polish_action_line(item)
                        if cleaned:
                            lines.append(cleaned)
                    for item in normalize_list(stage.get("reflection")):
                        cleaned = polish_action_line(item)
                        if cleaned:
                            lines.append(cleaned)
                    dedup = []
                    seen = set()
                    for line in lines:
                        key = self._normalize_template_name(line)
                        if key and key not in seen:
                            seen.add(key)
                            dedup.append(line)
                    return dedup

                def build_stage_node(stage, index):
                    raw_title = clean_text(stage.get("stage") or "") or "Етап уроку"
                    if re.match(r"^[IVX]+\.\s", raw_title, flags=re.IGNORECASE):
                        display_title = raw_title
                    else:
                        roman = ["I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X"]
                        prefix = roman[index] if index < len(roman) else f"{index + 1}"
                        display_title = f"{prefix}. {raw_title}"

                    script_lines = stage_script_lines(stage)
                    blocks = split_into_blocks(script_lines)
                    labels = stage_group_from_title(raw_title)
                    substeps = []
                    for block_index, block_lines in enumerate(blocks, start=1):
                        label = labels[(block_index - 1) % len(labels)]
                        substeps.append(
                            {
                                "title": f"{block_index}. {label}",
                                "display_title": f"{block_index}. {label}",
                                "style": "Metodist Substep",
                                "content_style": "Metodist Body",
                                "items": block_lines,
                                "sample_item_styles": [],
                                "substeps": [],
                                "children_order": [],
                            }
                        )
                    return {
                        "title": raw_title,
                        "display_title": display_title,
                        "style": "Metodist Stage",
                        "content_style": "Metodist Body",
                        "items": [],
                        "sample_item_styles": [],
                        "substeps": substeps,
                        "stages": [],
                        "children_order": [("substep", idx) for idx in range(len(substeps))],
                    }

                topic = clean_text(data.get("topic") or "")
                grade = clean_text(data.get("grade") or "")
                subject = clean_text(data.get("subject") or "")
                goal = self._join_items(data.get("goal"), separator=". ")
                equipment = self._join_items(self._prune_equipment_items(data.get("equipment")), separator="; ")
                lesson_type = clean_text(data.get("lesson_type") or "")
                expected_results = normalize_expected_results(data.get("expected_results") or {})
                expected_flat = normalize_list([
                    *(expected_results.get("knowledge") or []),
                    *(expected_results.get("skills") or []),
                    *(expected_results.get("values") or []),
                ])

                header_fields = [
                    {"label": "Тема", "value": topic or "—", "style": "Metodist Body"},
                    {"label": "Мета", "value": goal or "—", "style": "Metodist Body"},
                    {"label": "Обладнання", "value": equipment or "—", "style": "Metodist Body"},
                ]
                if lesson_type:
                    header_fields.append({"label": "Тип уроку", "value": lesson_type, "style": "Metodist Body"})
                if expected_flat:
                    header_fields.append(
                        {
                            "label": "Очікувані результати",
                            "value": "; ".join(expected_flat[:6]),
                            "style": "Metodist Body",
                        }
                    )
                if grade:
                    header_fields.append({"label": "Клас", "value": grade, "style": "Metodist Body"})
                if subject:
                    header_fields.append({"label": "Предмет", "value": subject, "style": "Metodist Body"})

                flow_stages = []
                for idx, stage in enumerate(data.get("lesson_flow") or []):
                    if isinstance(stage, dict):
                        flow_stages.append(build_stage_node(stage, idx))
                if not flow_stages:
                    flow_stages.append(
                        build_stage_node(
                            {
                                "stage": "Організаційний момент",
                                "teacher_actions": ["Сьогодні працюємо за темою уроку."],
                                "student_actions": ["(відповіді учнів)"],
                                "activities": ["Обчисліть усно: 6×3, 6×4, 24:6."],
                            },
                            0,
                        )
                    )

                sections = [
                    {
                        "title": "Хід уроку",
                        "display_title": "Хід уроку",
                        "style": "Metodist Section",
                        "content_style": "Metodist Body",
                        "items": [],
                        "sample_item_styles": [],
                        "substeps": [],
                        "stages": flow_stages,
                        "children_order": [("stage", idx) for idx in range(len(flow_stages))],
                    }
                ]
                return {"header_fields": header_fields, "sections": sections}

            if isinstance(data, dict) and isinstance(data.get("lesson_flow"), list) and data.get("lesson_flow"):
                document_model = make_nush_script_model()
            else:
                document_model = data if data.get("header_fields") is not None and data.get("sections") is not None else make_generic_model()

            def resolve_style(style_name):
                if not style_name:
                    return None
                try:
                    return doc.styles[style_name]
                except KeyError:
                    normalized = clean_text(style_name).lower().replace("_", " ")
                    style_aliases = [
                        ("heading 1", title_style),
                        ("heading 2", section_style),
                        ("heading 3", stage_style),
                        ("title", title_style),
                        ("subtitle", section_style),
                        ("lesson-tema2", title_style),
                        ("step-tema2", stage_style),
                        ("normal", body_style),
                    ]
                    for alias, fallback_style in style_aliases:
                        if normalized == alias:
                            return fallback_style
                    return None

            def should_avoid_numbering_style(style_name):
                normalized = clean_text(style_name).lower()
                if not normalized:
                    return False
                blocked_tokens = ("list", "bullet", "number", "спис", "маркер", "перелік")
                return any(token in normalized for token in blocked_tokens)

            def normalize_ordered_substeps(node):
                substeps = node.get("substeps") or []
                ordered_indexes = []
                for index, substep in enumerate(substeps):
                    raw_title = clean_text_preserve_prefix(substep.get("display_title") or substep.get("title") or "")
                    if LEADING_ORDERED_TITLE_RE.match(raw_title):
                        ordered_indexes.append(index)
                if ordered_indexes:
                    counter = 1
                    for index in ordered_indexes:
                        substep = substeps[index]
                        raw_title = clean_text_preserve_prefix(substep.get("display_title") or substep.get("title") or "")
                        match = LEADING_ORDERED_TITLE_RE.match(raw_title)
                        if match:
                            substep["display_title"] = f"{counter}. {clean_text_preserve_prefix(match.group(1))}"
                            counter += 1
                for substep in substeps:
                    normalize_ordered_substeps(substep)
                for stage in node.get("stages") or []:
                    normalize_ordered_substeps(stage)

            def normalize_document_model(node):
                if not isinstance(node, dict):
                    return

                node["items"] = [
                    item
                    for item in (
                        normalize_sentence_punctuation(raw_item)
                        for raw_item in normalize_block_list(node.get("items") or [])
                    )
                    if item and not is_placeholder_text(item)
                ]

                substeps = [item for item in (node.get("substeps") or []) if isinstance(item, dict)]
                stages = [item for item in (node.get("stages") or []) if isinstance(item, dict)]
                node["substeps"] = substeps
                node["stages"] = stages

                children_order = []
                for child_kind, child_index in node.get("children_order") or []:
                    if child_kind == "substep" and 0 <= child_index < len(substeps):
                        children_order.append((child_kind, child_index))
                    elif child_kind == "stage" and 0 <= child_index < len(stages):
                        children_order.append((child_kind, child_index))
                if not children_order and (substeps or stages):
                    children_order = [("substep", index) for index in range(len(substeps))]
                    children_order.extend(("stage", index) for index in range(len(stages)))
                node["children_order"] = children_order

                for substep in substeps:
                    normalize_document_model(substep)
                for stage in stages:
                    normalize_document_model(stage)

            def preferred_title_style(role, original_style):
                cleaned_style = clean_text(original_style)
                if cleaned_style and cleaned_style.lower() != "normal":
                    return cleaned_style
                if role == "section":
                    return "Metodist Section"
                if role == "stage":
                    return "Metodist Stage"
                if role == "substep":
                    return "Metodist Substep"
                return "Metodist Body"

            def expand_render_item(text):
                normalized = normalize_sentence_punctuation(text)
                if not normalized:
                    return []

                if normalized.count(";") >= 2:
                    label_match = re.match(r"^([^:]{2,60}):\s*(.+)$", normalized)
                    if label_match:
                        label = normalize_sentence_punctuation(label_match.group(1))
                        parts = [
                            normalize_sentence_punctuation(part)
                            for part in label_match.group(2).split(";")
                            if normalize_sentence_punctuation(part)
                        ]
                        if len(parts) >= 2:
                            return [f"{label}: {parts[0]}", *parts[1:]]
                    parts = [
                        normalize_sentence_punctuation(part)
                        for part in normalized.split(";")
                        if normalize_sentence_punctuation(part)
                    ]
                    if len(parts) >= 2:
                        return parts

                return [normalized]

            def add_paragraph_line(text, style_name=None, force_bold=False, align=None, space_after=4, space_before=0):
                text = normalize_sentence_punctuation(text)
                if not text:
                    return None
                paragraph = doc.add_paragraph()
                paragraph_style = None if should_avoid_numbering_style(style_name) else resolve_style(style_name)
                if paragraph_style is not None:
                    paragraph.style = paragraph_style
                paragraph.paragraph_format.space_before = DocxPt(space_before)
                paragraph.paragraph_format.space_after = DocxPt(space_after)
                if align is not None:
                    paragraph.paragraph_format.alignment = align
                run = paragraph.add_run(text)
                run.font.name = "Times New Roman"
                run.font.size = DocxPt(14)
                if force_bold:
                    run.bold = True
                return paragraph

            def add_labeled_paragraph_line(label, value, style_name=None, align=WD_ALIGN_PARAGRAPH.LEFT, space_after=2):
                clean_label = clean_text(label)
                clean_value = normalize_sentence_punctuation(value)
                if not clean_label or not clean_value:
                    return None
                paragraph = doc.add_paragraph()
                paragraph_style = resolve_style(style_name or "Metodist Body")
                if paragraph_style is not None:
                    paragraph.style = paragraph_style
                paragraph.paragraph_format.alignment = align
                paragraph.paragraph_format.space_before = DocxPt(0)
                paragraph.paragraph_format.space_after = DocxPt(space_after)

                label_run = paragraph.add_run(f"{clean_label}: ")
                label_run.font.name = "Times New Roman"
                label_run.font.size = DocxPt(14)
                label_run.bold = True

                value_run = paragraph.add_run(clean_value)
                value_run.font.name = "Times New Roman"
                value_run.font.size = DocxPt(14)
                return paragraph

            def normalize_heading_key(value):
                text = clean_text_preserve_prefix(value or "").strip()
                text = re.sub(r"^\d+\s*[.)]\s*", "", text)
                text = text.rstrip(":").strip().lower()
                return text

            MAIN_LESSON_HEADINGS = {
                "тема уроку",
                "клас",
                "предмет",
                "тип уроку",
                "мета уроку",
                "очікувані результати",
                "обладнання та матеріали",
                "обладнання",
                "хід уроку",
                "домашнє завдання",
            }
            LESSON_STAGE_HEADINGS = {
                "організаційний момент",
                "актуалізація опорних знань",
                "актуалізація знань",
                "мотивація навчальної діяльності",
                "вивчення нового матеріалу",
                "закріплення знань",
                "підсумок уроку",
            }

            def split_heading_line(line):
                text = normalize_sentence_punctuation(line)
                if not text:
                    return "", ""
                if ":" in text:
                    before, after = text.split(":", 1)
                    key = normalize_heading_key(before)
                    if key in MAIN_LESSON_HEADINGS:
                        return before.strip(), after.strip()
                key = normalize_heading_key(text)
                if key in MAIN_LESSON_HEADINGS or key in LESSON_STAGE_HEADINGS:
                    return text.strip().rstrip(":"), ""
                return "", text

            def split_role_line(line):
                text = normalize_sentence_punctuation(line)
                match = re.match(r"^\s*(Учитель|Вчитель|Учні|Завдання/вправа|Завдання|Вправа)\s*:\s*(.+?)\s*$", text, flags=re.IGNORECASE)
                if not match:
                    return "", text
                label = match.group(1)
                value = match.group(2).strip()
                if label.lower() == "вчитель":
                    label = "Учитель"
                if label.lower() in {"завдання", "вправа"}:
                    label = "Завдання"
                return label, value

            def add_docx_cell_text(cell, text, *, bold=False):
                cell.text = ""
                paragraph = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph()
                paragraph.paragraph_format.space_before = DocxPt(0)
                paragraph.paragraph_format.space_after = DocxPt(2)
                run = paragraph.add_run(normalize_sentence_punctuation(text))
                run.font.name = "Times New Roman"
                run.font.size = DocxPt(12)
                run.bold = bold

            def set_cell_border(cell):
                tc_pr = cell._tc.get_or_add_tcPr()
                tc_borders = tc_pr.first_child_found_in("w:tcBorders")
                if tc_borders is None:
                    tc_borders = OxmlElement("w:tcBorders")
                    tc_pr.append(tc_borders)
                for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
                    tag = f"w:{edge}"
                    element = tc_borders.find(qn(tag))
                    if element is None:
                        element = OxmlElement(tag)
                        tc_borders.append(element)
                    element.set(qn("w:val"), "single")
                    element.set(qn("w:sz"), "4")
                    element.set(qn("w:space"), "0")
                    element.set(qn("w:color"), "808080")

            def style_lesson_flow_table(table):
                try:
                    table.style = "Table Grid"
                except KeyError:
                    # Some user-provided lesson templates do not contain the built-in table style.
                    # Keep rendering stable and apply visible borders directly to cells instead.
                    pass
                for row in table.rows:
                    for cell in row.cells:
                        set_cell_border(cell)

            def append_stage_value(stage, label, value):
                if not value:
                    return
                key = {
                    "Учитель": "teacher",
                    "Вчитель": "teacher",
                    "Учні": "students",
                    "Завдання/вправа": "task",
                    "Завдання": "task",
                    "Вправа": "task",
                }.get(label, "task")
                stage[key].append(value)

            def parse_lesson_flow_table(lines):
                stages = []
                current = None
                for raw in lines:
                    line = normalize_sentence_punctuation(raw)
                    if not line:
                        continue
                    heading, tail = split_heading_line(line)
                    heading_key = normalize_heading_key(heading)
                    if heading_key in LESSON_STAGE_HEADINGS:
                        current = {"stage": heading, "teacher": [], "students": [], "task": []}
                        stages.append(current)
                        if tail:
                            current["task"].append(tail)
                        continue
                    if current is None:
                        current = {"stage": "Етап уроку", "teacher": [], "students": [], "task": []}
                        stages.append(current)
                    role, value = split_role_line(line)
                    if role:
                        append_stage_value(current, role, value)
                    else:
                        current["task"].append(line)
                return [
                    stage
                    for stage in stages
                    if stage.get("stage") or stage.get("teacher") or stage.get("students") or stage.get("task")
                ]

            def render_lesson_flow_table(flow_lines):
                stages = parse_lesson_flow_table(flow_lines)
                if not stages:
                    return False
                table = doc.add_table(rows=1, cols=4)
                style_lesson_flow_table(table)
                headers = ("Етап", "Учитель", "Учні", "Завдання")
                for index, header in enumerate(headers):
                    add_docx_cell_text(table.rows[0].cells[index], header, bold=True)
                for stage in stages:
                    cells = table.add_row().cells
                    for cell in cells:
                        set_cell_border(cell)
                    add_docx_cell_text(cells[0], stage.get("stage") or "")
                    add_docx_cell_text(cells[1], "\n".join(stage.get("teacher") or []))
                    add_docx_cell_text(cells[2], "\n".join(stage.get("students") or []))
                    add_docx_cell_text(cells[3], "\n".join(stage.get("task") or []))
                spacer = doc.add_paragraph()
                spacer.paragraph_format.space_after = DocxPt(8)
                return True

            def render_structured_lesson_items(items):
                normalized_items = [normalize_sentence_punctuation(item) for item in items if normalize_sentence_punctuation(item)]
                if not any(normalize_heading_key(item) == "хід уроку" for item in normalized_items):
                    return False

                index = 0
                while index < len(normalized_items):
                    line = normalized_items[index]
                    heading, tail = split_heading_line(line)
                    key = normalize_heading_key(heading)

                    if key == "хід уроку":
                        add_paragraph_line(
                            heading or "Хід уроку",
                            "Metodist Section",
                            force_bold=True,
                            align=WD_ALIGN_PARAGRAPH.LEFT,
                            space_before=14,
                            space_after=7,
                        )
                        flow_lines = []
                        index += 1
                        while index < len(normalized_items):
                            next_heading, _ = split_heading_line(normalized_items[index])
                            next_key = normalize_heading_key(next_heading)
                            if next_key in MAIN_LESSON_HEADINGS and next_key != "хід уроку":
                                break
                            flow_lines.append(normalized_items[index])
                            index += 1
                        if not render_lesson_flow_table(flow_lines):
                            for flow_line in flow_lines:
                                role, value = split_role_line(flow_line)
                                if role:
                                    add_labeled_paragraph_line(role, value, "Metodist Body", space_after=2)
                                else:
                                    add_paragraph_line(flow_line, "Metodist Body", align=WD_ALIGN_PARAGRAPH.JUSTIFY, space_after=3)
                        continue

                    if key in MAIN_LESSON_HEADINGS:
                        add_paragraph_line(
                            heading,
                            "Metodist Section",
                            force_bold=True,
                            align=WD_ALIGN_PARAGRAPH.LEFT,
                            space_before=12,
                            space_after=5,
                        )
                        if tail:
                            add_paragraph_line(tail, "Metodist Body", align=WD_ALIGN_PARAGRAPH.JUSTIFY, space_after=4)
                        index += 1
                        continue

                    role, value = split_role_line(line)
                    if role:
                        add_labeled_paragraph_line(role, value, "Metodist Body", space_after=2)
                    else:
                        add_paragraph_line(line, "Metodist Body", align=WD_ALIGN_PARAGRAPH.JUSTIFY, space_after=4)
                    index += 1
                return True

            def render_items(node):
                item_styles = list(node.get("sample_item_styles") or [])
                default_style = node.get("content_style") or "Normal"
                if render_structured_lesson_items(node.get("items") or []):
                    return
                for index, item in enumerate(node.get("items") or []):
                    style_name = item_styles[index] if index < len(item_styles) else default_style
                    expanded_items = expand_render_item(item)
                    for fragment_index, fragment in enumerate(expanded_items):
                        add_paragraph_line(
                            fragment,
                            style_name,
                            force_bold=False,
                            align=WD_ALIGN_PARAGRAPH.JUSTIFY,
                            space_after=3,
                            space_before=0 if fragment_index else 1,
                        )

            def render_node(node, role):
                title_text = clean_text_preserve_prefix(node.get("display_title") or node.get("title") or "")
                title_style = preferred_title_style(role, node.get("style") or "Normal")
                if title_text:
                    title_spacing = {"section": (8, 4), "stage": (6, 3), "substep": (4, 2)}.get(role, (3, 3))
                    add_paragraph_line(
                        title_text,
                        title_style,
                        force_bold=True,
                        align=WD_ALIGN_PARAGRAPH.LEFT,
                        space_after=title_spacing[1],
                        space_before=title_spacing[0],
                    )
                render_items(node)

                children_order = list(node.get("children_order") or [])
                if children_order:
                    for child_kind, child_index in children_order:
                        if child_kind == "substep":
                            substeps = node.get("substeps") or []
                            if 0 <= child_index < len(substeps):
                                render_node(substeps[child_index], "substep")
                        elif child_kind == "stage":
                            stages = node.get("stages") or []
                            if 0 <= child_index < len(stages):
                                render_node(stages[child_index], "stage")
                    return

                for substep in node.get("substeps") or []:
                    render_node(substep, "substep")
                for stage in node.get("stages") or []:
                    render_node(stage, "stage")

            header_fields = document_model.get("header_fields") or []
            sections = document_model.get("sections") or []
            for section in sections:
                normalize_document_model(section)
                normalize_ordered_substeps(section)

            if not selected_template and not header_fields:
                add_paragraph_line("Конспект уроку", "Metodist Title", force_bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, space_after=10)

            for field in header_fields:
                label = clean_text(field.get("label") or "")
                value = normalize_sentence_punctuation(field.get("value") or "")
                if not label or not value:
                    continue
                field_fragments = expand_render_item(value)
                if not field_fragments:
                    continue
                add_labeled_paragraph_line(
                    label,
                    field_fragments[0],
                    field.get("style") or "Metodist Body",
                    align=WD_ALIGN_PARAGRAPH.LEFT,
                    space_after=2,
                )
                for fragment in field_fragments[1:]:
                    add_paragraph_line(
                        fragment,
                        field.get("style") or "Metodist Body",
                        force_bold=False,
                        align=WD_ALIGN_PARAGRAPH.LEFT,
                        space_after=2,
                    )

            for section in sections:
                render_node(section, "section")

            doc.save(path)
            return path

        return await asyncio.to_thread(build)

    async def create_pptx(self, slides_data, filename):
        path = f"storage/{filename}.pptx"
        
        def build():
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
            renderer = SlideRenderer(prs, random.choice(PALETTES))
            
            slides = slides_data or [{"type": "title", "title": "Урок згенеровано"}]
            total = len(slides)
            for i, s in enumerate(slides, start=1):
                if s.get("type") == "title": 
                    renderer.render_title(s)
                else: 
                    renderer.render_bullets(s, i, total)
            
            prs.save(path)
            return path
            
        return await asyncio.to_thread(build)
