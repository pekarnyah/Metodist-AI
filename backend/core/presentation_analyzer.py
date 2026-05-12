from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PresentationAnalyzer:
    @staticmethod
    def _clean_text(value: str) -> str:
        return " ".join(str(value or "").replace(" ", " ").split()).strip()

    @classmethod
    def _extract_table_text(cls, shape) -> list[str]:
        if not getattr(shape, "has_table", False):
            return []

        lines: list[str] = []
        for row in shape.table.rows:
            cells = [cls._clean_text(cell.text) for cell in row.cells]
            cells = [cell for cell in cells if cell]
            if cells:
                lines.append(" | ".join(cells))
        return lines

    @classmethod
    def _extract_text_blocks(cls, shape) -> tuple[list[str], list[str]]:
        if not getattr(shape, "has_text_frame", False):
            return [], []

        non_empty = [cls._clean_text(paragraph.text) for paragraph in shape.text_frame.paragraphs]
        non_empty = [text for text in non_empty if text]
        if not non_empty:
            return [], []

        text_blocks: list[str] = []
        bullet_points: list[str] = []
        multi_paragraph_frame = len(non_empty) > 1

        for index, paragraph in enumerate(shape.text_frame.paragraphs):
            text = cls._clean_text(paragraph.text)
            if not text:
                continue

            is_bullet = getattr(paragraph, "level", 0) > 0
            if not is_bullet and multi_paragraph_frame and index > 0:
                is_bullet = True
            if not is_bullet and text[:1] in {"•", "-", "—"}:
                is_bullet = True
            if not is_bullet and len(text) > 2 and text[0].isdigit() and text[1] in {".", ")"}:
                is_bullet = True

            if is_bullet:
                bullet_points.append(text)
            else:
                text_blocks.append(text)

        if not text_blocks and bullet_points:
            text_blocks.append(bullet_points[0])
            bullet_points = bullet_points[1:]

        return text_blocks, bullet_points

    @classmethod
    def _extract_notes(cls, slide) -> str:
        try:
            notes_frame = slide.notes_slide.notes_text_frame
        except Exception:
            return ""

        parts: list[str] = []
        for paragraph in notes_frame.paragraphs:
            text = cls._clean_text(paragraph.text)
            if not text or text.lower().startswith("click to add notes"):
                continue
            parts.append(text)
        return " ".join(parts)

    @classmethod
    def _extract_image_label(cls, shape) -> str:
        candidates = [
            getattr(shape, "name", ""),
            getattr(getattr(shape, "image", None), "filename", ""),
        ]
        for candidate in candidates:
            cleaned = cls._clean_text(candidate)
            lowered = cleaned.lower()
            if cleaned and not lowered.startswith("picture") and lowered not in {"image", "jpg", "png"}:
                return cleaned
        return ""

    @classmethod
    def analyze(cls, file_path: str) -> list[dict]:
        prs = Presentation(file_path)
        slides_data = []

        for index, slide in enumerate(prs.slides, start=1):
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = cls._clean_text(slide.shapes.title.text)

            text_blocks: list[str] = []
            bullet_points: list[str] = []
            image_labels: list[str] = []
            images_count = 0

            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                if getattr(shape, "has_table", False):
                    text_blocks.extend(cls._extract_table_text(shape))
                    continue

                if getattr(shape, "has_text_frame", False):
                    blocks, bullets = cls._extract_text_blocks(shape)
                    text_blocks.extend(blocks)
                    bullet_points.extend(bullets)
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images_count += 1
                    image_label = cls._extract_image_label(shape)
                    if image_label:
                        image_labels.append(image_label)

            notes = cls._extract_notes(slide)
            full_text_parts = [part for part in text_blocks + bullet_points + ([notes] if notes else []) if part]
            is_title_slide = index == 1 and len(full_text_parts) <= 1 and images_count == 0

            slides_data.append({
                "number": index,
                "title": title,
                "text": " ".join(full_text_parts).strip(),
                "text_blocks": text_blocks,
                "bullet_points": bullet_points,
                "notes": notes,
                "images": images_count,
                "image_labels": image_labels,
                "is_title_slide": is_title_slide,
            })

        return slides_data
